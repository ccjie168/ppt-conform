"""模板格式提取器：从模板的Master/Layout中提取占位符的格式规范"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from lxml import etree


class TemplateFormatExtractor:
    """从模板提取标题/正文占位符的字体、字号、颜色规范"""

    _NSMAP = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    def extract_placeholder_formats(self, template_path: str) -> dict:
        """提取模板中各占位符的格式规范
        
        返回结构:
        {
            "title": {"font_name": "...", "font_size": Pt, "color": "RRGGBB", "bold": bool, "alignment": ...},
            "body": {"font_name": "...", "font_size": Pt, "color": "RRGGBB", ...},
            "subtitle": {...},
        }
        """
        prs = Presentation(template_path)
        formats = {}

        # 优先从第一个Master提取
        if prs.slide_masters:
            master = prs.slide_masters[0]
            master_formats = self._extract_formats_from_shapes(master.shapes)
            for k, v in master_formats.items():
                formats.setdefault(k, v)

        # 从Layouts补充（覆盖Master的设置）
        for layout in prs.slide_layouts:
            layout_formats = self._extract_formats_from_shapes(layout.shapes)
            for k, v in layout_formats.items():
                formats[k] = v

        return formats

    def _extract_formats_from_shapes(self, shapes) -> dict:
        formats = {}
        for shape in shapes:
            try:
                if not shape.is_placeholder:
                    continue
                phf = shape.placeholder_format
                ph_type = phf.type

                if ph_type == 0 or (phf.idx == 0 and ph_type is None):
                    formats["title"] = self._extract_text_format(shape)
                elif ph_type == 2 or (phf.idx == 1 and "title" not in formats):
                    formats["body"] = self._extract_text_format(shape)
                elif ph_type == 4:
                    formats["subtitle"] = self._extract_text_format(shape)
            except Exception:
                continue
        return formats

    def _extract_text_format(self, shape) -> dict:
        """从形状提取文本格式，包括继承自XML的格式"""
        fmt = {
            "font_name": None,
            "font_size": None,
            "bold": None,
            "italic": None,
            "color": None,
            "alignment": None,
        }

        try:
            tf = shape.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]
                if p.alignment is not None:
                    fmt["alignment"] = str(p.alignment)

                if p.runs:
                    run = p.runs[0]
                    font = run.font
                    if font.name:
                        fmt["font_name"] = font.name
                    if font.size:
                        fmt["font_size"] = font.size
                    if font.bold is not None:
                        fmt["bold"] = font.bold
                    if font.italic is not None:
                        fmt["italic"] = font.italic
                    try:
                        if font.color and font.color.rgb:
                            fmt["color"] = str(font.color.rgb)
                    except Exception:
                        pass
        except Exception:
            pass

        # 从XML深层解析（处理继承的格式）
        try:
            xml_format = self._parse_format_from_xml(shape)
            for key, value in xml_format.items():
                if value is not None and fmt.get(key) is None:
                    fmt[key] = value
        except Exception:
            pass

        return fmt

    def _parse_format_from_xml(self, shape) -> dict:
        """直接从XML元素解析格式，包括布局样式中的默认格式"""
        fmt = {}
        try:
            xml_elem = shape.element

            # 查找 txBody 中的段落属性和run属性
            lst_style = xml_elem.find(".//p:txBody/a:lstStyle", self._NSMAP)
            if lst_style is not None:
                # 标题级别样式
                for level in range(5):
                    lvl_style = lst_style.find(f"a:lvl{level+1}pPr", self._NSMAP)
                    if lvl_style is not None:
                        def_rpr = lvl_style.find("a:defRPr", self._NSMAP)
                        if def_rpr is not None:
                            if fmt.get("font_name") is None:
                                latin = def_rpr.find("a:latin", self._NSMAP)
                                if latin is not None:
                                    typeface = latin.get("typeface")
                                    if typeface:
                                        fmt["font_name"] = typeface
                            if fmt.get("font_size") is None:
                                sz = def_rpr.get("sz")
                                if sz:
                                    fmt["font_size"] = Pt(int(sz) / 100)
                            if fmt.get("bold") is None:
                                b = def_rpr.get("b")
                                if b:
                                    fmt["bold"] = b == "1"
                            if fmt.get("italic") is None:
                                i = def_rpr.get("i")
                                if i:
                                    fmt["italic"] = i == "1"
                            if fmt.get("color") is None:
                                solid_fill = def_rpr.find("a:solidFill", self._NSMAP)
                                if solid_fill is not None:
                                    srgb = solid_fill.find("a:srgbClr", self._NSMAP)
                                    if srgb is not None:
                                        fmt["color"] = srgb.get("val")
                            break
        except Exception:
            pass
        return fmt

    def extract_theme_fonts(self, template_path: str, master_index: int = 0) -> dict:
        """从模板主题提取字体方案"""
        prs = Presentation(template_path)
        if not prs.slide_masters or master_index >= len(prs.slide_masters):
            return {}
        theme_part = prs.slide_masters[master_index].part

        theme_fonts = {"major": None, "minor": None}
        try:
            for rel in theme_part.rels.values():
                if "theme" in rel.reltype:
                    theme_element = etree.fromstring(rel.target_part.blob)
                    font_scheme = theme_element.find(
                        ".//a:fontScheme", self._NSMAP
                    )
                    if font_scheme is not None:
                        major = font_scheme.find("a:majorFont", self._NSMAP)
                        minor = font_scheme.find("a:minorFont", self._NSMAP)
                        if major is not None:
                            latin = major.find("a:latin", self._NSMAP)
                            if latin is not None:
                                theme_fonts["major"] = latin.get("typeface")
                        if minor is not None:
                            latin = minor.find("a:latin", self._NSMAP)
                            if latin is not None:
                                theme_fonts["minor"] = latin.get("typeface")
                    break
        except Exception:
            pass

        return theme_fonts

    def extract_theme_colors(self, template_path: str, master_index: int = 0) -> dict:
        """从模板主题提取颜色方案"""
        prs = Presentation(template_path)
        if not prs.slide_masters or master_index >= len(prs.slide_masters):
            return {}
        theme_part = prs.slide_masters[master_index].part

        colors = {}
        try:
            for rel in theme_part.rels.values():
                if "theme" in rel.reltype:
                    theme_element = etree.fromstring(rel.target_part.blob)
                    color_scheme = theme_element.find(".//a:clrScheme", self._NSMAP)
                    if color_scheme is not None:
                        for child in color_scheme:
                            tag = child.tag.split('}')[-1]
                            srgb = child.find("a:srgbClr", self._NSMAP)
                            if srgb is not None:
                                colors[tag] = srgb.get("val")
                            sys_clr = child.find("a:sysClr", self._NSMAP)
                            if sys_clr is not None:
                                colors[tag] = sys_clr.get("lastClr", "")
                    break
        except Exception:
            pass

        return colors

    def get_text_color_for_master(self, template_path: str, master_index: int = 0) -> str | None:
        """根据Master背景获取合适的文字颜色

        浅色背景返回深色文字（dk1），深色背景返回浅色文字（lt1）
        """
        try:
            from core.analyzer.template_analyzer import TemplateAnalyzer
            analyzer = TemplateAnalyzer()
            result = analyzer.analyze(template_path)
            masters = result.get("masters", [])
            if master_index < len(masters):
                master_info = masters[master_index]
                bg = master_info.get("background", {})
                bg_type = bg.get("type", "solid")
                bg_color = bg.get("color")
                
                colors = self.extract_theme_colors(template_path, master_index)
                
                # 渐变背景：取渐变中最浅的颜色判断，或默认深色背景用白色
                if bg_type == "gradient":
                    gradient_colors = bg.get("display_color") or bg.get("gradient_colors") or []
                    if gradient_colors:
                        # 检查渐变中是否有深色
                        has_dark = any(self._is_dark_color(c) for c in gradient_colors if c)
                        if has_dark:
                            return colors.get("lt1", "FFFFFF")
                        else:
                            return colors.get("dk1", "000000")
                    # 默认渐变背景用白色文字
                    return colors.get("lt1", "FFFFFF")
                
                # 纯色背景
                if bg_color:
                    if self._is_dark_color(bg_color):
                        return colors.get("lt1", "FFFFFF")
                    else:
                        return colors.get("dk1", "000000")
        except Exception:
            pass
        return None

    def extract_placeholder_mapping(self, template_path: str, master_index: int = 0) -> dict:
        """提取模板中占位符的语义映射表
        
        返回结构:
        {
            "title": {"placeholder_type": 1, "placeholder_idx": 0, "name": "Title 1", "left": ..., "top": ...},
            "body_main": {"placeholder_type": 2, "placeholder_idx": 1, "name": "Content Placeholder 2", ...},
            "body_sidebar": {"placeholder_type": 7, "placeholder_idx": 2, "name": "Content Placeholder 3", ...},
            "subtitle": {"placeholder_type": 4, "placeholder_idx": 1, "name": "Subtitle 2", ...},
        }
        
        该映射表用于回填时根据内容的 semantic_role 匹配模板占位符
        """
        prs = Presentation(template_path)
        mapping = {}
        
        if not prs.slide_masters or master_index >= len(prs.slide_masters):
            return mapping
        
        master = prs.slide_masters[master_index]
        slide_width = prs.slide_width
        
        # 收集所有正文占位符（用于区分主正文和侧边栏）
        body_placeholders = []
        
        # 从Master的Layouts中提取占位符
        for layout in master.slide_layouts:
            for shape in layout.placeholders:
                try:
                    phf = shape.placeholder_format
                    ph_type = phf.type
                    ph_idx = phf.idx
                    
                    # 1=title, 3=ctrTitle
                    if ph_type in (1, 3):
                        mapping.setdefault("title", {
                            "placeholder_type": ph_type,
                            "placeholder_idx": ph_idx,
                            "name": shape.name,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                        })
                    # 4=subtitle
                    elif ph_type == 4:
                        mapping.setdefault("subtitle", {
                            "placeholder_type": ph_type,
                            "placeholder_idx": ph_idx,
                            "name": shape.name,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                        })
                    # 2=body, 7=text, 8=content, 9=object
                    elif ph_type in (2, 7, 8, 9, 10):
                        body_placeholders.append({
                            "placeholder_type": ph_type,
                            "placeholder_idx": ph_idx,
                            "name": shape.name,
                            "left": shape.left or 0,
                            "top": shape.top or 0,
                            "width": shape.width or 0,
                            "height": shape.height or 0,
                        })
                except Exception:
                    continue
        
        # 区分主正文和侧边栏
        if len(body_placeholders) == 1:
            mapping["body_main"] = body_placeholders[0]
        elif len(body_placeholders) >= 2:
            # 根据位置排序：左边的为主正文，右边的为侧边栏
            body_placeholders.sort(key=lambda x: x["left"])
            mapping["body_main"] = body_placeholders[0]
            mapping["body_sidebar"] = body_placeholders[1]
        # 如果没有body占位符，使用第一个非标题占位符
        if "body_main" not in mapping:
            mapping["body_main"] = {
                "placeholder_type": 2,
                "placeholder_idx": 1,
                "name": "Content Placeholder 2",
                "left": None,
                "top": None,
                "width": None,
                "height": None,
            }
        
        return mapping

    def _is_dark_color(self, hex_color: str) -> bool:
        """判断颜色是否为深色（基于亮度）"""
        try:
            hex_color = hex_color.lstrip('#')
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            # 计算相对亮度
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            return luminance < 0.5
        except Exception:
            return False
