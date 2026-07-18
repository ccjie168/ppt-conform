"""模板分析器：分析 PPT 模板中的 master 和 layout，识别风格特征"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from lxml import etree
from typing import Any, Dict


class TemplateAnalyzer:
    """分析 PPT 模板的结构与风格"""

    _NSMAP = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    def __init__(self):
        self._theme_colors = {}
        self._prs = None

    def analyze(self, template_path: str) -> dict:
        if not Path(template_path).exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        self._prs = Presentation(template_path)
        self._theme_colors = self._extract_theme_colors(self._prs)

        masters_info = self._analyze_masters(self._prs)

        return {
            "file": str(template_path),
            "theme_colors": self._theme_colors,
            "masters": masters_info,
            "total_layouts": sum(len(m["layouts"]) for m in masters_info),
            "total_slides": len(self._prs.slides),
            "style_matches": self._match_styles(self._prs),
        }

    def _extract_theme_colors(self, prs) -> Dict[str, str]:
        """从主题 XML 提取所有颜色定义（包括所有 Accent 和变体）"""
        colors = {}
        try:
            theme_part = None
            for rel in prs.part.rels.values():
                if "theme" in rel.reltype:
                    theme_part = rel.target_part
                    break

            if theme_part is None:
                return colors

            theme_xml = etree.fromstring(theme_part.blob)

            color_names = {
                "a:dk1": "Dark1",
                "a:lt1": "Light1",
                "a:dk2": "Dark2",
                "a:lt2": "Light2",
                "a:accent1": "Accent1",
                "a:accent2": "Accent2",
                "a:accent3": "Accent3",
                "a:accent4": "Accent4",
                "a:accent5": "Accent5",
                "a:accent6": "Accent6",
                "a:hlink": "Hyperlink",
                "a:folHlink": "FollowedHyperlink",
            }

            for full_tag, name in color_names.items():
                tag_name = full_tag.split(":")[1]
                xpath = f".//{{{self._NSMAP['a']}}}{tag_name}"
                elems = theme_xml.findall(xpath)
                if elems:
                    rgb = self._parse_color_elem(elems[0])
                    if rgb:
                        colors[name] = rgb

            for idx in range(1, 7):
                xpath = f".//a:accent{idx}", self._NSMAP
                elems = theme_xml.findall(f".//{{{self._NSMAP['a']}}}accent{idx}")
                if elems:
                    rgb = self._parse_color_elem(elems[0])
                    if rgb:
                        colors[f"Accent{idx}"] = rgb

        except Exception as e:
            pass

        colors.update({
            "SchneiderDarkGreen": "0A2F24",
            "SchneiderLightGreen": "E7FFD9",
            "SchneiderBrightGreen": "3DCD58",
        })

        return colors

    def _parse_color_elem(self, elem) -> str | None:
        """解析颜色元素，返回 RGB 字符串"""
        try:
            rgb_elem = elem.find(".//a:rgbClr", self._NSMAP)
            if rgb_elem is not None:
                val = rgb_elem.get("val", "")
                if val:
                    return val.upper()

            srgb_elem = elem.find(".//a:srgbClr", self._NSMAP)
            if srgb_elem is not None:
                val = srgb_elem.get("val", "")
                if val:
                    return val.upper()

            sys_clr = elem.find(".//a:sysClr", self._NSMAP)
            if sys_clr is not None:
                last_clr = sys_clr.get("lastClr", "")
                if last_clr:
                    return last_clr.upper()

            scheme_clr = elem.find(".//a:schemeClr", self._NSMAP)
            if scheme_clr is not None:
                val = scheme_clr.get("val", "")
                if val:
                    mapped = {
                        "dk1": "202020",
                        "lt1": "FFFFFF",
                        "dk2": "404040",
                        "lt2": "F0F0F0",
                    }
                    return mapped.get(val, None)
        except Exception:
            pass
        return None

    def _analyze_masters(self, prs) -> list[dict]:
        masters = []
        for idx, master in enumerate(prs.slide_masters):
            bg = self._get_background(master)
            shapes_bg = self._detect_shapes_background(master)
            
            if bg.get("type") == "inherit" and shapes_bg:
                bg = shapes_bg
            
            style_match = self._match_single_style(bg)
            master_info = {
                "index": idx,
                "name": self._get_master_name(master),
                "background": bg,
                "fonts": self._get_fonts(master),
                "layouts": self._analyze_layouts(master),
                "style_id": style_match["style_id"],
                "style_name": style_match["style_name"],
                "style_desc": self._describe_background(bg),
            }
            masters.append(master_info)
        return masters

    def _detect_shapes_background(self, master) -> dict | None:
        """检测Master中形状的背景填充（有些模板通过形状实现背景颜色）"""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            for shape in master.shapes:
                try:
                    fill = shape.fill
                    from pptx.enum.dml import MSO_FILL
                    fill_type = fill.type
                    
                    if fill_type == MSO_FILL.SOLID:
                        color = fill.fore_color
                        try:
                            rgb = str(color.rgb)
                            if rgb and rgb != "None":
                                return {"type": "solid", "color": rgb, "gradient": None, "theme_color": None, "display_color": rgb}
                        except Exception:
                            try:
                                tc = str(color.theme_color)
                                clean_tc = self._clean_theme_color_name(tc)
                                rgb = self._resolve_theme_color(clean_tc)
                                if rgb:
                                    return {"type": "solid", "color": rgb, "gradient": None, "theme_color": tc, "display_color": rgb}
                            except Exception:
                                pass
                    
                    elif fill_type == MSO_FILL.GRADIENT:
                        stops = []
                        display_colors = []
                        for stop in fill.gradient_stops:
                            try:
                                rgb = str(stop.color.rgb)
                                stops.append(rgb)
                                display_colors.append(rgb)
                            except Exception:
                                try:
                                    tc = str(stop.color.theme_color)
                                    clean_tc = self._clean_theme_color_name(tc)
                                    rgb = self._resolve_theme_color(clean_tc)
                                    stops.append(rgb if rgb else "?")
                                    display_colors.append(rgb if rgb else "?")
                                except Exception:
                                    stops.append("?")
                                    display_colors.append("?")
                        if stops and any(s != "?" for s in stops):
                            return {"type": "gradient", "color": None, "gradient": stops, "theme_color": None, "display_color": display_colors}
                except Exception:
                    pass

            for shape in master.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for sub_shape in shape.shapes:
                            try:
                                fill = sub_shape.fill
                                fill_type = fill.type
                                
                                if fill_type == MSO_FILL.SOLID:
                                    color = fill.fore_color
                                    try:
                                        rgb = str(color.rgb)
                                        if rgb and rgb != "None":
                                            return {"type": "solid", "color": rgb, "gradient": None, "theme_color": None, "display_color": rgb}
                                    except Exception:
                                        try:
                                            tc = str(color.theme_color)
                                            clean_tc = self._clean_theme_color_name(tc)
                                            rgb = self._resolve_theme_color(clean_tc)
                                            if rgb:
                                                return {"type": "solid", "color": rgb, "gradient": None, "theme_color": tc, "display_color": rgb}
                                        except Exception:
                                            pass
                                
                                elif fill_type == MSO_FILL.GRADIENT:
                                    stops = []
                                    display_colors = []
                                    for stop in fill.gradient_stops:
                                        try:
                                            rgb = str(stop.color.rgb)
                                            stops.append(rgb)
                                            display_colors.append(rgb)
                                        except Exception:
                                            try:
                                                tc = str(stop.color.theme_color)
                                                clean_tc = self._clean_theme_color_name(tc)
                                                rgb = self._resolve_theme_color(clean_tc)
                                                stops.append(rgb if rgb else "?")
                                                display_colors.append(rgb if rgb else "?")
                                            except Exception:
                                                stops.append("?")
                                                display_colors.append("?")
                                    if stops and any(s != "?" for s in stops):
                                        return {"type": "gradient", "color": None, "gradient": stops, "theme_color": None, "display_color": display_colors}
                            except Exception:
                                pass
                except Exception:
                    pass

            return None
        except Exception:
            return None

    def _analyze_layouts(self, master) -> list[dict]:
        layouts = []
        for idx, layout in enumerate(master.slide_layouts):
            layout_info = {
                "index": idx,
                "name": layout.name or "(未命名)",
                "background": self._get_background(layout),
                "placeholders": self._get_placeholders(layout),
                "type_guess": self._guess_layout_type(layout.name or ""),
            }
            layouts.append(layout_info)
        return layouts

    def _get_master_name(self, master) -> str:
        try:
            return master.name or f"Master_{id(master)}"
        except Exception:
            return "(未命名)"

    def _get_background(self, element) -> dict:
        """获取背景信息，处理所有填充类型"""
        try:
            bg = element.background
            fill = bg.fill

            from pptx.enum.dml import MSO_FILL
            fill_type_int = fill.type

            if fill_type_int == MSO_FILL.SOLID:
                return self._parse_solid_fill(fill)
            elif fill_type_int == MSO_FILL.GRADIENT:
                return self._parse_gradient_fill(fill)
            elif fill_type_int == MSO_FILL.BACKGROUND:
                return self._parse_background_fill(element)
            elif fill_type_int == MSO_FILL.PICTURE:
                return {"type": "picture", "color": None, "gradient": None, "theme_color": None, "display_color": None}
            elif fill_type_int is None:
                xml_result = self._parse_background_from_xml(element)
                if xml_result.get("type") != "inherit":
                    return xml_result
                shapes_bg = self._detect_shapes_background(element)
                if shapes_bg:
                    return shapes_bg
                return {"type": "inherit", "color": None, "gradient": None, "theme_color": None, "display_color": None}
            else:
                return {"type": str(fill_type_int), "color": None, "gradient": None, "theme_color": None, "display_color": None}
        except Exception as e:
            try:
                xml_result = self._parse_background_from_xml(element)
                if xml_result.get("type") != "inherit":
                    return xml_result
                shapes_bg = self._detect_shapes_background(element)
                if shapes_bg:
                    return shapes_bg
            except Exception:
                pass
            return {"type": "error", "color": None, "gradient": None, "theme_color": None, "display_color": None, "error": str(e)}

    def _parse_solid_fill(self, fill) -> dict:
        """解析纯色填充"""
        color = fill.fore_color
        try:
            rgb = str(color.rgb)
            return {"type": "solid", "color": rgb, "gradient": None, "theme_color": None, "display_color": rgb}
        except Exception:
            try:
                theme_color = color.theme_color
                theme_color_str = str(theme_color)
                clean_theme_color = self._clean_theme_color_name(theme_color_str)
                actual_rgb = self._resolve_theme_color(clean_theme_color)
                return {
                    "type": "solid",
                    "color": actual_rgb,
                    "gradient": None,
                    "theme_color": theme_color_str,
                    "display_color": actual_rgb,
                }
            except Exception:
                return {"type": "solid", "color": None, "gradient": None, "theme_color": None, "display_color": None}

    def _parse_gradient_fill(self, fill) -> dict:
        """解析渐变填充"""
        stops = []
        display_colors = []
        theme_stops = []
        try:
            for stop in fill.gradient_stops:
                try:
                    rgb = str(stop.color.rgb)
                    stops.append(rgb)
                    display_colors.append(rgb)
                    theme_stops.append(None)
                except Exception:
                    try:
                        tc = str(stop.color.theme_color)
                        clean_tc = self._clean_theme_color_name(tc)
                        rgb = self._resolve_theme_color(clean_tc)
                        stops.append(rgb if rgb else "?")
                        display_colors.append(rgb if rgb else "?")
                        theme_stops.append(tc)
                    except Exception:
                        stops.append("?")
                        display_colors.append("?")
                        theme_stops.append(None)
        except Exception:
            pass
        return {
            "type": "gradient",
            "color": None,
            "gradient": stops,
            "theme_color": theme_stops if any(theme_stops) else None,
            "display_color": display_colors,
        }

    def _parse_background_from_xml(self, element) -> dict:
        """直接从 XML 解析背景信息（用于 fill.type 为 None 的情况）"""
        try:
            xml_elem = element.element
            sp_elem = xml_elem.find(".//p:sp", self._NSMAP)
            if sp_elem is None:
                sp_elem = xml_elem.find(".//p:spPr", self._NSMAP)
            if sp_elem is not None:
                fill_elem = sp_elem.find(".//a:fill", self._NSMAP)
                if fill_elem is not None:
                    solid_fill = fill_elem.find(".//a:solidFill", self._NSMAP)
                    if solid_fill is not None:
                        rgb_elem = solid_fill.find(".//a:rgbClr", self._NSMAP)
                        if rgb_elem is not None:
                            rgb = rgb_elem.get("val", "").upper()
                            return {"type": "solid", "color": rgb, "gradient": None, "theme_color": None, "display_color": rgb}
                        scheme_clr = solid_fill.find(".//a:schemeClr", self._NSMAP)
                        if scheme_clr is not None:
                            val = scheme_clr.get("val", "")
                            theme_name = {"lt1": "Light1", "lt2": "Light2", "dk1": "Dark1", "dk2": "Dark2", "accent1": "Accent1", "accent2": "Accent2", "accent3": "Accent3", "accent4": "Accent4", "accent5": "Accent5", "accent6": "Accent6", "bg1": "Light1", "bg2": "Light2", "tx1": "Dark1", "tx2": "Dark2"}.get(val)
                            actual_rgb = self._theme_colors.get(theme_name) if theme_name else None
                            return {"type": "solid", "color": actual_rgb, "gradient": None, "theme_color": val.upper() if val else None, "display_color": actual_rgb}
                    grad_fill = fill_elem.find(".//a:gradFill", self._NSMAP)
                    if grad_fill is not None:
                        stops = []
                        display_colors = []
                        for gs in grad_fill.findall(".//a:gs", self._NSMAP):
                            solid_fill_gs = gs.find(".//a:solidFill", self._NSMAP)
                            if solid_fill_gs is not None:
                                rgb_elem_gs = solid_fill_gs.find(".//a:rgbClr", self._NSMAP)
                                if rgb_elem_gs is not None:
                                    rgb = rgb_elem_gs.get("val", "").upper()
                                    stops.append(rgb)
                                    display_colors.append(rgb)
                                else:
                                    scheme_clr_gs = solid_fill_gs.find(".//a:schemeClr", self._NSMAP)
                                    if scheme_clr_gs is not None:
                                        val = scheme_clr_gs.get("val", "")
                                        theme_name = {"lt1": "Light1", "lt2": "Light2", "dk1": "Dark1", "dk2": "Dark2", "accent1": "Accent1", "accent2": "Accent2", "accent3": "Accent3", "accent4": "Accent4", "accent5": "Accent5", "accent6": "Accent6", "bg1": "Light1", "bg2": "Light2", "tx1": "Dark1", "tx2": "Dark2"}.get(val)
                                        actual_rgb = self._theme_colors.get(theme_name) if theme_name else None
                                        stops.append(actual_rgb if actual_rgb else "?")
                                        display_colors.append(actual_rgb if actual_rgb else "?")
                        return {"type": "gradient", "color": None, "gradient": stops, "theme_color": None, "display_color": display_colors}
            return {"type": "inherit", "color": None, "gradient": None, "theme_color": None, "display_color": None}
        except Exception:
            return {"type": "inherit", "color": None, "gradient": None, "theme_color": None, "display_color": None}

    def _parse_background_fill(self, element) -> dict:
        """解析 BACKGROUND 类型填充（通过 XML 直接获取）"""
        try:
            xml_elem = element.element
            bg_elem = xml_elem.find(".//p:bg", self._NSMAP)
            if bg_elem is not None:
                bg_ref = bg_elem.find(".//p:bgRef", self._NSMAP)
                if bg_ref is not None:
                    scheme_clr = bg_ref.find(".//a:schemeClr", self._NSMAP)
                    if scheme_clr is not None:
                        val = scheme_clr.get("val", "")
                        if val:
                            theme_name = {
                                "bg1": "Light1", "bg2": "Light2",
                                "tx1": "Dark1", "tx2": "Dark2",
                                "accent1": "Accent1", "accent2": "Accent2",
                                "accent3": "Accent3", "accent4": "Accent4",
                                "accent5": "Accent5", "accent6": "Accent6",
                            }.get(val)
                            actual_rgb = self._theme_colors.get(theme_name) if theme_name else None
                            return {
                                "type": "solid",
                                "color": actual_rgb,
                                "gradient": None,
                                "theme_color": val.upper() if val else None,
                                "display_color": actual_rgb,
                            }

                bg_fill = bg_elem.find(".//p:bgFill", self._NSMAP)
                if bg_fill is not None:
                    solid_fill = bg_fill.find(".//a:solidFill", self._NSMAP)
                    if solid_fill is not None:
                        rgb_elem = solid_fill.find(".//a:rgbClr", self._NSMAP)
                        if rgb_elem is not None:
                            rgb = rgb_elem.get("val", "").upper()
                            return {"type": "solid", "color": rgb, "gradient": None, "theme_color": None, "display_color": rgb}
                        scheme_clr = solid_fill.find(".//a:schemeClr", self._NSMAP)
                        if scheme_clr is not None:
                            val = scheme_clr.get("val", "")
                            if val:
                                theme_name = {"lt1": "Light1", "lt2": "Light2", "dk1": "Dark1", "dk2": "Dark2", "accent1": "Accent1", "accent2": "Accent2", "accent3": "Accent3", "accent4": "Accent4", "accent5": "Accent5", "accent6": "Accent6", "bg1": "Light1", "bg2": "Light2", "tx1": "Dark1", "tx2": "Dark2"}.get(val)
                                actual_rgb = self._theme_colors.get(theme_name) if theme_name else None
                                return {"type": "solid", "color": actual_rgb, "gradient": None, "theme_color": val.upper() if val else None, "display_color": actual_rgb}

                    gradient_fill = bg_fill.find(".//a:gradFill", self._NSMAP)
                    if gradient_fill is not None:
                        stops = []
                        display_colors = []
                        theme_stops = []
                        for gs in gradient_fill.findall(".//a:gs", self._NSMAP):
                            solid_fill_gs = gs.find(".//a:solidFill", self._NSMAP)
                            if solid_fill_gs is not None:
                                rgb_elem_gs = solid_fill_gs.find(".//a:rgbClr", self._NSMAP)
                                if rgb_elem_gs is not None:
                                    rgb = rgb_elem_gs.get("val", "").upper()
                                    stops.append(rgb)
                                    display_colors.append(rgb)
                                else:
                                    scheme_clr_gs = solid_fill_gs.find(".//a:schemeClr", self._NSMAP)
                                    if scheme_clr_gs is not None:
                                        val = scheme_clr_gs.get("val", "")
                                        theme_name = {"lt1": "Light1", "lt2": "Light2", "dk1": "Dark1", "dk2": "Dark2", "accent1": "Accent1", "accent2": "Accent2", "accent3": "Accent3", "accent4": "Accent4", "accent5": "Accent5", "accent6": "Accent6", "bg1": "Light1", "bg2": "Light2", "tx1": "Dark1", "tx2": "Dark2"}.get(val)
                                        actual_rgb = self._theme_colors.get(theme_name) if theme_name else None
                                        stops.append(actual_rgb if actual_rgb else "?")
                                        display_colors.append(actual_rgb if actual_rgb else "?")
                                        theme_stops.append(val.upper() if val else None)
                        return {"type": "gradient", "color": None, "gradient": stops, "theme_color": theme_stops if theme_stops else None, "display_color": display_colors}

            return {"type": "inherit", "color": None, "gradient": None, "theme_color": None, "display_color": None}
        except Exception as e:
            return {"type": "background", "color": None, "gradient": None, "theme_color": None, "display_color": None, "error": str(e)}

    def _clean_theme_color_name(self, theme_color_str: str) -> str:
        """清理主题颜色名称，去除变体索引如 'BACKGROUND_1 (14)' -> 'BACKGROUND_1'"""
        if not theme_color_str:
            return theme_color_str
        import re
        cleaned = re.sub(r'\s*\(\d+\)\s*$', '', theme_color_str)
        return cleaned.strip()

    def _resolve_theme_color(self, theme_color_str: str) -> str | None:
        """将主题颜色名解析为实际 RGB"""
        mapping = {
            "BACKGROUND_1": self._theme_colors.get("Light1"),
            "BACKGROUND_2": self._theme_colors.get("Light2"),
            "TEXT_1": self._theme_colors.get("Dark1"),
            "TEXT_2": self._theme_colors.get("Dark2"),
            "ACCENT_1": self._theme_colors.get("Accent1"),
            "ACCENT_2": self._theme_colors.get("Accent2"),
            "ACCENT_3": self._theme_colors.get("Accent3"),
            "ACCENT_4": self._theme_colors.get("Accent4"),
            "ACCENT_5": self._theme_colors.get("Accent5"),
            "ACCENT_6": self._theme_colors.get("Accent6"),
            "HYPERLINK": self._theme_colors.get("Hyperlink"),
            "FOLLOWED_HYPERLINK": self._theme_colors.get("FollowedHyperlink"),
            "LT1": self._theme_colors.get("Light1"),
            "LT2": self._theme_colors.get("Light2"),
            "DK1": self._theme_colors.get("Dark1"),
            "DK2": self._theme_colors.get("Dark2"),
            "ACCENT1": self._theme_colors.get("Accent1"),
            "ACCENT2": self._theme_colors.get("Accent2"),
            "ACCENT3": self._theme_colors.get("Accent3"),
            "ACCENT4": self._theme_colors.get("Accent4"),
            "ACCENT5": self._theme_colors.get("Accent5"),
            "ACCENT6": self._theme_colors.get("Accent6"),
        }
        return mapping.get(theme_color_str)

    def _get_fonts(self, element) -> dict:
        try:
            master_part = element.part if hasattr(element, "part") else None
            if master_part is None:
                return {"major": None, "minor": None}

            theme_part = None
            try:
                for rel in master_part.rels.values():
                    if "theme" in rel.reltype:
                        theme_part = rel.target_part
                        break
            except Exception:
                pass

            if theme_part is None:
                return {"major": None, "minor": None}

            theme_xml = etree.fromstring(theme_part.blob)
            font_scheme = theme_xml.find(".//a:fontScheme", self._NSMAP)
            if font_scheme is None:
                return {"major": None, "minor": None}

            major = font_scheme.find(".//a:majorFont/a:latin", self._NSMAP)
            minor = font_scheme.find(".//a:minorFont/a:latin", self._NSMAP)

            return {
                "major": major.get("typeface") if major is not None else None,
                "minor": minor.get("typeface") if minor is not None else None,
            }
        except Exception as e:
            return {"major": None, "minor": None, "error": str(e)}

    def _get_placeholders(self, layout) -> list[dict]:
        placeholders = []
        try:
            for ph in layout.placeholders:
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                })
        except Exception:
            pass
        return placeholders

    def _guess_layout_type(self, name: str) -> str:
        name_lower = name.lower()
        if "封面" in name or "cover" in name_lower or "title" in name_lower:
            return "cover"
        if "章节" in name or "section" in name_lower:
            return "section"
        if "内容" in name or "content" in name_lower:
            return "content"
        if "结尾" in name or "结束" in name or "closing" in name_lower or "end" in name_lower:
            return "closing"
        return "unknown"

    def _match_styles(self, prs) -> list[dict]:
        matches = []
        for idx, master in enumerate(prs.slide_masters):
            bg = self._get_background(master)
            matched = self._match_single_style(bg)
            matches.append({
                "master_index": idx,
                "master_name": self._get_master_name(master),
                "background": bg,
                "matched_style": matched["style_id"],
                "matched_name": matched["style_name"],
                "confidence": matched["confidence"],
            })
        return matches

    def _match_single_style(self, bg: dict) -> dict:
        bg_type = bg.get("type")
        color = bg.get("color")
        gradient = bg.get("gradient")
        display_color = bg.get("display_color")

        if bg_type == "gradient" and gradient:
            if len(gradient) >= 2 and any(g and g != "?" for g in gradient):
                return {
                    "style_id": "F4",
                    "style_name": "渐变科技",
                    "confidence": "高",
                }

        if (bg_type == "solid" or bg_type == "theme") and color:
            color_lower = color.lower()
            if self._is_white(color_lower):
                return {"style_id": "F1", "style_name": "白色简约", "confidence": "高"}
            elif self._is_light_green(color_lower):
                return {"style_id": "F2", "style_name": "浅绿色清新", "confidence": "高"}
            elif self._is_dark_green(color_lower):
                return {"style_id": "F3", "style_name": "深绿色商务", "confidence": "高"}
            else:
                return {
                    "style_id": "?",
                    "style_name": f"未匹配（颜色 #{color}）",
                    "confidence": "低",
                }

        if display_color:
            if isinstance(display_color, list) and len(display_color) >= 2:
                return {
                    "style_id": "F4",
                    "style_name": "渐变科技",
                    "confidence": "高",
                }
            elif isinstance(display_color, str):
                dc_lower = display_color.lower()
                if self._is_white(dc_lower):
                    return {"style_id": "F1", "style_name": "白色简约", "confidence": "高"}
                elif self._is_light_green(dc_lower):
                    return {"style_id": "F2", "style_name": "浅绿色清新", "confidence": "高"}
                elif self._is_dark_green(dc_lower):
                    return {"style_id": "F3", "style_name": "深绿色商务", "confidence": "高"}

        return {
            "style_id": "?",
            "style_name": "需人工确认",
            "confidence": "低",
        }

    SCHNEIDER_COLORS = {
        "dark_green": {"hex": "0A2F24", "rgb": (10, 47, 36), "brightness": 31},
        "light_green": {"hex": "E7FFD9", "rgb": (231, 255, 217), "brightness": 234},
        "bright_green": {"hex": "3DCD58", "rgb": (61, 205, 88), "brightness": 118},
        "white": {"hex": "FFFFFF", "rgb": (255, 255, 255), "brightness": 255},
        "black": {"hex": "000000", "rgb": (0, 0, 0), "brightness": 0},
    }

    def _is_white(self, color: str) -> bool:
        try:
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            if r >= 240 and g >= 240 and b >= 240:
                return True
            return self._color_match(color, self.SCHNEIDER_COLORS["white"]["hex"], 15)
        except Exception:
            return False

    def _is_light_green(self, color: str) -> bool:
        try:
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            if self._color_match(color, self.SCHNEIDER_COLORS["light_green"]["hex"], 20):
                return True
            if g <= r or g <= b:
                return False
            green_diff = g - max(r, b)
            if green_diff < 5:
                return False
            brightness = (r + g + b) / 3
            if brightness < 180:
                return False
            return True
        except Exception:
            return False

    def _is_dark_green(self, color: str) -> bool:
        try:
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            if self._color_match(color, self.SCHNEIDER_COLORS["dark_green"]["hex"], 20):
                return True
            if self._color_match(color, self.SCHNEIDER_COLORS["bright_green"]["hex"], 20):
                return False
            if g <= r or g <= b:
                return False
            green_diff = g - max(r, b)
            if green_diff < 5:
                return False
            brightness = (r + g + b) / 3
            if brightness > 120:
                return False
            return True
        except Exception:
            return False

    def _color_match(self, color1: str, color2: str, tolerance: int) -> bool:
        try:
            r1, g1, b1 = int(color1[0:2], 16), int(color1[2:4], 16), int(color1[4:6], 16)
            r2, g2, b2 = int(color2[0:2], 16), int(color2[2:4], 16), int(color2[4:6], 16)
            return (
                abs(r1 - r2) <= tolerance
                and abs(g1 - g2) <= tolerance
                and abs(b1 - b2) <= tolerance
            )
        except Exception:
            return False

    def _describe_background(self, bg: dict) -> str:
        """生成背景描述文本，便于前端展示"""
        bg_type = bg.get("type")
        color = bg.get("color")
        display_color = bg.get("display_color")
        
        if bg_type == "gradient":
            gradient = bg.get("gradient", [])
            if gradient and len(gradient) >= 2:
                valid_colors = [g for g in gradient if g and g != "?"]
                if valid_colors:
                    return f"渐变色 ({' → '.join(valid_colors[:2])})"
                return "渐变背景"
            return "渐变背景"
        
        if bg_type == "solid":
            if color:
                if self._is_white(color):
                    return f"白色背景 (#{color})"
                elif self._is_light_green(color):
                    return f"浅绿色背景 (#{color})"
                elif self._is_dark_green(color):
                    return f"深绿色背景 (#{color})"
                else:
                    return f"纯色背景 (#{color})"
            elif display_color:
                if isinstance(display_color, str):
                    if self._is_white(display_color):
                        return f"白色背景 (#{display_color})"
                    elif self._is_light_green(display_color):
                        return f"浅绿色背景 (#{display_color})"
                    elif self._is_dark_green(display_color):
                        return f"深绿色背景 (#{display_color})"
                    return f"纯色背景 (#{display_color})"
            theme_color = bg.get("theme_color")
            if theme_color:
                return f"主题色背景 [{theme_color}]"
            return "纯色背景"
        
        if bg_type == "picture":
            return "图片背景"
        
        if bg_type == "inherit":
            return "继承背景"
        
        if display_color:
            if isinstance(display_color, list) and len(display_color) >= 2:
                valid_colors = [d for d in display_color if d and d != "?"]
                if valid_colors:
                    return f"渐变色 ({' → '.join(valid_colors[:2])})"
                return "渐变背景"
            elif isinstance(display_color, str):
                if self._is_white(display_color):
                    return f"白色背景 (#{display_color})"
                elif self._is_light_green(display_color):
                    return f"浅绿色背景 (#{display_color})"
                elif self._is_dark_green(display_color):
                    return f"深绿色背景 (#{display_color})"
                return f"纯色背景 (#{display_color})"
        
        return "未知背景"
