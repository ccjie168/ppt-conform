import yaml
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from core.models import ValidationIssue


class ValidationRule:
    def __init__(self, rule_id: str, name: str, level: str, description: str):
        self.rule_id = rule_id
        self.name = name
        self.level = level
        self.description = description

    def check(self, pptx_path: str) -> list[ValidationIssue]:
        raise NotImplementedError


class WatermarkTextRule(ValidationRule):
    def __init__(self):
        super().__init__("R040", "水印文本检查", "fail", "产物中不得存在已知水印文本")
        self.keywords = []
        self._load_keywords()

    def _load_keywords(self):
        config_path = Path(__file__).parent.parent.parent / "config" / "validation_rules.yaml"
        if config_path.exists():
            with open(config_path, "r", encoding="utf-8") as f:
                config = yaml.safe_load(f)
                self.keywords = config.get("watermark_keywords", [])

    def check(self, pptx_path: str) -> list[ValidationIssue]:
        issues = []
        prs = Presentation(pptx_path)

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text
                    text_stripped = text.strip()
                    for keyword in self.keywords:
                        if keyword in text:
                            # 区分水印和业务内容
                            # 水印通常是独立的短文本，如果文本远长于关键词，很可能是正文内容
                            if len(text_stripped) <= len(keyword) + 10:
                                issues.append(ValidationIssue(
                                    level=self.level,
                                    rule_id=self.rule_id,
                                    message=f"发现水印文本: {keyword}",
                                    slide_index=slide_idx + 1
                                ))
                            else:
                                # 长文本：只有当关键词占文本比例超过50%才判定为水印
                                if len(keyword) / len(text_stripped) > 0.5:
                                    issues.append(ValidationIssue(
                                        level=self.level,
                                        rule_id=self.rule_id,
                                        message=f"发现水印文本: {keyword}",
                                        slide_index=slide_idx + 1
                                    ))

        return issues


class FontWhitelistRule(ValidationRule):
    def __init__(self):
        super().__init__("R010", "字体白名单检查", "warning", "仅允许使用模板主题字体")
        self.font_whitelist = []
        self._load_fonts()

    def _load_fonts(self):
        config_path = Path(__file__).parent.parent.parent / "config" / "validation_rules.yaml"
        if config_path.exists():
            with open(config_path, "r", encoding="utf-8") as f:
                config = yaml.safe_load(f)
                self.font_whitelist = config.get("font_whitelist", [])

    def check(self, pptx_path: str) -> list[ValidationIssue]:
        issues = []
        prs = Presentation(pptx_path)

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            font_name = run.font.name or "Unknown"
                            if font_name not in self.font_whitelist:
                                issues.append(ValidationIssue(
                                    level=self.level,
                                    rule_id=self.rule_id,
                                    message=f"使用非白名单字体: {font_name}",
                                    slide_index=slide_idx
                                ))

        return issues


class ContentOverflowRule(ValidationRule):
    def __init__(self):
        super().__init__("R020", "内容溢出检查", "fail", "内容不得超出幻灯片边界")
        self.margin = 91440

    def check(self, pptx_path: str) -> list[ValidationIssue]:
        issues = []
        prs = Presentation(pptx_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                try:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    if left < -self.margin:
                        issues.append(ValidationIssue(
                            level=self.level,
                            rule_id=self.rule_id,
                            message=f"形状{shape_idx}左侧溢出: left={left/914400:.1f}英寸",
                            slide_index=slide_idx
                        ))
                    if top < -self.margin:
                        issues.append(ValidationIssue(
                            level=self.level,
                            rule_id=self.rule_id,
                            message=f"形状{shape_idx}顶部溢出: top={top/914400:.1f}英寸",
                            slide_index=slide_idx
                        ))
                    if left + width > slide_width + self.margin:
                        issues.append(ValidationIssue(
                            level=self.level,
                            rule_id=self.rule_id,
                            message=f"形状{shape_idx}右侧溢出: right={(left+width)/914400:.1f}英寸",
                            slide_index=slide_idx
                        ))
                    if top + height > slide_height + self.margin:
                        issues.append(ValidationIssue(
                            level=self.level,
                            rule_id=self.rule_id,
                            message=f"形状{shape_idx}底部溢出: bottom={(top+height)/914400:.1f}英寸",
                            slide_index=slide_idx
                        ))
                except Exception:
                    pass

        return issues


class TextOverflowRule(ValidationRule):
    """文本溢出精确检测：估算文本行数是否超出文本框容量"""

    def __init__(self):
        super().__init__("R030", "文本容量检查", "warning", "文本内容可能超出文本框容量")

    def check(self, pptx_path: str) -> list[ValidationIssue]:
        issues = []
        prs = Presentation(pptx_path)

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                # 跳过空文本框
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                try:
                    tf = shape.text_frame
                    if not tf.word_wrap:
                        continue

                    box_width = shape.width
                    box_height = shape.height
                    if box_width <= 0 or box_height <= 0:
                        continue

                    total_text_height = 0
                    for para in tf.paragraphs:
                        text = para.text
                        if not text:
                            total_text_height += Emu(914400 * 0.2)
                            continue

                        font_size = Emu(914400 * 0.18)  # 默认14pt
                        for run in para.runs:
                            if run.font.size:
                                font_size = run.font.size
                                break

                        char_width = font_size * 0.6
                        chars_per_line = max(1, int(box_width / char_width))
                        num_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)

                        line_height = font_size * 1.2
                        total_text_height += num_lines * line_height

                    if total_text_height > box_height:
                        ratio = total_text_height / box_height
                        issues.append(ValidationIssue(
                            level=self.level,
                            rule_id=self.rule_id,
                            message=f"形状{shape_idx}文本可能溢出: 预估高度={total_text_height/914400:.1f}英寸, 容器高度={box_height/914400:.1f}英寸 ({ratio:.0%})",
                            slide_index=slide_idx
                        ))
                except Exception:
                    pass

        return issues


class SourceValidationRule(ValidationRule):
    def __init__(self):
        super().__init__("R050", "源文件检查", "warning", "检查源PPT文件的基本结构")

    def check(self, pptx_path: str) -> list[ValidationIssue]:
        issues = []
        prs = Presentation(pptx_path)

        if len(prs.slides) == 0:
            issues.append(ValidationIssue(
                level="fail",
                rule_id=self.rule_id,
                message="源文件没有任何幻灯片",
                slide_index=-1
            ))

        if len(prs.slide_masters) == 0:
            issues.append(ValidationIssue(
                level="warning",
                rule_id=self.rule_id,
                message="源文件没有母版幻灯片",
                slide_index=-1
            ))

        for slide_idx, slide in enumerate(prs.slides):
            empty = True
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    empty = False
                    break
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    empty = False
                    break
                if shape.has_table:
                    empty = False
                    break

            if empty:
                issues.append(ValidationIssue(
                    level="warning",
                    rule_id=self.rule_id,
                    message=f"第{slide_idx}页为空幻灯片",
                    slide_index=slide_idx
                ))

        return issues