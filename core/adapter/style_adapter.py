"""PPT风格贴合转换器

核心设计思想：不换模板，直接在原PPT上进行风格修改（技术贴合）

主要功能：
1. 水印检测和清理
2. 页脚处理（去除原页脚，添加模板页脚和图标）
3. 风格统一（背景色、文字颜色、强调色的视觉贴合）
4. 保留原PPT版式结构，避免版式错乱

与传统套用模板的区别：
- 传统方式：提取内容 → 创建新PPT（套用模板）→ 填入内容 → 版式容易错乱
- 技术贴合：打开原PPT → 直接修改 → 保存为新文件 → 完全保留原版式
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from core.watermark.detector import WatermarkDetector
from core.registry.template_registry import TemplateRegistry

class StyleAdapter:
    """风格贴合转换器"""
    
    def __init__(self, template_registry: TemplateRegistry = None):
        self.watermark_detector = WatermarkDetector()
        self.registry = template_registry or TemplateRegistry()
        
        # 模板配置
        self.master_style = None
        self.background_config = None
        self.text_color_rules = None
        self.footer_config = None
        self.icon_config = None
        
        # 状态变量
        self.slide_width = Emu(12192000)
        self.slide_height = Emu(6858000)
        self.header_threshold = None
        self.footer_threshold = None
    
    def configure(self, master_style: str):
        """配置转换参数"""
        self.master_style = master_style.upper()
        
        # 获取模板配置
        style_config = self.registry.get_master_style(self.master_style)
        if style_config:
            # 处理背景配置（支持background对象或background_color字符串）
            bg_config = style_config.get("background", {})
            if not bg_config and style_config.get("background_color"):
                bg_config = {
                    "type": "solid",
                    "color": style_config["background_color"]
                }
            self.background_config = bg_config
            
            self.text_color_rules = {
                "title_color": style_config.get("text_color", "FFFFFF"),
                "body_color": style_config.get("body_text_color", "E0E0E0"),
                "accent_color": style_config.get("accent_color", "3DCD58"),
            }
            self.footer_config = {
                "font_name": "Calibri",
                "font_size": Pt(10),
                "font_color": style_config.get("text_color", "FFFFFF"),
                "text": "KunPeng Testing Agent v0.1",
            }
            self.icon_config = {
                "position": {"right": Emu(914400), "bottom": Emu(914400)},
                "size": {"width": Emu(1270000), "height": Emu(508000)},
            }
    
    def adapt(self, input_path: str, output_path: str) -> str:
        """执行风格贴合转换
        
        Args:
            input_path: 原PPT文件路径
            output_path: 输出PPT文件路径
        
        Returns:
            输出文件路径
        """
        prs = Presentation(input_path)
        
        # 初始化页面尺寸参数
        self.slide_width = prs.slide_width
        self.slide_height = prs.slide_height
        self.header_threshold = self.slide_height * 0.08
        self.footer_threshold = self.slide_height * 0.88
        
        # 逐页处理
        for slide in prs.slides:
            self._process_slide(slide)
        
        prs.save(output_path)
        return output_path
    
    def _process_slide(self, slide):
        """处理单页幻灯片"""
        # 1. 删除水印
        self._remove_watermarks(slide)
        
        # 2. 删除原页脚（不添加新页脚，模板页脚由用户配置决定）
        self._remove_original_footer(slide)
        
        # 3. 修改背景色
        self._apply_background(slide)
        
        # 4. 根据位置调整文字颜色
        self._adjust_text_colors(slide)
        
        # 5. 修改强调色（边框、装饰元素）
        self._apply_accent_colors(slide)
    
    def _remove_watermarks(self, slide):
        """删除水印文字和形状"""
        indices_to_remove = []
        
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                result = self.watermark_detector.detect_text(text, 0)
                if result and result.detected:
                    indices_to_remove.append(i)
        
        for i in reversed(indices_to_remove):
            shape = slide.shapes[i]
            sp = shape._element
            sp.getparent().remove(sp)
    
    def _remove_original_footer(self, slide):
        """删除原PPT的页脚文字和署名"""
        indices_to_remove = []
        
        for i, shape in enumerate(slide.shapes):
            top = shape.top or 0
            if top > self.footer_threshold:
                is_ph = shape.is_placeholder if hasattr(shape, 'is_placeholder') else False
                if not is_ph:
                    indices_to_remove.append(i)
        
        for i in reversed(indices_to_remove):
            shape = slide.shapes[i]
            sp = shape._element
            sp.getparent().remove(sp)
    
    def _apply_background(self, slide):
        """应用背景色或渐变"""
        if not self.background_config:
            return
        
        bg_type = self.background_config.get("type", "solid")
        
        # 删除原有的背景形状（铺满整页的矩形）
        self._remove_full_page_background(slide)
        
        # 应用新背景
        if bg_type == "gradient":
            self._apply_gradient_background(slide)
        elif bg_type == "solid":
            self._apply_solid_background(slide)
        elif bg_type == "radial":
            self._apply_radial_background(slide)
    
    def _remove_full_page_background(self, slide):
        """删除铺满整页的背景矩形
        
        识别规则：
        1. 必须是AUTO_SHAPE类型
        2. 必须铺满整页（left≈0, top≈0, width≈slide_width, height≈slide_height）
        3. 必须是纯色填充
        4. 必须没有文本内容（或文本为空）
        
        这样可以避免误删有文本的内容形状
        """
        indices_to_remove = []
        
        for i, shape in enumerate(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            
            left = shape.left or Emu(0)
            top = shape.top or Emu(0)
            width = shape.width or Emu(0)
            height = shape.height or Emu(0)
            
            # 判断是否是铺满整页的背景
            is_full_page = (
                left <= Emu(1000) and
                top <= Emu(1000) and
                width >= self.slide_width - Emu(2000) and
                height >= self.slide_height - Emu(2000)
            )
            
            if not is_full_page:
                continue
            
            # 判断是否是纯色填充（背景矩形通常是纯色）
            has_solid_fill = False
            try:
                if shape.fill.type == 1:  # solid
                    has_solid_fill = True
            except:
                pass
            
            if not has_solid_fill:
                continue
            
            # 判断是否有文本内容（有文本的不是背景）
            has_text = False
            try:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        has_text = True
            except:
                pass
            
            if has_text:
                continue
            
            # 满足所有条件：是背景矩形
            indices_to_remove.append(i)
        
        for i in reversed(indices_to_remove):
            shape = slide.shapes[i]
            sp = shape._element
            sp.getparent().remove(sp)
    
    def _apply_gradient_background(self, slide):
        """应用渐变背景"""
        colors = self.background_config.get("colors", [])
        if len(colors) >= 2:
            # 创建一个铺满整页的矩形作为背景
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.oxml.xmlchemy import OxmlElement
            
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                left=Emu(0),
                top=Emu(0),
                width=self.slide_width,
                height=self.slide_height
            )
            
            # 设置渐变填充
            try:
                bg_shape.fill.gradient()
                grad = bg_shape.fill.gradient
                grad.angle = self.background_config.get("angle", 90)
                
                # 设置渐变颜色
                stops = grad.stops
                for i, color in enumerate(colors):
                    stop = stops.add_position(i / (len(colors) - 1))
                    stop.color.rgb = RGBColor.from_string(color.lstrip("#"))
            except Exception:
                # 如果渐变设置失败，使用纯色填充
                bg_shape.fill.solid()
                bg_shape.fill.fore_color.rgb = RGBColor.from_string(colors[0].lstrip("#"))
            
            # 确保背景在最底层
            bg_shape.zorder = 0
    
    def _apply_solid_background(self, slide):
        """应用纯色背景"""
        color = self.background_config.get("color", "#1F3A68")
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
    
    def _apply_radial_background(self, slide):
        """应用径向渐变背景（中心光晕效果）"""
        colors = self.background_config.get("colors", ["#3DCD58", "#0A1929"])
        
        # 创建背景矩形
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            left=Emu(0),
            top=Emu(0),
            width=self.slide_width,
            height=self.slide_height
        )
        
        try:
            bg_shape.fill.gradient()
            grad = bg_shape.fill.gradient
            grad.type = 1  # 径向渐变
            
            # 设置渐变颜色
            stops = grad.stops
            stops.add_position(0).color.rgb = RGBColor.from_string(colors[0].lstrip("#"))
            stops.add_position(1).color.rgb = RGBColor.from_string(colors[1].lstrip("#"))
            
            # 设置中心位置（偏上方）
            grad.center_left = 0.5
            grad.center_top = 0.3
        except Exception:
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor.from_string(colors[1].lstrip("#"))
        
        bg_shape.zorder = 0
    
    def _adjust_text_colors(self, slide):
        """根据文字位置调整颜色 - 直接在XML层面修改确保生效
        
        规则：
        - 文字框压在浅色卡片上 → 保持深色（或根据卡片颜色调整）
        - 文字框落在深色背景上 → 改为浅色
        """
        if not self.text_color_rules:
            return
        
        # 先找出所有有填充色的形状（卡片）
        cards = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                fill_color = None
                try:
                    if shape.fill.type == 1:  # solid fill
                        fill_color = str(shape.fill.fore_color.rgb).upper()
                except:
                    pass
                
                if fill_color and fill_color != "000000":
                    left = shape.left or Emu(0)
                    top = shape.top or Emu(0)
                    width = shape.width or Emu(0)
                    height = shape.height or Emu(0)
                    cards.append({
                        "left": left,
                        "top": top,
                        "right": left + width,
                        "bottom": top + height,
                        "fill_color": fill_color,
                        "is_light": self._is_light_color(fill_color),
                    })
        
        # 判断文字框是否落在浅色卡片上
        for shape in slide.shapes:
            if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                continue
            
            text_left = shape.left or Emu(0)
            text_top = shape.top or Emu(0)
            text_width = shape.width or Emu(0)
            text_height = shape.height or Emu(0)
            text_center_x = text_left + text_width / 2
            text_center_y = text_top + text_height / 2
            
            # 检查文字框中心是否落在某个浅色卡片上
            on_light_card = False
            for card in cards:
                if (card["left"] < text_center_x < card["right"] and
                    card["top"] < text_center_y < card["bottom"] and
                    card["is_light"]):
                    on_light_card = True
                    break
            
            # 根据位置调整文字颜色
            # 先设置shape级别的默认颜色
            tf = shape.text_frame
            self._set_xml_font_color(tf._element, "FFFFFF" if not on_light_card else "333333")
            
            for paragraph in tf.paragraphs:
                # 修改paragraph级别的默认颜色
                self._set_xml_font_color(paragraph._p, "FFFFFF" if not on_light_card else "333333")
                
                for run in paragraph.runs:
                    if run.text.strip():
                        if on_light_card:
                            # 在浅色卡片上：保持深色
                            target_color = "333333"
                        else:
                            # 在深色背景上：改为浅色
                            # 判断是否是标题（字号大或加粗）
                            font_size = run.font.size
                            is_bold = run.font.bold
                            if font_size and font_size >= Pt(24) or is_bold:
                                target_color = self.text_color_rules["title_color"]
                            else:
                                target_color = self.text_color_rules["body_color"]
                        
                        # 使用XML级别修改确保生效
                        self._set_xml_font_color(run._r, target_color)
                        # 同时用API修改
                        self._set_font_color(run.font, target_color)
    
    def _apply_accent_colors(self, slide):
        """修改强调色（边框、装饰元素）"""
        accent_color = self.text_color_rules.get("accent_color", "3DCD58")
        
        for shape in slide.shapes:
            # 修改边框颜色
            try:
                if shape.line.color is not None:
                    shape.line.color.rgb = RGBColor.from_string(accent_color)
            except:
                pass
            
            # 修改装饰性小元素的颜色
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # 判断是否是装饰性元素（小尺寸）
                width = shape.width or Emu(0)
                height = shape.height or Emu(0)
                if width < Emu(500000) and height < Emu(500000):
                    try:
                        if shape.fill.type == 1:
                            shape.fill.fore_color.rgb = RGBColor.from_string(accent_color)
                    except:
                        pass
    
    def _add_template_footer(self, slide):
        """添加模板页脚"""
        if not self.footer_config:
            return
        
        # 添加页码
        footer_text = self.footer_config.get("text", "")
        
        if footer_text:
            textbox = slide.shapes.add_textbox(
                left=Emu(4572000),
                top=self.slide_height - Emu(508000),
                width=Emu(3048000),
                height=Emu(406400)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()
            
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = footer_text
            run.font.name = self.footer_config.get("font_name", "Calibri")
            run.font.size = self.footer_config.get("font_size", Pt(10))
            self._set_font_color(run.font, self.footer_config.get("font_color", "FFFFFF"))
            para.alignment = PP_ALIGN.CENTER
    
    def _add_template_icon(self, slide):
        """添加模板图标（右下角）"""
        # 这里需要实际的图标图片路径
        # 目前先跳过，因为没有具体的图标文件
        pass
    
    def _is_light_color(self, hex_color: str) -> bool:
        """判断颜色是否为浅色"""
        if not hex_color or len(hex_color) != 6:
            return False
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
            return luminance > 0.5
        except:
            return False
    
    def _set_font_color(self, font, hex_color: str):
        """设置字体颜色"""
        try:
            font.color.rgb = RGBColor.from_string(hex_color.lstrip("#"))
        except:
            pass
    
    def _set_xml_font_color(self, element, hex_color: str):
        """在XML层面直接设置字体颜色，确保不被主题覆盖
        
        支持处理不同层级的元素：run(r), paragraph(p), text body(txBody)
        """
        try:
            from lxml import etree
            
            ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            nsmap = {'a': ns_a}
            
            tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
            
            if tag == 'r':
                # run级别：直接设置rPr
                rPr = element.find('a:rPr', nsmap)
                if rPr is None:
                    rPr = etree.SubElement(element, f'{{{ns_a}}}rPr')
                target = rPr
            elif tag == 'p':
                # paragraph级别：设置pPr/defRPr
                pPr = element.find('a:pPr', nsmap)
                if pPr is None:
                    pPr = etree.SubElement(element, f'{{{ns_a}}}pPr')
                defRPr = pPr.find('a:defRPr', nsmap)
                if defRPr is None:
                    defRPr = etree.SubElement(pPr, f'{{{ns_a}}}defRPr')
                target = defRPr
            elif tag == 'txBody':
                # text body级别：设置bodyPr/defPPr/defRPr
                defPPr = element.find('a:defPPr', nsmap)
                if defPPr is None:
                    defPPr = etree.SubElement(element, f'{{{ns_a}}}defPPr')
                defRPr = defPPr.find('a:defRPr', nsmap)
                if defRPr is None:
                    defRPr = etree.SubElement(defPPr, f'{{{ns_a}}}defRPr')
                target = defRPr
            else:
                return
            
            # 移除现有的颜色设置
            for color_tag in ['solidFill', 'noFill', 'gradFill', 'pattFill', 'blipFill']:
                existing = target.find(f'a:{color_tag}', nsmap)
                if existing is not None:
                    target.remove(existing)
            
            # 创建solidFill
            solidFill = etree.SubElement(target, f'{{{ns_a}}}solidFill')
            
            # 创建srgbClr
            srgbClr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
            srgbClr.set('val', hex_color.lstrip('#').upper())
            
        except Exception:
            pass