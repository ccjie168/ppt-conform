"""
ObjectMigrator: migrate non-placeholder shapes from source slide to target slide.
Handles pictures, tables, text boxes, auto-shapes, lines, etc.
Filters out decorative objects (old logos, watermarks, footers).

NEW FEATURES:
1. Position-based placeholder matching for free text boxes
2. Format normalization based on matched placeholder type
3. Overflow detection and adjustment
4. Full object preservation (no skipping title/subtitle)
"""

from io import BytesIO
from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor

from core.migrator.position_matcher import PositionMatcher
from core.migrator.overflow_adjuster import OverflowAdjuster
from core.migrator.coordinate_mapper import CoordinateMapper


class ObjectMigrator:
    """Migrate non-placeholder objects with position matching and overflow handling."""

    FOOTER_RATIO = 0.82
    LOGO_MAX_W, LOGO_MAX_H = 500_000, 200_000
    TINY_SIZE = 50_000

    FOOTER_KEYWORDS = ["se.com", "schneider", "copyright", "\u00a9", "page", "footer",
                       "kunpeng", "version", "v.", "rev.", "date", "draft", "confidential",
                       "slide", "number", "\u4f4d\u7f6e", "\u7b2c", "of", "\u9875", "ppt", "powerpoint"]
    WATERMARK_KEYWORDS = ["trademark", "confidential", "draft", "watermark", "trae ai", "ai generated", "ai 生成"]

    SIZE_NORMALIZATION = {
        0: 24,
        4: 16,
        2: 11,
    }

    COLOR_MAPPING = {
        "#1f3a68": "#0A2F24",
        "#000000": "#0A2F24",
        "#666666": "#718096",
        "#ffffff": "#FFFFFF",
        "#cccccc": "#A0AEC0",
        "#3dcd58": "#3DCD58",
        "#f39200": "#3DCD58",
        "#cc0000": "#DC2626",
    }

    def __init__(self, text_color="#0A2F24", bg_dark=False, position_matcher=None, overflow_adjuster=None, coordinate_mapper=None):
        self.text_color = text_color
        self.bg_dark = bg_dark
        self.position_matcher = position_matcher
        self.overflow_adjuster = overflow_adjuster
        self.coordinate_mapper = coordinate_mapper

    def migrate_objects(self, source_slide, new_slide, slide_width=None, slide_height=None) -> tuple[int, int, dict]:
        """
        Migrate all non-placeholder, non-decorative objects.
        Position-match free text boxes to target placeholders.
        Adjust overflowing objects.
        Returns (migrated_count, skipped_count, semantic_info).
        """
        migrated = 0
        skipped = 0
        semantic_info = {"title_text": "", "subtitle_text": "", "body_texts": []}

        try:
            if slide_width is None or slide_height is None:
                prs = source_slide.part.package.presentation_part.presentation
                slide_width = prs.slide_width
                slide_height = prs.slide_height
        except Exception:
            slide_width = slide_width or 12192000
            slide_height = slide_height or 6858000

        if self.overflow_adjuster is None:
            self.overflow_adjuster = OverflowAdjuster(slide_width, slide_height)

        classified_text_boxes = self._classify_text_boxes(source_slide)

        for shape in source_slide.shapes:
            if shape.is_placeholder:
                continue
            if self.is_decorative(shape, source_slide):
                skipped += 1
                continue

            try:
                if self._migrate_single(shape, new_slide, slide_width, slide_height):
                    migrated += 1
                else:
                    skipped += 1
            except Exception:
                skipped += 1

        semantic_info.update(classified_text_boxes)
        return migrated, skipped, semantic_info

    def _classify_text_boxes(self, source_slide) -> dict:
        """Classify text boxes by semantics."""
        title_text = ""
        subtitle_text = ""
        body_texts = []

        for shape in source_slide.shapes:
            if shape.is_placeholder:
                continue
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            matched_zone = self.position_matcher.match_shape(shape) if self.position_matcher else None
            ph_type = matched_zone.ph_type if matched_zone else None

            if ph_type == 0 and not title_text:
                title_text = text
            elif ph_type == 4 and not subtitle_text:
                subtitle_text = text
            elif ph_type == 2 or ph_type is None:
                body_texts.append(text)

        return {"title_text": title_text, "subtitle_text": subtitle_text, "body_texts": body_texts}

    def _migrate_single(self, shape, new_slide, slide_width, slide_height) -> bool:
        """Dispatch to type-specific migration with position matching and overflow handling."""
        st = shape.shape_type

        if st == MSO_SHAPE_TYPE.PICTURE:
            self._migrate_picture(shape, new_slide)
            self._adjust_overflow(new_slide.shapes[-1])
            return True

        if st == MSO_SHAPE_TYPE.TABLE:
            self._migrate_table(shape, new_slide)
            self._adjust_overflow(new_slide.shapes[-1])
            return True

        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            matched_zone = self.position_matcher.match_shape(shape) if self.position_matcher else None
            ph_type = matched_zone.ph_type if matched_zone else None
            
            if ph_type == 0 or ph_type == 4:
                return False
            
            self._migrate_text_box_with_position(shape, new_slide)
            self._adjust_overflow(new_slide.shapes[-1])
            return True

        if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
            self._migrate_shape_by_xml(shape, new_slide)
            self._adjust_overflow(new_slide.shapes[-1])
            return True

        if st == MSO_SHAPE_TYPE.LINE:
            self._migrate_line(shape, new_slide)
            return True

        return False

    def _adjust_overflow(self, shape):
        """Adjust shape if it overflows page boundaries."""
        if self.overflow_adjuster:
            self.overflow_adjuster.adjust_shape(shape)

    def _get_mapped_position(self, shape):
        """Get mapped position using coordinate mapper if available."""
        if self.coordinate_mapper:
            return self.coordinate_mapper.map_shape(shape, mode="fit_width")
        return shape.left, shape.top, shape.width, shape.height

    def _migrate_picture(self, shape, new_slide):
        image_stream = BytesIO(shape.image.blob)
        left, top, width, height = self._get_mapped_position(shape)
        new_slide.shapes.add_picture(
            image_stream, left, top, width, height
        )

    def _migrate_table(self, shape, new_slide):
        table = shape.table
        rows, cols = len(table.rows), len(table.columns)
        left, top, width, height = self._get_mapped_position(shape)
        new_table = new_slide.shapes.add_table(
            rows, cols, left, top, width, height
        ).table
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                new_table.cell(i, j).text = cell.text

    def _migrate_text_box_with_position(self, shape, new_slide):
        """Migrate text box with format based on position-matched placeholder."""
        matched_zone = self.position_matcher.match_shape(shape) if self.position_matcher else None
        ph_type = matched_zone.ph_type if matched_zone else None

        left, top, width, height = self._get_mapped_position(shape)

        if ph_type == 2 and matched_zone:
            left = matched_zone.left
            top = matched_zone.top
            width = matched_zone.width
            height = matched_zone.height

        txBox = new_slide.shapes.add_textbox(
            left, top, width, height
        )
        tf = txBox.text_frame
        tf.word_wrap = shape.text_frame.word_wrap

        target_size = self.SIZE_NORMALIZATION.get(ph_type, 11)

        for i, para in enumerate(shape.text_frame.paragraphs):
            new_para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            new_para.text = para.text
            new_para.level = para.level
            new_para.font.size = Pt(target_size)
            new_para.font.name = "Poppins"
            try:
                new_para.font.color.rgb = RGBColor.from_string(self.text_color.lstrip("#"))
            except Exception:
                pass

    def _migrate_shape_by_xml(self, shape, new_slide):
        new_sp = deepcopy(shape._element)
        nsmap = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }
        nvSpPr = new_sp.find(".//p:nvSpPr", nsmap)
        if nvSpPr is not None:
            ph = nvSpPr.find(".//p:ph", nsmap)
            if ph is not None:
                nvSpPr.remove(ph)

        if self.coordinate_mapper:
            left, top, width, height = self.coordinate_mapper.map_shape(shape, mode="fit_width")
            xfrm = new_sp.find(".//a:xfrm", nsmap)
            if xfrm is not None:
                off = xfrm.find(".//a:off", nsmap)
                if off is not None:
                    off.set("x", str(int(left)))
                    off.set("y", str(int(top)))
                ext = xfrm.find(".//a:ext", nsmap)
                if ext is not None:
                    ext.set("cx", str(int(width)))
                    ext.set("cy", str(int(height)))

        new_slide.shapes._spTree.insert_element_before(new_sp, "p:extLst")

    def _migrate_line(self, shape, new_slide):
        from pptx.enum.shapes import MSO_SHAPE
        left, top, width, height = self._get_mapped_position(shape)
        new_slide.shapes.add_shape(
            MSO_SHAPE.LINE_INVERSE, left, top, width, height
        )

    def is_decorative(self, shape, source_slide) -> bool:
        """True if shape is an old logo, watermark, footer, etc."""
        try:
            prs = source_slide.part.package.presentation_part.presentation
            slide_height = prs.slide_height
        except Exception:
            slide_height = 6858000

        if shape.top > slide_height * self.FOOTER_RATIO:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if shape.width < self.LOGO_MAX_W and shape.height < self.LOGO_MAX_H:
                    return True
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                if any(kw in text for kw in self.FOOTER_KEYWORDS):
                    return True
                text_length = len(text.strip())
                if text_length > 0 and text_length < 80:
                    return True
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                return True

        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if any(kw in text for kw in self.WATERMARK_KEYWORDS):
                return True

        try:
            if shape.width < self.TINY_SIZE and shape.height < self.TINY_SIZE:
                if not (shape.has_text_frame and shape.text_frame.text.strip()):
                    return True
        except Exception:
            pass

        return False
