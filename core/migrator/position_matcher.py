"""
PositionMatcher: match source text boxes to target template placeholders based on position.

Analyzes target template placeholder positions, then matches source text boxes
by their position on the slide to determine which placeholder format to apply.
"""

from pptx.enum.shapes import MSO_SHAPE_TYPE


class PlaceholderZone:
    """Defines a placeholder zone with position and type."""
    def __init__(self, name, ph_type, left, top, width, height):
        self.name = name
        self.ph_type = ph_type
        self.left = left
        self.top = top
        self.width = width
        self.height = height

    def contains_point(self, x, y, tolerance=0.15) -> bool:
        """Check if point is within zone (with tolerance)."""
        expanded_left = self.left - self.width * tolerance
        expanded_right = self.left + self.width * (1 + tolerance)
        expanded_top = self.top - self.height * tolerance
        expanded_bottom = self.top + self.height * (1 + tolerance)
        return (expanded_left <= x <= expanded_right and
                expanded_top <= y <= expanded_bottom)

    def overlaps_shape(self, shape, overlap_ratio=0.3) -> bool:
        """Check if shape overlaps this zone by at least overlap_ratio."""
        shape_left = shape.left
        shape_top = shape.top
        shape_right = shape_left + shape.width
        shape_bottom = shape_top + shape.height

        zone_right = self.left + self.width
        zone_bottom = self.top + self.height

        overlap_left = max(shape_left, self.left)
        overlap_top = max(shape_top, self.top)
        overlap_right = min(shape_right, zone_right)
        overlap_bottom = min(shape_bottom, zone_bottom)

        if overlap_right <= overlap_left or overlap_bottom <= overlap_top:
            return False

        overlap_area = (overlap_right - overlap_left) * (overlap_bottom - overlap_top)
        shape_area = shape.width * shape.height

        return overlap_area / max(shape_area, 1) >= overlap_ratio


class PositionMatcher:
    """Match source text boxes to target placeholder zones based on position."""

    def __init__(self, target_prs, master_index=0):
        self.zones = []
        self.slide_width = 0
        self.slide_height = 0
        self._analyze_template(target_prs, master_index)

    def _analyze_template(self, target_prs, master_index):
        """Analyze target template to extract placeholder zones."""
        self.slide_width = target_prs.slide_width
        self.slide_height = target_prs.slide_height

        if master_index >= len(target_prs.slide_masters):
            master = target_prs.slide_masters[0]
        else:
            master = target_prs.slide_masters[master_index]

        for layout in master.slide_layouts:
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type == 0:  # TITLE
                    self.zones.append(PlaceholderZone(
                        "Title", ph_type, ph.left, ph.top, ph.width, ph.height
                    ))
                elif ph_type == 4:  # SUBTITLE
                    self.zones.append(PlaceholderZone(
                        "Subtitle", ph_type, ph.left, ph.top, ph.width, ph.height
                    ))
                elif ph_type == 2:  # BODY
                    self.zones.append(PlaceholderZone(
                        "Body", ph_type, ph.left, ph.top, ph.width, ph.height
                    ))

        if not self.zones:
            self._create_default_zones()

    def _create_default_zones(self):
        """Create default zones based on typical slide layout."""
        w, h = self.slide_width, self.slide_height
        self.zones = [
            PlaceholderZone("Title", 0, 0, 0, w, h * 0.15),
            PlaceholderZone("Subtitle", 4, 0, h * 0.15, w, h * 0.1),
            PlaceholderZone("Body", 2, w * 0.05, h * 0.28, w * 0.9, h * 0.6),
        ]

    def match_shape(self, shape) -> PlaceholderZone | None:
        """Match a shape to the best matching placeholder zone."""
        if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE,
                               MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.MEDIA):
            return None

        shape_center_x = shape.left + shape.width / 2
        shape_center_y = shape.top + shape.height / 2

        for zone in self.zones:
            if zone.overlaps_shape(shape):
                return zone

        if shape.top < self.slide_height * 0.18:
            for zone in self.zones:
                if zone.ph_type == 0:
                    return zone
        elif shape.top < self.slide_height * 0.3:
            for zone in self.zones:
                if zone.ph_type == 4:
                    return zone

        return None

    def get_zone_by_name(self, name) -> PlaceholderZone | None:
        """Get zone by name."""
        for zone in self.zones:
            if zone.name.lower() == name.lower():
                return zone
        return None
