from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from core.clrmap.resolver import ClrMapResolver
from core.migrator.object_migrator import ObjectMigrator


class SlideMigrator:
    """Content Migration: create new slide from template layout, migrate text and objects."""

    def __init__(self, template_path: str, master_index: int):
        self.template_path = template_path
        self.master_index = master_index
        self.clr_resolver = ClrMapResolver(template_path, master_index)

    def migrate_slide(self, source_slide, target_prs, slide_type: str, layout_index: int):
        """
        Create a new slide based on specified layout, migrate objects with position matching.
        Returns the new slide object.
        """
        from core.migrator.position_matcher import PositionMatcher
        from core.migrator.overflow_adjuster import OverflowAdjuster
        from core.migrator.coordinate_mapper import CoordinateMapper

        if self.master_index >= len(target_prs.slide_masters):
            master = target_prs.slide_masters[0]
        else:
            master = target_prs.slide_masters[self.master_index]

        if layout_index >= len(master.slide_layouts):
            layout_index = 0
        target_layout = master.slide_layouts[layout_index]
        new_slide = target_prs.slides.add_slide(target_layout)

        self._clear_placeholder_defaults(new_slide)

        bg_dark = self.master_index == 2
        text_color = "#FFFFFF" if bg_dark else "#0A2F24"

        tgt_width = target_prs.slide_width
        tgt_height = target_prs.slide_height

        try:
            src_prs = source_slide.part.package.presentation_part.presentation
            src_width = src_prs.slide_width
            src_height = src_prs.slide_height
        except Exception:
            src_width = tgt_width
            src_height = tgt_height

        position_matcher = PositionMatcher(target_prs, self.master_index)
        overflow_adjuster = OverflowAdjuster(tgt_width, tgt_height)
        coordinate_mapper = CoordinateMapper(src_width, src_height, tgt_width, tgt_height)

        object_migrator = ObjectMigrator(
            text_color=text_color,
            bg_dark=bg_dark,
            position_matcher=position_matcher,
            overflow_adjuster=overflow_adjuster,
            coordinate_mapper=coordinate_mapper
        )
        migrated, skipped, semantic_info = object_migrator.migrate_objects(source_slide, new_slide, tgt_width, tgt_height)

        title_text = semantic_info.get("title_text", "") or self._extract_title(source_slide)
        subtitle_text = semantic_info.get("subtitle_text", "") or self._extract_subtitle(source_slide)

        if new_slide.shapes.title and title_text:
            new_slide.shapes.title.text = title_text
        for ph in new_slide.placeholders:
            ph_type = ph.placeholder_format.type
            if ph_type == 4 and ph.has_text_frame and subtitle_text:
                ph.text_frame.text = subtitle_text

        return new_slide

    def _clear_placeholder_defaults(self, slide):
        """Clear default placeholder text like 'Click to edit master title style'."""
        for ph in slide.placeholders:
            if ph.has_text_frame:
                tf = ph.text_frame
                for para in tf.paragraphs:
                    para.text = ""

    def _extract_title(self, slide) -> str:
        if slide.shapes.title:
            return slide.shapes.title.text or ""
        # Fallback: find largest text shape
        max_shape = None
        max_area = 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                area = shape.width * shape.height
                if area > max_area:
                    max_area = area
                    max_shape = shape
        if max_shape:
            return max_shape.text_frame.text.strip()
        return ""

    def _extract_subtitle(self, slide) -> str:
        # Try subtitle placeholder (type 4 in python-pptx)
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            # PP_PLACEHOLDER.SUBTITLE = 4
            if ph_type == 4 and shape.has_text_frame:
                return shape.text_frame.text.strip()
        # Fallback: second largest text shape
        text_shapes = [
            s for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        text_shapes.sort(key=lambda s: s.width * s.height, reverse=True)
        if len(text_shapes) >= 2:
            return text_shapes[1].text_frame.text.strip()
        return ""

    def _extract_body_paragraphs(self, slide) -> list[str]:
        paragraphs = []
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            if title_shape and shape == title_shape:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragraphs.append(text)
        return paragraphs

    def _fill_title(self, slide, text: str):
        if slide.shapes.title and text:
            slide.shapes.title.text = text
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Poppins"

    def _fill_subtitle(self, slide, text: str):
        if not text:
            return
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            if ph_type == 4 and shape.has_text_frame:
                shape.text_frame.text = text
                return
        # No subtitle placeholder, try body placeholder
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            if ph_type == 2 and shape.has_text_frame:  # BODY = 2
                shape.text_frame.text = text
                return

    def _fill_body(self, slide, paragraphs: list[str]):
        if not paragraphs:
            return
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            if ph_type == 2 and shape.has_text_frame:  # BODY = 2
                tf = shape.text_frame
                tf.clear()
                for i, para_text in enumerate(paragraphs):
                    if i == 0:
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()
                    para.text = para_text
                    for run in para.runs:
                        run.font.name = "Poppins"
                return
