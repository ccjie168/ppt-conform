import yaml
import re
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from core.models import SlideClassification


CONFIG_PATH = Path(__file__).parent.parent.parent / "config" / "layout_mapping.yaml"

SUMMARY_KEYWORDS = [
    "executive summary", "key messages", "key takeaways",
    "总结", "摘要", "概述", "overview", "agenda",
]
RISK_KEYWORDS = [
    "risk", "issue", "mitigation", "挑战", "风险", "问题", "应对",
]
TIMELINE_KEYWORDS = [
    "timeline", "milestone", "roadmap", "时间线", "里程碑", "路线图", "进度",
]
CLOSING_KEYWORDS = [
    "thank you", "thanks", "se.com", "closing", "q&a",
    "谢谢", "感谢", "结束",
]


class SlideClassifier:
    def __init__(self, config_path: str | None = None):
        self.layout_mappings: dict = {}
        self.migration_types: set = set()
        self.adaptation_types: set = set()
        self._load_config(config_path)

    def _load_config(self, config_path: str | None):
        path = Path(config_path) if config_path else CONFIG_PATH
        if path.exists():
            with open(path, "r", encoding="utf-8") as f:
                config = yaml.safe_load(f)
            self.layout_mappings = config.get("layout_mappings", {})
            self.migration_types = set(config.get("migration_types", []))
            self.adaptation_types = set(config.get("adaptation_types", []))

    def classify_all(self, presentation) -> list[SlideClassification]:
        if isinstance(presentation, str):
            prs = Presentation(presentation)
        else:
            prs = presentation

        results = []
        total = len(prs.slides)
        for idx, slide in enumerate(prs.slides):
            cls = self.classify_slide(slide, idx, total)
            results.append(cls)
        return results

    def classify_slide(self, slide, slide_index: int, total_slides: int) -> SlideClassification:
        slide_type = self._detect_type(slide, slide_index, total_slides)
        migration_mode = self._decide_mode(slide_type)
        layout_info = self.layout_mappings.get(slide_type, {
            "layout_index": 1, "layout_name": "One column"
        })

        return SlideClassification(
            slide_index=slide_index,
            slide_type=slide_type,
            migration_mode=migration_mode,
            target_layout_index=layout_info.get("layout_index", 1),
            target_layout_name=layout_info.get("layout_name", "One column"),
            confidence=0.7,
        )

    def _detect_type(self, slide, slide_index: int, total_slides: int) -> str:
        title_text = self._get_title_text(slide)
        all_text = self._get_all_text(slide).lower()

        # 1. Cover
        if slide_index == 0 and self._is_cover_like(slide):
            return "Cover"

        # 2. Closing
        if slide_index == total_slides - 1:
            if any(kw in all_text for kw in CLOSING_KEYWORDS):
                return "Closing"

        # 3. Section Divider
        if self._is_section_divider(slide):
            return "Section_Divider"

        # 4. Executive Summary
        if any(kw in title_text.lower() for kw in SUMMARY_KEYWORDS):
            return "Executive_Summary"

        # 5. Risk
        if any(kw in title_text.lower() for kw in RISK_KEYWORDS):
            return "Risk_Matrix"

        # 6. Timeline
        if any(kw in title_text.lower() for kw in TIMELINE_KEYWORDS):
            return "Timeline"

        # 7. Table
        if self._has_table(slide):
            return "Table_Page"

        # 8. Chart
        chart_count = self._count_charts(slide)
        if chart_count >= 2:
            return "KPI_Dashboard"
        if chart_count == 1:
            return "Chart_Page"

        # 9. Image dominant
        if self._is_image_dominant(slide):
            return "Image_Page"

        return "Content"

    def _decide_mode(self, slide_type: str) -> str:
        if slide_type in self.migration_types:
            return "migration"
        return "adaptation"

    def _get_title_text(self, slide) -> str:
        if slide.shapes.title:
            return slide.shapes.title.text or ""
        return ""

    def _get_all_text(self, slide) -> str:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        return " ".join(texts)

    def _is_cover_like(self, slide) -> bool:
        # 封面页不应包含表格、图表等复杂对象
        for shape in slide.shapes:
            if shape.shape_type in (MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART,
                                    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT):
                return False
        # 检查非标题文本框的内容量：如果有较多正文内容，不是封面
        total_body_chars = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            total_body_chars += len(shape.text_frame.text.strip())
        # 正文内容超过 30 字符，不是封面
        if total_body_chars > 30:
            return False
        return True

    def _is_section_divider(self, slide) -> bool:
        title = self._get_title_text(slide)
        if not title:
            return False
        # 有表格/图表/媒体等复杂对象，不是章节分隔页
        for shape in slide.shapes:
            if shape.shape_type in (MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART,
                                    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
                                    MSO_SHAPE_TYPE.MEDIA):
                return False
        body_shapes = [
            s for s in slide.shapes
            if s.has_text_frame and s != slide.shapes.title
        ]
        if len(body_shapes) == 0:
            return True
        total_body_chars = sum(
            len(s.text_frame.text) for s in body_shapes if s.has_text_frame
        )
        if total_body_chars < 20:
            return True
        return False

    def _has_table(self, slide) -> bool:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return True
        return False

    def _count_charts(self, slide) -> int:
        count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                count += 1
        return count

    def _is_image_dominant(self, slide) -> bool:
        image_count = 0
        text_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            if shape.has_text_frame and len(shape.text_frame.text) > 10:
                text_count += 1
        return image_count >= 1 and text_count <= 1
