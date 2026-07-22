import zipfile
from lxml import etree
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from core.models import PreCheckResult, PreCheckIssue

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

SE_KEYWORDS = ["schneider", "se.com", "Schneider Electric", "施耐德"]


class PreCheckAnalyzer:
    def analyze(self, pptx_path: str) -> PreCheckResult:
        result = PreCheckResult()
        prs = Presentation(pptx_path)

        # 1. 页面数量
        result.slide_count = len(prs.slides)

        # 2. 母版数量
        result.master_count = len(prs.slide_masters)
        if result.master_count > 3:
            result.issues.append(PreCheckIssue(
                level="warning",
                rule_id="too_many_masters",
                message=f"母版数量为 {result.master_count}，可能存在模板污染",
            ))

        # 3. 尺寸比例
        ratio = prs.slide_width / prs.slide_height if prs.slide_height else 0
        result.is_4_3_ratio = abs(ratio - 4/3) < 0.02
        if result.is_4_3_ratio:
            result.issues.append(PreCheckIssue(
                level="warning",
                rule_id="4_3_ratio",
                message="PPT 为 4:3 比例，将自动调整为 16:9",
            ))

        # 4. 字体清单
        fonts = set()
        for slide in prs.slides:
            self._collect_fonts_from_slide(slide, fonts)
        result.fonts_used = sorted(fonts)

        # 检查旧字体
        old_fonts = [f for f in fonts if any(
            kw in f.lower() for kw in ["microsoft yahei", "calibri", "微软雅黑"]
        )]
        if old_fonts:
            result.issues.append(PreCheckIssue(
                level="info",
                rule_id="old_fonts",
                message=f"检测到旧字体: {', '.join(old_fonts)}，将自动替换",
            ))

        # 5. 其他检查
        for slide in prs.slides:
            self._check_slide_media(slide, result)
            self._check_slide_overflow(slide, prs, result)

        # 6. 检查动画
        result.has_animation = self._check_animation(pptx_path)

        # 7. 检查 SmartArt
        result.has_smartart = self._check_smartart(pptx_path)

        # 8. 检查是否含旧 SE 模板（通过主题/母版名称）
        result.has_old_se_template = self._check_old_se_template(pptx_path)

        # 9. 外部主题
        result.has_external_theme = result.master_count > 2

        return result

    def _collect_fonts_from_slide(self, slide, fonts: set):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    fonts.add(run.font.name)

    def _check_slide_media(self, slide, result):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                result.has_media = True
            if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                result.has_embedded_chart = True

    def _check_slide_overflow(self, slide, prs, result):
        for shape in slide.shapes:
            try:
                if shape.left + shape.width > prs.slide_width + 10000:
                    result.overflow_objects_count += 1
                if shape.top + shape.height > prs.slide_height + 10000:
                    result.overflow_objects_count += 1
            except Exception:
                pass

    def _check_animation(self, pptx_path: str) -> bool:
        try:
            with zipfile.ZipFile(pptx_path, "r") as zf:
                for name in zf.namelist():
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                        xml = zf.read(name)
                        if b"<p:timing" in xml or b"<p:anim" in xml:
                            return True
        except Exception:
            pass
        return False

    def _check_smartart(self, pptx_path: str) -> bool:
        try:
            with zipfile.ZipFile(pptx_path, "r") as zf:
                for name in zf.namelist():
                    if "diagrams" in name or "smartArt" in name.lower():
                        return True
        except Exception:
            pass
        return False

    def _check_old_se_template(self, pptx_path: str) -> bool:
        try:
            with zipfile.ZipFile(pptx_path, "r") as zf:
                for name in zf.namelist():
                    if name.startswith("ppt/theme/theme") and name.endswith(".xml"):
                        xml = zf.read(name).decode("utf-8", errors="ignore")
                        if any(kw.lower() in xml.lower() for kw in SE_KEYWORDS):
                            return True
        except Exception:
            pass
        return False
