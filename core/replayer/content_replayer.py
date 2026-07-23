from pathlib import Path
from io import BytesIO
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from core.models import SlideContentModel, UserConfig
from core.registry.template_registry import TemplateRegistry
from core.analyzer.template_format_extractor import TemplateFormatExtractor


class ContentReplayer:
    """内容重放器：将抽取的内容重放到目标模板，标题/正文字体遵循模板要求，保留源内容格式"""

    def __init__(self, registry: TemplateRegistry, template_path: str | None = None):
        self.registry = registry
        self.template_path = template_path
        self.template_formats: dict = {}
        self.theme_fonts: dict = {}
        self.default_text_color: str | None = None
        self.title_text_color: str | None = None
        self.template_width: int = 0
        self.template_height: int = 0
        self.footer_shapes: list[dict] = []
        self.background_image: dict | None = None
        self.background_color: str | None = None
        self.background_theme_color: str | None = None
        self.placeholder_mapping: dict = {}  # 模板占位符语义映射
        self.color_mapping: dict = {}  # 颜色映射表：原颜色 -> 模板颜色
        self.template_accent_color: str | None = None  # 模板强调色
        self.template_title_font: str | None = None  # 模板标题字体 (major)
        self.template_body_font: str | None = None  # 模板正文字体 (minor)
        self.force_template_font: bool = True  # 是否强制使用模板字体

        if template_path and Path(template_path).exists():
            extractor = TemplateFormatExtractor()
            try:
                self.template_formats = extractor.extract_placeholder_formats(template_path)
                self.theme_fonts = extractor.extract_theme_fonts(template_path)
                self._analyze_template_footer(template_path)
            except Exception:
                pass

    def _analyze_template_footer(self, template_path: str, master_index: int = 0) -> None:
        """分析指定Master中的footer元素（图标、页脚文本等）和背景信息，用于复制到新slide"""
        self.footer_shapes = []
        self.master_brand_images = []
        self.background_image = None
        self.background_color = None
        self.background_theme_color = None
        try:
            from pptx import Presentation
            from lxml import etree

            prs = Presentation(template_path)
            self.template_width = prs.slide_width
            self.template_height = prs.slide_height

            if not prs.slide_masters or master_index >= len(prs.slide_masters):
                return

            master = prs.slide_masters[master_index]
            footer_threshold = self.template_height * 0.85

            nsmap = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            }

            # 获取母版关联的主题文件路径
            import zipfile
            self.template_theme_path = None
            try:
                master_idx = master_index + 1
                master_rels_path = f'ppt/slideMasters/_rels/slideMaster{master_idx}.xml.rels'
                with zipfile.ZipFile(template_path, 'r') as zf:
                    if master_rels_path in zf.namelist():
                        master_rels_xml = zf.read(master_rels_path)
                        rels_elem = etree.fromstring(master_rels_xml)
                        ns_rels = 'http://schemas.openxmlformats.org/package/2006/relationships'
                        for rel in rels_elem.findall(f'.//{{{ns_rels}}}Relationship'):
                            if 'theme' in rel.get('Type', ''):
                                self.template_theme_path = rel.get('Target', '')
                                break
            except Exception:
                pass

            # 分析Master背景色
            bg = master._element.find('.//p:bg', nsmap)
            if bg is not None:
                bgPr = bg.find('.//p:bgPr', nsmap)
                if bgPr is not None:
                    solidFill = bgPr.find('.//a:solidFill', nsmap)
                    if solidFill is not None:
                        schemeClr = solidFill.find('.//a:schemeClr', nsmap)
                        if schemeClr is not None:
                            self.background_theme_color = schemeClr.get("val")
                        srgbClr = solidFill.find('.//a:srgbClr', nsmap)
                        if srgbClr is not None:
                            self.background_color = srgbClr.get("val")
                
                # 尝试从渐变背景获取颜色
                gradFill = bgPr.find('.//a:gradFill', nsmap)
                if gradFill is not None and not self.background_color:
                    stop_colors = []
                    for stop in gradFill.findall('.//a:stop', nsmap):
                        srgb = stop.find('.//a:srgbClr', nsmap)
                        if srgb is not None:
                            stop_colors.append(srgb.get("val"))
                        scheme = stop.find('.//a:schemeClr', nsmap)
                        if scheme is not None:
                            stop_colors.append(scheme.get("val"))
                    if stop_colors:
                        self.background_color = stop_colors[0]
                        self.background_theme_color = stop_colors[0] if any(c in ['bg1', 'bg2', 'dk1', 'dk2', 'lt1', 'lt2'] for c in stop_colors) else None
            
            # 如果只有主题颜色，从主题文件中提取实际RGB值
            if self.background_theme_color and not self.background_color:
                self.background_color = self._get_color_from_theme(self.background_theme_color)

            for shape in master.shapes:
                try:
                    is_footer = False
                    shape_type = shape.shape_type

                    # 检查是否是背景图（铺满整个页面的图片）
                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            left = shape.left or 0
                            top = shape.top or 0
                            width = shape.width or 0
                            height = shape.height or 0
                            is_full_bg = (abs(left) < Emu(20000) and abs(top) < Emu(20000) and
                                          abs(width - self.template_width) < Emu(50000) and
                                          abs(height - self.template_height) < Emu(50000))
                            if is_full_bg:
                                # 保存背景图
                                self.background_image = {
                                    "image_blob": shape.image.blob,
                                    "image_ext": shape.image.ext,
                                    "left": left,
                                    "top": top,
                                    "width": width,
                                    "height": height,
                                }
                                continue
                        except Exception:
                            pass

                    # 底部区域的图片（施耐德图标等）- 需要从母版删除，但保留在footer_shapes中以便重新添加
                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if shape.top and shape.top > footer_threshold:
                            is_footer = True
                            # 记录母版中的品牌图片，后续需要删除
                            self.master_brand_images.append({
                                "name": shape.name,
                                "left": shape.left,
                                "top": shape.top,
                                "width": shape.width,
                                "height": shape.height,
                            })

                    # 底部区域的文本框（页脚文本等）
                    if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        if shape.top and shape.top > footer_threshold:
                            is_footer = True

                    # 底部区域的占位符（页脚、页码等）
                    if shape.is_placeholder:
                        try:
                            ph_type = int(shape.placeholder_format.type)
                            # 13=SLIDE_NUMBER, 14=HEADER, 15=FOOTER, 16=DATE
                            if ph_type in (13, 14, 15, 16) and shape.top and shape.top > footer_threshold:
                                is_footer = True
                        except Exception:
                            pass

                    if is_footer:
                        right_margin = self.template_width - shape.left - shape.width
                        bottom_margin = self.template_height - shape.top - shape.height

                        shape_data = {
                            "type": shape_type,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                            "right_margin": right_margin,
                            "bottom_margin": bottom_margin,
                        }

                        # 图片：保存图片二进制数据
                        if shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                shape_data["image_blob"] = shape.image.blob
                                shape_data["image_ext"] = shape.image.ext
                            except Exception:
                                continue

                        # 文本框：保存文本内容
                        if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                            try:
                                shape_data["text"] = shape.text_frame.text
                                # 保存字体样式
                                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                                    run = shape.text_frame.paragraphs[0].runs[0]
                                    shape_data["font_name"] = run.font.name
                                    shape_data["font_size"] = run.font.size
                                    shape_data["font_bold"] = run.font.bold
                                    try:
                                        if run.font.color and run.font.color.rgb:
                                            shape_data["font_color"] = run.font.color.rgb
                                    except (AttributeError, TypeError):
                                        pass
                            except Exception:
                                shape_data["text"] = ""

                        # 占位符（页脚、页码等）：保存文本内容和类型
                        if shape.is_placeholder:
                            try:
                                shape_data["text"] = shape.text_frame.text
                                shape_data["ph_type"] = int(shape.placeholder_format.type)
                                # 保存字体样式
                                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                                    run = shape.text_frame.paragraphs[0].runs[0]
                                    shape_data["font_name"] = run.font.name
                                    shape_data["font_size"] = run.font.size
                                    shape_data["font_bold"] = run.font.bold
                                    try:
                                        if run.font.color and run.font.color.rgb:
                                            shape_data["font_color"] = run.font.color.rgb
                                    except (AttributeError, TypeError):
                                        pass
                            except Exception:
                                shape_data["text"] = ""

                        self.footer_shapes.append(shape_data)
                except Exception:
                    continue
        except Exception:
            pass

    # 标准16:9尺寸 (13.333 x 7.5 英寸)
    WIDESCREEN_16_9 = Emu(12192000)
    WIDESCREEN_16_9_H = Emu(6858000)

    def _remove_master_brand_images(self, prs) -> None:
        """从母版中删除品牌图片，避免与代码添加的页脚图标重复"""
        if not self.master_brand_images:
            return

        try:
            for master in prs.slide_masters:
                shapes_to_remove = []

                for brand_img in self.master_brand_images:
                    brand_name = brand_img.get("name", "")
                    brand_left = brand_img.get("left")
                    brand_top = brand_img.get("top")
                    brand_width = brand_img.get("width")
                    brand_height = brand_img.get("height")

                    for shape in master.shapes:
                        try:
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                if brand_name and shape.name == brand_name:
                                    shapes_to_remove.append(shape)
                                    break
                                if (brand_left is not None and brand_top is not None and
                                    brand_width is not None and brand_height is not None):
                                    if (abs(shape.left - brand_left) < Emu(1000) and
                                        abs(shape.top - brand_top) < Emu(1000) and
                                        abs(shape.width - brand_width) < Emu(1000) and
                                        abs(shape.height - brand_height) < Emu(1000)):
                                        shapes_to_remove.append(shape)
                                        break
                        except Exception:
                            continue

                for shape in shapes_to_remove:
                    try:
                        sp = shape._element
                        sp.getparent().remove(sp)
                    except Exception:
                        continue

                # 同时检查所有幻灯片布局
                for layout in master.slide_layouts:
                    layout_shapes_to_remove = []
                    for brand_img in self.master_brand_images:
                        brand_name = brand_img.get("name", "")
                        brand_left = brand_img.get("left")
                        brand_top = brand_img.get("top")
                        brand_width = brand_img.get("width")
                        brand_height = brand_img.get("height")

                        for shape in layout.shapes:
                            try:
                                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    if brand_name and shape.name == brand_name:
                                        layout_shapes_to_remove.append(shape)
                                        break
                                    if (brand_left is not None and brand_top is not None and
                                        brand_width is not None and brand_height is not None):
                                        if (abs(shape.left - brand_left) < Emu(1000) and
                                            abs(shape.top - brand_top) < Emu(1000) and
                                            abs(shape.width - brand_width) < Emu(1000) and
                                            abs(shape.height - brand_height) < Emu(1000)):
                                            layout_shapes_to_remove.append(shape)
                                            break
                            except Exception:
                                continue

                    for shape in layout_shapes_to_remove:
                        try:
                            sp = shape._element
                            sp.getparent().remove(sp)
                        except Exception:
                            continue
        except Exception:
            pass

    def _adapt_master_background_size(self, prs) -> None:
        """适配母版背景尺寸，确保在页面尺寸变更后背景能铺满整个页面
        
        当强制设置16:9页面尺寸后，母版中的背景形状（矩形、图片等）尺寸不会自动更新，
        导致背景只占部分区域。此方法遍历所有母版和版式，更新背景形状尺寸。
        """
        target_width = prs.slide_width
        target_height = prs.slide_height

        try:
            for master in prs.slide_masters:
                for shape in master.shapes:
                    try:
                        left = shape.left or 0
                        top = shape.top or 0
                        width = shape.width or 0
                        height = shape.height or 0

                        is_full_bg = (
                            abs(left) < Emu(20000) and
                            abs(top) < Emu(20000) and
                            width > 0 and height > 0
                        )

                        if is_full_bg:
                            if abs(width - target_width) > Emu(50000) or abs(height - target_height) > Emu(50000):
                                shape.width = target_width
                                shape.height = target_height
                                shape.left = Emu(0)
                                shape.top = Emu(0)
                    except Exception:
                        continue

                for layout in master.slide_layouts:
                    for shape in layout.shapes:
                        try:
                            left = shape.left or 0
                            top = shape.top or 0
                            width = shape.width or 0
                            height = shape.height or 0

                            is_full_bg = (
                                abs(left) < Emu(20000) and
                                abs(top) < Emu(20000) and
                                width > 0 and height > 0
                            )

                            if is_full_bg:
                                if abs(width - target_width) > Emu(50000) or abs(height - target_height) > Emu(50000):
                                    shape.width = target_width
                                    shape.height = target_height
                                    shape.left = Emu(0)
                                    shape.top = Emu(0)
                        except Exception:
                            continue
        except Exception:
            pass

    def replay(self, content_models: list[SlideContentModel], config: UserConfig) -> str:
        source_prs = Presentation(config.input_path)
        output_prs = Presentation(config.input_path)

        selected_master_index = int(config.master_style) if config.master_style.isdigit() else 0

        # 从目标模板只提取选中母版的样式信息（技术适配模式）
        self._extract_selected_master_styles(selected_master_index)

        # 从目标模板复制选中的母版和版式到输出PPT（保留原母版）
        output_prs = self._copy_selected_master_to_output(output_prs, selected_master_index)

        # 删除原PPT母版中的页脚元素，避免继承到幻灯片
        self._remove_master_footer_elements(output_prs)

        self.source_width = source_prs.slide_width
        self.source_height = source_prs.slide_height
        self.target_width = output_prs.slide_width
        self.target_height = output_prs.slide_height
        self.scale_x = 1.0
        self.scale_y = 1.0

        try:
            master_style = config.master_style.upper()
            master_config = self.registry.get_master_style(master_style)
            if master_config:
                self.template_accent_color = master_config.get("accent_color", "").lstrip("#")
            self.color_mapping = self._build_color_mapping(content_models, self.template_accent_color)
        except Exception:
            self.color_mapping = {}

        for slide_idx, slide in enumerate(output_prs.slides):
            self._apply_style_adaptation(slide, slide_idx)

        output_prs.save(config.output_path)
        return config.output_path

    def _copy_selected_master_to_output(self, output_prs, selected_master_index: int):
        """从目标模板只复制选中的母版和版式到输出PPT，保留原母版

        使用zipfile完整复制母版相关的所有文件，确保PPT结构完整不损坏。
        同时修改输出PPT尺寸为模板尺寸，并将所有幻灯片指向新母版的版式。
        返回新的Presentation对象。
        """
        if not self.template_path or not Path(self.template_path).exists():
            return output_prs

        try:
            import zipfile
            import tempfile
            from lxml import etree

            template_prs = Presentation(self.template_path)
            if selected_master_index >= len(template_prs.slide_masters):
                return output_prs

            template_master = template_prs.slide_masters[selected_master_index]
            template_master_idx = selected_master_index + 1
            template_width = template_prs.slide_width
            template_height = template_prs.slide_height

            temp_path = '/tmp/tmp_output_master_src.pptx'
            output_prs.save(temp_path)

            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                tmp_path = tmp.name

            ns_ct = 'http://schemas.openxmlformats.org/package/2006/content-types'
            ns_rels = 'http://schemas.openxmlformats.org/package/2006/relationships'
            ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

            with zipfile.ZipFile(self.template_path, 'r') as tpl_zip, \
                 zipfile.ZipFile(temp_path, 'r') as src_zip, \
                 zipfile.ZipFile(tmp_path, 'w') as dst_zip:

                src_files = set(src_zip.namelist())

                existing_layouts = len([f for f in src_files if 'ppt/slideLayouts/slideLayout' in f and '_rels' not in f])
                existing_masters = len([f for f in src_files if 'ppt/slideMasters/slideMaster' in f and '_rels' not in f])
                existing_themes = len([f for f in src_files if 'ppt/theme/theme' in f and '_rels' not in f])

                new_master_idx = existing_masters + 1

                tpl_master_xml = tpl_zip.read(f'ppt/slideMasters/slideMaster{template_master_idx}.xml')
                tpl_master_rels = tpl_zip.read(f'ppt/slideMasters/_rels/slideMaster{template_master_idx}.xml.rels')

                master_rels_elem = etree.fromstring(tpl_master_rels)
                layout_rels = master_rels_elem.findall(f'{{{ns_rels}}}Relationship[@Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"]')
                theme_rels = master_rels_elem.findall(f'{{{ns_rels}}}Relationship[@Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"]')

                layout_mapping = {}
                new_layout_idx = existing_layouts + 1
                for rel in layout_rels:
                    target = rel.get('Target')
                    layout_num = target.split('slideLayout')[-1].split('.')[0]
                    layout_mapping[layout_num] = new_layout_idx
                    new_layout_idx += 1

                theme_mapping = {}
                new_theme_idx = existing_themes + 1
                for rel in theme_rels:
                    target = rel.get('Target')
                    theme_num = target.split('theme')[-1].split('.')[0]
                    theme_mapping[theme_num] = new_theme_idx
                    new_theme_idx += 1

                files_to_copy = {}

                new_master_filename = f'ppt/slideMasters/slideMaster{new_master_idx}.xml'
                files_to_copy[new_master_filename] = tpl_master_xml

                new_master_rels_filename = f'ppt/slideMasters/_rels/slideMaster{new_master_idx}.xml.rels'
                files_to_copy[new_master_rels_filename] = tpl_master_rels

                for old_layout_num, new_layout_num in layout_mapping.items():
                    old_layout_path = f'ppt/slideLayouts/slideLayout{old_layout_num}.xml'
                    new_layout_path = f'ppt/slideLayouts/slideLayout{new_layout_num}.xml'
                    if old_layout_path in tpl_zip.namelist():
                        files_to_copy[new_layout_path] = tpl_zip.read(old_layout_path)

                    old_layout_rels_path = f'ppt/slideLayouts/_rels/slideLayout{old_layout_num}.xml.rels'
                    new_layout_rels_path = f'ppt/slideLayouts/_rels/slideLayout{new_layout_num}.xml.rels'
                    if old_layout_rels_path in tpl_zip.namelist():
                        layout_rels_xml = tpl_zip.read(old_layout_rels_path)
                        layout_rels_elem = etree.fromstring(layout_rels_xml)
                        for rel in layout_rels_elem.findall(f'{{{ns_rels}}}Relationship'):
                            target = rel.get('Target', '')
                            if 'theme' in target:
                                old_theme_num = target.split('theme')[-1].split('.')[0]
                                if old_theme_num in theme_mapping:
                                    rel.set('Target', f'../theme/theme{theme_mapping[old_theme_num]}.xml')
                            elif 'slideMaster' in target:
                                rel.set('Target', f'../slideMasters/slideMaster{new_master_idx}.xml')
                        files_to_copy[new_layout_rels_path] = etree.tostring(layout_rels_elem, xml_declaration=True, encoding='UTF-8', standalone=True)

                for old_theme_num, new_theme_num in theme_mapping.items():
                    old_theme_path = f'ppt/theme/theme{old_theme_num}.xml'
                    new_theme_path = f'ppt/theme/theme{new_theme_num}.xml'
                    if old_theme_path in tpl_zip.namelist():
                        files_to_copy[new_theme_path] = tpl_zip.read(old_theme_path)

                master_rels_new = etree.fromstring(tpl_master_rels)
                for rel in master_rels_new.findall(f'{{{ns_rels}}}Relationship'):
                    target = rel.get('Target', '')

                    if 'slideLayout' in target:
                        old_num = target.split('slideLayout')[-1].split('.')[0]
                        if old_num in layout_mapping:
                            rel.set('Target', f'../slideLayouts/slideLayout{layout_mapping[old_num]}.xml')
                    elif 'theme' in target:
                        old_num = target.split('theme')[-1].split('.')[0]
                        if old_num in theme_mapping:
                            rel.set('Target', f'../theme/theme{theme_mapping[old_num]}.xml')

                files_to_copy[new_master_rels_filename] = etree.tostring(master_rels_new, xml_declaration=True, encoding='UTF-8', standalone=True)

                # 新母版的第一个版式索引，用于后续幻灯片指向
                first_new_layout_num = min(layout_mapping.values()) if layout_mapping else (existing_layouts + 1)

                for item in src_zip.infolist():
                    data = src_zip.read(item.filename)

                    if item.filename == '[Content_Types].xml':
                        ct_elem = etree.fromstring(data)

                        for old_layout_num, new_layout_num in layout_mapping.items():
                            override = etree.SubElement(ct_elem, f'{{{ns_ct}}}Override')
                            override.set('PartName', f'/ppt/slideLayouts/slideLayout{new_layout_num}.xml')
                            override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml')

                        for old_theme_num, new_theme_num in theme_mapping.items():
                            override = etree.SubElement(ct_elem, f'{{{ns_ct}}}Override')
                            override.set('PartName', f'/ppt/theme/theme{new_theme_num}.xml')
                            override.set('ContentType', 'application/vnd.openxmlformats-officedocument.theme+xml')

                        override = etree.SubElement(ct_elem, f'{{{ns_ct}}}Override')
                        override.set('PartName', f'/ppt/slideMasters/slideMaster{new_master_idx}.xml')
                        override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml')

                        data = etree.tostring(ct_elem, xml_declaration=True, encoding='UTF-8', standalone=True)

                    elif item.filename == 'ppt/presentation.xml':
                        pres_elem = etree.fromstring(data)
                        sld_master_id_lst = pres_elem.find(f'.//{{{ns_p}}}sldMasterIdLst')
                        if sld_master_id_lst is None:
                            sld_master_id_lst = etree.SubElement(
                                pres_elem,
                                f'{{{ns_p}}}sldMasterIdLst'
                            )

                        new_master_id = etree.SubElement(
                            sld_master_id_lst,
                            f'{{{ns_p}}}sldMasterId'
                        )
                        new_master_id.set(f'{{{ns_r}}}id', f'rIdMaster{new_master_idx}')

                        # 修改页面尺寸为模板尺寸
                        sldSz = pres_elem.find(f'.//{{{ns_p}}}sldSz')
                        if sldSz is not None:
                            sldSz.set('cx', str(template_width))
                            sldSz.set('cy', str(template_height))
                            sldSz.set('type', 'screen16x9')

                        data = etree.tostring(pres_elem, xml_declaration=True, encoding='UTF-8', standalone=True)

                    elif item.filename == 'ppt/_rels/presentation.xml.rels':
                        pres_rels_elem = etree.fromstring(data)

                        new_rel = etree.SubElement(pres_rels_elem, f'{{{ns_rels}}}Relationship')
                        new_rel.set('Id', f'rIdMaster{new_master_idx}')
                        new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster')
                        new_rel.set('Target', f'slideMasters/slideMaster{new_master_idx}.xml')

                        data = etree.tostring(pres_rels_elem, xml_declaration=True, encoding='UTF-8', standalone=True)

                    elif item.filename.startswith('ppt/slides/_rels/slide') and item.filename.endswith('.xml.rels'):
                        # 将所有幻灯片的版式引用指向新母版的第一个版式
                        slide_rels_elem = etree.fromstring(data)
                        modified = False
                        for rel in slide_rels_elem.findall(f'{{{ns_rels}}}Relationship'):
                            if 'slideLayout' in rel.get('Type', ''):
                                rel.set('Target', f'../slideLayouts/slideLayout{first_new_layout_num}.xml')
                                modified = True
                        if modified:
                            data = etree.tostring(slide_rels_elem, xml_declaration=True, encoding='UTF-8', standalone=True)

                    dst_zip.writestr(item.filename, data)

                for filename, content in files_to_copy.items():
                    if filename not in src_files:
                        dst_zip.writestr(filename, content)

            return Presentation(tmp_path)

        except Exception:
            return output_prs

    def _resize_master_backgrounds(self, prs) -> None:
        """调整母版背景尺寸，确保与页面尺寸匹配"""
        page_width = prs.slide_width
        page_height = prs.slide_height

        for master in prs.slide_masters:
            try:
                for shape in master.shapes:
                    try:
                        left = shape.left or Emu(0)
                        top = shape.top or Emu(0)

                        is_full_page = (
                            left <= Emu(1000) and
                            top <= Emu(1000)
                        )

                        if is_full_page:
                            shape.width = page_width
                            shape.height = page_height
                    except Exception:
                        continue

                for layout in master.slide_layouts:
                    for layout_shape in layout.shapes:
                        try:
                            layout_left = layout_shape.left or Emu(0)
                            layout_top = layout_shape.top or Emu(0)

                            is_full_page = (
                                layout_left <= Emu(1000) and
                                layout_top <= Emu(1000)
                            )

                            if is_full_page:
                                layout_shape.width = page_width
                                layout_shape.height = page_height
                        except Exception:
                            continue
            except Exception:
                continue

    def _extract_selected_master_styles(self, selected_master_index: int) -> None:
        """从目标模板只提取选中母版的样式信息，不引入其他模板"""
        if not self.template_path or not Path(self.template_path).exists():
            return

        try:
            extractor = TemplateFormatExtractor()
            self.default_text_color = extractor.get_text_color_for_master(
                self.template_path, selected_master_index
            )

            master_style = str(selected_master_index).upper()
            self.master_style = master_style

            # 将数字索引映射到 F1-F4 风格ID
            style_id_map = {"0": "F1", "1": "F2", "2": "F3", "3": "F4"}
            style_id = style_id_map.get(master_style, master_style)

            colors = extractor.extract_theme_colors(self.template_path)

            if master_style in ("F1", "F2", "0", "1"):
                self.title_text_color = colors.get("dk2", "3DCD58")
                if not self.default_text_color or self._is_light_color(self.default_text_color):
                    self.default_text_color = "333333"
            else:
                self.title_text_color = colors.get("lt1", "FFFFFF")
                self.default_text_color = "FFFFFF"

            if self.default_text_color:
                if "body" not in self.template_formats:
                    self.template_formats["body"] = {}
                if not self.template_formats["body"].get("color"):
                    self.template_formats["body"]["color"] = self.default_text_color
            if self.title_text_color:
                if "title" not in self.template_formats:
                    self.template_formats["title"] = {}
                self.template_formats["title"]["color"] = self.title_text_color

            self.theme_fonts = extractor.extract_theme_fonts(
                self.template_path, selected_master_index
            )

            self._analyze_template_footer(self.template_path, selected_master_index)

            # 根据 registry 配置强制覆盖背景色，确保与 UI 选择的风格一致
            # （模板母版实际背景可能与风格定义不符，如 F3 深绿色在模板中可能对应白色母版）
            try:
                master_config = self.registry.get_master_style(style_id)
                if master_config:
                    bg_config = master_config.get("background")
                    if bg_config and isinstance(bg_config, dict):
                        # 渐变背景
                        if bg_config.get("type") == "gradient":
                            self.background_color = "gradient"
                            self.background_gradient = {
                                "start_color": bg_config.get("start_color", "#1A237E").lstrip("#"),
                                "end_color": bg_config.get("end_color", "#0D47A1").lstrip("#"),
                                "direction": bg_config.get("direction", "vertical"),
                            }
                        else:
                            self.background_color = bg_config.get("color", self.background_color)
                            self.background_gradient = None
                    else:
                        bg_color = master_config.get("background_color")
                        if bg_color:
                            self.background_color = bg_color.lstrip("#")
                        self.background_gradient = None
            except Exception:
                pass

            self.placeholder_mapping = extractor.extract_placeholder_mapping(
                self.template_path, selected_master_index
            )

            fonts = extractor.extract_theme_fonts(self.template_path, selected_master_index)
            self.template_title_font = fonts.get("major")
            self.template_body_font = fonts.get("minor")

        except Exception:
            pass

    def _apply_style_adaptation(self, slide, slide_idx: int) -> None:
        """对单页幻灯片应用技术适配样式调整
        
        技术适配原则：
        1. 有占位符定义的组件（标题、副标题、正文等）→ 比照目标模板的占位符做转换（字体和颜色）
        2. 没有占位符定义的组件（text box, shape等）→ 保留，分两种情况：
           a) 本身有背景颜色的 → 字体转变成模板字体，颜色不变
           b) 本身没有背景颜色的 → 字体和颜色要符合目标模板的color pairing
        3. 识别为页脚的占位符 → 删除，并完全应用模板的页脚定义
        """
        # 1. 删除水印
        self._remove_watermarks_from_slide(slide)

        # 2. 删除原PPT页脚（占位符+普通文本框/图片），然后应用模板页脚
        self._remove_original_footer(slide)
        self._apply_template_footer(slide)
        
        # 3. 应用背景（基于目标模板）
        self._apply_background_to_slide(slide)
        
        # 4. 按规则转换组件样式
        self._convert_components_style(slide)
        
        # 5. 检查并修复溢出
        self._check_and_fix_overflow(slide)

    def _is_footer_shape(self, shape, footer_threshold) -> bool:
        """判断一个形状是否是页脚元素（位于底部区域）
        
        更精确的判断逻辑：
        1. 页脚占位符（13=SLIDE_NUMBER, 14=HEADER, 15=FOOTER, 16=DATE）总是删除
        2. 底部区域的小文本框（高度<100pt）可能是页脚
        3. 底部区域的图片可能是品牌图标
        4. 排除大尺寸的形状（可能是正常内容）
        """
        try:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type in (13, 14, 15, 16):
                    return True
            
            top = shape.top or Emu(0)
            height = shape.height or Emu(0)
            bottom = top + height

            if top > footer_threshold or bottom > footer_threshold:
                # 小尺寸元素更可能是页脚
                if height < Emu(360000):  # < 100pt
                    if shape.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PICTURE):
                        return True
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                            text = shape.text_frame.text.strip()
                            if text and len(text) < 50:
                                return True
            return False
        except Exception:
            return False

    def _collect_footer_indices(self, shapes, footer_threshold) -> list:
        """收集底部区域所有需要删除的形状索引（递归处理group）"""
        indices_to_remove = []

        for i, shape in enumerate(shapes):
            try:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type in (13, 14, 15, 16):
                        indices_to_remove.append(i)
                        continue

                if self._is_footer_shape(shape, footer_threshold):
                    indices_to_remove.append(i)
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        group_top = shape.top or Emu(0)
                        group_height = shape.height or Emu(0)
                        group_bottom = group_top + group_height
                        if group_top > footer_threshold or group_bottom > footer_threshold:
                            indices_to_remove.append(i)
                    except Exception:
                        pass
            except Exception:
                continue

        return indices_to_remove

    def _remove_original_footer(self, slide) -> None:
        """删除原PPT的页脚元素（包括占位符、文本框、图片、自选图形、组）

        删除底部区域（>80%页面高度）的所有页脚相关元素：
        - 页脚占位符（13=SLIDE_NUMBER, 14=HEADER, 15=FOOTER, 16=DATE）
        - 底部区域的文本框（如版权信息）
        - 底部区域的图片（如图标）
        - 底部区域带文本的自选图形
        - 底部区域的组形状
        """
        footer_threshold = self.target_height * 0.80

        indices_to_remove = self._collect_footer_indices(slide.shapes, footer_threshold)

        for i in reversed(indices_to_remove):
            try:
                shape = slide.shapes[i]
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception:
                continue

    def _remove_master_footer_elements(self, prs) -> None:
        """删除母版中的页脚元素，确保原PPT的页脚不影响输出"""
        page_height = prs.slide_height
        footer_threshold = page_height * 0.80

        for master in prs.slide_masters:
            indices_to_remove = self._collect_footer_indices(master.shapes, footer_threshold)

            for i in reversed(indices_to_remove):
                try:
                    shape = master.shapes[i]
                    sp = shape._element
                    sp.getparent().remove(sp)
                except Exception:
                    continue

    def _apply_template_footer(self, slide) -> None:
        """完全应用模板的页脚定义，保持相对边距（右下角对齐）"""
        if not self.footer_shapes:
            return

        try:
            for footer_info in self.footer_shapes:
                try:
                    from io import BytesIO

                    shape_type = footer_info.get("type")
                    width = footer_info.get("width", Emu(3048000))
                    height = footer_info.get("height", Emu(406400))

                    # 根据 right_margin 和 bottom_margin 重新计算位置
                    # 保持相对边距，确保图标/文本始终在右下角
                    right_margin = footer_info.get("right_margin")
                    bottom_margin = footer_info.get("bottom_margin")

                    if right_margin is not None and bottom_margin is not None:
                        left = self.target_width - width - right_margin
                        top = self.target_height - height - bottom_margin
                    else:
                        left = footer_info.get("left", Emu(0))
                        top = footer_info.get("top", Emu(0))

                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_blob = footer_info.get("image_blob")
                        if image_blob:
                            stream = BytesIO(image_blob)
                            slide.shapes.add_picture(
                                stream,
                                left=int(left),
                                top=int(top),
                                width=int(width),
                                height=int(height),
                            )
                    else:
                        textbox = slide.shapes.add_textbox(
                            int(left), int(top), int(width), int(height)
                        )
                        text = footer_info.get("text", "")
                        
                        # 强制使用 Poppins 6号字（Pt(6)）作为页脚字体
                        font_name = "Poppins"
                        font_size = Pt(6)
                        font_color = footer_info.get("font_color")
                        font_bold = footer_info.get("font_bold")

                        if text:
                            for paragraph in textbox.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.text = ""

                            p = textbox.text_frame.paragraphs[0]
                            run = p.add_run()
                            run.text = text

                            run.font.name = font_name
                            run.font.size = font_size
                            if font_bold is not None:
                                run.font.bold = font_bold
                            if font_color:
                                run.font.color.rgb = font_color
                        else:
                            p = textbox.text_frame.paragraphs[0]
                            run = p.add_run()
                            run.text = ""
                            run.font.name = font_name
                            run.font.size = font_size
                            if font_bold is not None:
                                run.font.bold = font_bold
                            if font_color:
                                run.font.color.rgb = font_color
                except Exception:
                    continue
        except Exception:
            pass

    def _convert_components_style(self, slide) -> None:
        """按规则转换幻灯片上所有组件的样式
        
        规则：
        1. 有占位符定义的组件（标题、副标题、正文等）→ 完全应用目标模板样式
        2. 页脚位置内的元素（占位符和非占位符）→ 应用模板页脚定义的颜色
        3. 没有占位符定义的组件（text box, shape等）→ 字体改模板字体，颜色保持原样
        """
        if not self.template_title_font and not self.template_body_font:
            return
        
        title_font = self.template_title_font or "Arial"
        body_font = self.template_body_font or "Arial"
        
        bg_is_dark = self._detect_slide_background_darkness(slide)
        if bg_is_dark:
            template_title_color = self.title_text_color or "FFFFFF"
            template_body_color = self.default_text_color or "E8F5E9"
        else:
            template_title_color = self.title_text_color or "3DCD58"
            template_body_color = self.default_text_color or "555555"
        
        footer_threshold = self.target_height * 0.85
        
        for shape in slide.shapes:
            if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                continue
            
            is_footer_region = self._is_shape_in_footer_region(shape, footer_threshold)
            is_placeholder = self._is_content_placeholder(shape)
            
            if is_placeholder:
                # 规则1：有占位符定义的组件 → 完全应用目标模板样式（字体和颜色）
                self._convert_placeholder_style(shape, title_font, body_font, template_title_color, template_body_color)
            elif is_footer_region:
                # 规则2：页脚位置内的元素 → 应用模板定义的颜色
                self._convert_footer_style(shape, body_font, template_body_color)
            else:
                # 规则3：没有占位符定义的组件 → 字体改模板字体，颜色保持原样
                has_bg_color = self._shape_has_background_color(shape)
                if has_bg_color:
                    self._convert_non_placeholder_with_bg(shape, title_font, body_font, template_title_color, template_body_color)
                else:
                    self._convert_non_placeholder_no_bg(shape, title_font, body_font, template_title_color, template_body_color, bg_is_dark)

    def _is_shape_in_footer_region(self, shape, footer_threshold) -> bool:
        """判断形状是否位于页脚区域（底部15%）"""
        try:
            top = shape.top or Emu(0)
            height = shape.height or Emu(0)
            bottom = top + height
            return top > footer_threshold or bottom > footer_threshold
        except Exception:
            return False

    def _convert_footer_style(self, shape, font: str, color: str) -> None:
        """页脚位置内的元素：应用模板定义的字体和颜色"""
        try:
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        run.font.name = font
                        run.font._element.set('eastAsian', font)
                        try:
                            run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
                        except Exception:
                            pass
        except Exception:
            pass

    def _is_content_placeholder(self, shape) -> bool:
        """判断是否是内容占位符（标题、副标题、正文等）
        
        页眉/页脚/页码/日期占位符不算内容占位符
        """
        header_footer_types = (13, 14, 15, 16)  # SLIDE_NUMBER, HEADER, FOOTER, DATE
        try:
            if not shape.is_placeholder:
                return False
            phf = shape.placeholder_format
            if phf.type in header_footer_types:
                return False
            # 1=TITLE, 3=CENTER_TITLE, 4=SUBTITLE, 2=BODY
            return phf.type in (1, 2, 3, 4, 7, 8, 9, 10, 18)
        except Exception:
            return False

    def _convert_placeholder_style(self, shape, title_font: str, body_font: str,
                                     title_color: str, body_color: str) -> None:
        """有占位符定义的组件：应用目标模板样式
        
        规则：
        - 标题：字体改模板字体，颜色从模板主题色中提取（浅色背景用dk2深色，深色背景用lt1浅色）
        - 正文：字体改模板字体，颜色改为模板定义的颜色
        """
        try:
            phf = shape.placeholder_format
            is_title = phf.type in (1, 3)  # TITLE or CENTER_TITLE
            target_font = title_font if is_title else body_font
            target_color = title_color if is_title else body_color
            
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        run.font.name = target_font
                        run.font._element.set('eastAsian', target_font)
                        try:
                            run.font.color.rgb = RGBColor.from_string(target_color.lstrip("#"))
                        except Exception:
                            pass
        except Exception:
            pass

    def _convert_non_placeholder_with_bg(self, shape, title_font: str, body_font: str,
                                          title_color: str, body_color: str) -> None:
        """无占位符但有背景颜色的组件：字体转模板字体，颜色保持原样"""
        try:
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        font_size = run.font.size
                        is_bold = run.font.bold
                        is_title = font_size and font_size >= Pt(24) or is_bold
                        target_font = title_font if is_title else body_font
                        
                        run.font.name = target_font
                        run.font._element.set('eastAsian', target_font)
                        # 颜色保持原样
        except Exception:
            pass

    def _detect_shape_background_darkness(self, shape) -> bool:
        """检测形状自身背景色是否为深色"""
        try:
            fill_type = shape.fill.type
            if fill_type == 1:  # solid fill
                fill_color = str(shape.fill.fore_color.rgb).upper()
                if fill_color and fill_color != "NONE":
                    return self._is_dark_color(fill_color)
            elif fill_type == 2:  # gradient
                try:
                    colors = [str(s.color.rgb).upper() for s in shape.fill.gradient.stops]
                    dark_count = sum(1 for c in colors if c and c != "NONE" and self._is_dark_color(c))
                    return dark_count >= len(colors) / 2
                except Exception:
                    pass
        except Exception:
            pass
        return True

    def _color_has_enough_contrast_with_shape(self, hex_color: str, shape) -> bool:
        """判断文字颜色与形状自身背景色是否有足够对比度"""
        try:
            bg_is_dark = self._detect_shape_background_darkness(shape)
            return self._color_has_enough_contrast(hex_color, bg_is_dark)
        except Exception:
            return False

    def _convert_non_placeholder_no_bg(self, shape, title_font: str, body_font: str,
                                          title_color: str, body_color: str,
                                          bg_is_dark: bool = True) -> None:
        """无占位符且无背景颜色的组件：字体改模板字体，颜色保持原样"""
        try:
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        font_size = run.font.size
                        is_bold = run.font.bold
                        is_title = font_size and font_size >= Pt(24) or is_bold
                        
                        target_font = title_font if is_title else body_font
                        
                        run.font.name = target_font
                        run.font._element.set('eastAsian', target_font)
                        # 颜色保持原样
        except Exception:
            pass

    def _color_has_enough_contrast(self, hex_color: str, bg_is_dark: bool) -> bool:
        """判断文字颜色与背景是否有足够对比度（不违反color pairing）
        
        简化判断：
        - 深色背景 + 浅色文字（亮度>=0.5）：对比度足够
        - 浅色背景 + 深色文字（亮度<0.5）：对比度足够
        """
        if not hex_color or len(hex_color) != 6:
            return False
        try:
            color_is_dark = self._is_dark_color(hex_color)
            if bg_is_dark:
                return not color_is_dark
            else:
                return color_is_dark
        except Exception:
            return False

    def _shape_has_background_color(self, shape) -> bool:
        """判断形状是否有背景颜色（用于text box/shape处理）"""
        try:
            fill_type = shape.fill.type
            if fill_type == 1:  # solid fill
                fill_color = str(shape.fill.fore_color.rgb).upper()
                if fill_color not in ["000000", "NONE", None]:
                    return True
            elif fill_type == 2:  # gradient
                return True
        except Exception:
            pass
        return False

    def _detect_slide_background_darkness(self, slide) -> bool:
        """检测幻灯片背景是否为深色
        
        注意：不能访问 slide.background.fill，因为这会自动添加背景元素并中断母版继承
        而是从模板母版的背景颜色来判断
        """
        try:
            if self.background_color:
                return self._is_dark_color(self.background_color)
        except Exception:
            pass
        
        try:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    left = shape.left or Emu(0)
                    top = shape.top or Emu(0)
                    width = shape.width or Emu(0)
                    height = shape.height or Emu(0)
                    
                    is_full_page = (
                        left <= Emu(1000) and
                        top <= Emu(1000) and
                        width >= self.target_width - Emu(2000) and
                        height >= self.target_height - Emu(2000)
                    )
                    
                    if is_full_page:
                        try:
                            if shape.fill.type == 1:  # solid
                                color = str(shape.fill.fore_color.rgb).upper()
                                return self._is_dark_color(color)
                            elif shape.fill.type == 2:  # gradient
                                colors = [str(s.color.rgb).upper() for s in shape.fill.gradient.stops]
                                dark_count = sum(1 for c in colors if self._is_dark_color(c))
                                return dark_count >= len(colors) / 2
                        except Exception:
                            pass
        except Exception:
            pass
        
        return True

    def _is_dark_color(self, hex_color: str) -> bool:
        """判断颜色是否为深色"""
        if not hex_color or len(hex_color) != 6:
            return False
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
            return luminance < 0.5
        except Exception:
            return False

    def _adapt_position(self, value, axis: str = "x") -> int:
        """根据源/目标尺寸比例适配位置或大小"""
        if value is None:
            return None
        scale = self.scale_x if axis == "x" else self.scale_y
        if abs(scale - 1.0) < 0.01:
            return value
        return int(value * scale)

    def _clear_slides(self, prs) -> None:
        sldIdLst = prs.slides._sldIdLst
        rIds = [sldId.rId for sldId in list(sldIdLst)]
        for rId in rIds:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
        for child in list(sldIdLst):
            sldIdLst.remove(child)

    def _clear_placeholders(self, slide) -> None:
        """清空幻灯片上的内容占位符的文字内容，保留占位符的样式（位置、字体、颜色）
        
        - 页眉/页脚/页码占位符保留不动
        - 标题/正文/副标题占位符只清空文字，保留样式
        - 非占位符形状都删掉（背景图和footer后面手动加）
        """
        header_footer_types = (13, 14, 15, 16)  # slide_number, header, footer, date

        shapes_to_remove = []
        for shape in slide.shapes:
            try:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type in header_footer_types:
                        # 页眉页脚页码占位符保留不动
                        continue
                    # 标题/正文/副标题占位符：只清空文字内容，保留样式
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ""
                else:
                    # 非占位符形状都删掉（背景图和footer后面手动加）
                    shapes_to_remove.append(shape)
            except Exception:
                shapes_to_remove.append(shape)

        spTree = slide.shapes._spTree
        for shape in shapes_to_remove:
            try:
                spTree.remove(shape._element)
            except Exception:
                pass

    def _add_background_image(self, slide) -> None:
        """添加背景：优先使用背景图，否则使用纯色背景（与模板Master一致）"""
        if self.background_image:
            try:
                from io import BytesIO

                target_width = self.target_width
                target_height = self.target_height

                image_blob = self.background_image["image_blob"]
                stream = BytesIO(image_blob)

                pic = slide.shapes.add_picture(
                    stream,
                    left=0,
                    top=0,
                    width=target_width,
                    height=target_height,
                )

                spTree = slide.shapes._spTree
                spTree.remove(pic._element)
                spTree.insert(2, pic._element)
            except Exception:
                pass
        elif self.background_theme_color:
            try:
                from pptx.enum.dml import MSO_THEME_COLOR
                theme_color_map = {
                    'bg1': MSO_THEME_COLOR.BACKGROUND_1,
                    'bg2': MSO_THEME_COLOR.BACKGROUND_2,
                    'tx1': MSO_THEME_COLOR.TEXT_1,
                    'tx2': MSO_THEME_COLOR.TEXT_2,
                    'accent1': MSO_THEME_COLOR.ACCENT_1,
                    'accent2': MSO_THEME_COLOR.ACCENT_2,
                    'accent3': MSO_THEME_COLOR.ACCENT_3,
                    'accent4': MSO_THEME_COLOR.ACCENT_4,
                    'accent5': MSO_THEME_COLOR.ACCENT_5,
                    'accent6': MSO_THEME_COLOR.ACCENT_6,
                    'lt1': MSO_THEME_COLOR.LIGHT_1,
                    'lt2': MSO_THEME_COLOR.LIGHT_2,
                    'dk1': MSO_THEME_COLOR.DARK_1,
                    'dk2': MSO_THEME_COLOR.DARK_2,
                }
                mso_color = theme_color_map.get(self.background_theme_color.lower())
                if mso_color is not None:
                    bg = slide.background
                    fill = bg.fill
                    fill.solid()
                    fill.fore_color.theme_color = mso_color
            except Exception:
                pass
        elif self.background_color:
            try:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                from pptx.dml.color import RGBColor
                fill.fore_color.rgb = RGBColor.from_string(self.background_color)
            except Exception:
                pass

    def _add_footer_shapes(self, slide) -> None:
        """将footer元素（施耐德图标、页脚文本等）添加到slide上，并适配16:9位置"""
        if not self.footer_shapes:
            return

        try:
            target_width = self.target_width
            target_height = self.target_height
        except Exception:
            target_width = Emu(12192000)
            target_height = Emu(6858000)

        # 判断是否需要适配（尺寸不同时按比例缩放位置）
        needs_adapt = (self.template_width > 0 and self.template_height > 0 and
                       (target_width != self.template_width or target_height != self.template_height))

        for footer_info in self.footer_shapes:
            try:
                from io import BytesIO

                shape_type = footer_info["type"]
                shape_w = footer_info["width"]
                shape_h = footer_info["height"]
                orig_right = footer_info["right_margin"]
                orig_bottom = footer_info["bottom_margin"]

                if needs_adapt:
                    # 计算缩放比例
                    scale_x = target_width / self.template_width
                    scale_y = target_height / self.template_height
                    # 保持右边距和下边距的比例
                    new_right = orig_right * scale_x
                    new_bottom = orig_bottom * scale_y
                    new_w = shape_w * scale_x
                    new_h = shape_h * scale_y
                    new_left = target_width - new_w - new_right
                    new_top = target_height - new_h - new_bottom
                else:
                    new_left = footer_info["left"]
                    new_top = footer_info["top"]
                    new_w = shape_w
                    new_h = shape_h

                # 确保位置合理
                if new_left < 0:
                    new_left = 0
                if new_top < 0:
                    new_top = 0

                if shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # 添加图片
                    image_blob = footer_info.get("image_blob")
                    if image_blob:
                        stream = BytesIO(image_blob)
                        slide.shapes.add_picture(
                            stream,
                            left=int(new_left),
                            top=int(new_top),
                            width=int(new_w),
                            height=int(new_h),
                        )
                else:
                    # 文本框或占位符（页脚、页码等），都用文本框添加
                    textbox = slide.shapes.add_textbox(
                        int(new_left), int(new_top), int(new_w), int(new_h)
                    )
                    text = footer_info.get("text", "")
                    if text:
                        textbox.text_frame.text = text
                        # 应用字体样式：所有页脚文本强制使用 Poppins 6号字
                        try:
                            font_name = "Poppins"
                            font_size = Pt(6)
                            font_bold = footer_info.get("font_bold")
                            font_color = footer_info.get("font_color")
                            
                            para = textbox.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                run.font.name = font_name
                                run.font.size = font_size
                                if font_bold is not None:
                                    run.font.bold = font_bold
                                if font_color:
                                    run.font.color.rgb = font_color
                        except Exception:
                            pass
            except Exception:
                continue

    def _build_layout_map_from_master(self, master) -> dict:
        """从Master中构建布局名称到索引的映射
        
        智能识别各种布局类型，优先选择最常用的内容页布局。
        """
        mapping = {}
        layout_names = []
        
        for idx, layout in enumerate(master.slide_layouts):
            name = (layout.name or "").lower()
            layout_names.append((idx, name, layout))
            
            # 直接用名称映射
            mapping[name] = idx
        
        # 按优先级匹配各种布局类型
        # 封面布局
        cover_priorities = [
            "title slide simple", "title slide", "cover", "封面",
            "title slide with image", "title slide offer",
        ]
        for kw in cover_priorities:
            for idx, name, layout in layout_names:
                if kw in name and "cover" not in mapping:
                    mapping["cover"] = idx
                    break
            if "cover" in mapping:
                break
        
        # 章节分隔页
        section_priorities = [
            "section break for longer copy", "section break with subhead",
            "section break", "section", "章节",
        ]
        for kw in section_priorities:
            for idx, name, layout in layout_names:
                if kw in name and "section" not in mapping:
                    mapping["section"] = idx
                    break
            if "section" in mapping:
                break
        
        # 结尾页
        closing_priorities = [
            "closing slide", "closing", "end", "结尾", "结束", "thank you",
        ]
        for kw in closing_priorities:
            for idx, name, layout in layout_names:
                if kw in name and "closing" not in mapping:
                    mapping["closing"] = idx
                    break
            if "closing" in mapping:
                break
        
        # 议程页
        agenda_priorities = ["agenda", "议程"]
        for kw in agenda_priorities:
            for idx, name, layout in layout_names:
                if kw in name and "agenda" not in mapping:
                    mapping["agenda"] = idx
                    break
            if "agenda" in mapping:
                break
        
        # 内容页布局（最常用的标准内容页）
        # 优先级：one column > blank slide with title > 其他单列布局
        content_priorities = [
            "one column", "one column two line headline",
            "one column with sketches", "blank slide with title",
            "content", "正文",
        ]
        for kw in content_priorities:
            for idx, name, layout in layout_names:
                if kw in name and "content" not in mapping:
                    mapping["content"] = idx
                    break
            if "content" in mapping:
                break
        
        # 双列布局
        two_col_priorities = [
            "two columns", "two column", "2 column",
            "two columns with subtitle",
        ]
        for kw in two_col_priorities:
            for idx, name, layout in layout_names:
                if kw in name and "two_column" not in mapping:
                    mapping["two_column"] = idx
                    break
            if "two_column" in mapping:
                break
        
        # 三列布局
        three_col_priorities = ["three columns", "three column", "3 column"]
        for kw in three_col_priorities:
            for idx, name, layout in layout_names:
                if kw in name and "three_column" not in mapping:
                    mapping["three_column"] = idx
                    break
            if "three_column" in mapping:
                break
        
        # 表格页（使用内容页布局即可）
        if "table" not in mapping and "content" in mapping:
            mapping["table"] = mapping["content"]
        
        # 图片页
        if "image" not in mapping and "content" in mapping:
            mapping["image"] = mapping["content"]
        
        return mapping

    def _resolve_layout_index(self, layout_name: str, master, master_style: str) -> int:
        layout_map = self._build_layout_map_from_master(master)
        if layout_name in layout_map:
            return layout_map[layout_name]

        layout_info = self.registry.get_layout_by_name(master_style, layout_name)
        if layout_info:
            idx = layout_info.get("index", 0)
            if idx < len(list(master.slide_layouts)):
                return idx

        if layout_name == "cover":
            return 0
        return min(1, len(list(master.slide_layouts)) - 1) if master else 1

    def _copy_slide_content(self, slide, model: SlideContentModel, slide_idx: int) -> None:
        """将内容回填到模板slide中，模板格式优先，原格式兜底
        
        核心原则（解决内容与格式割裂）：
        1. 标题填入模板的title占位符（保留模板样式）
        2. 副标题填入模板的subtitle占位符
        3. 主正文填入模板的body占位符（如果有且只有一个主正文形状）
        4. 其他内容作为额外形状按原位置回填，保留原PPT布局结构
        5. 用 raw_shape_id 精确去重，避免内容重复处理
        """
        # 已处理的 raw_shape_id 集合（用于精确去重）
        processed_shape_ids = set()
        
        # 如果没有占位符标题，从 body_blocks 中找 semantic_role="title" 的 block 作为标题
        title_text = model.title
        title_format = model.title_format
        # 记录被用作标题的 block（只取第一个，其余的作为副标题处理）
        title_block_processed = None
        if not title_text:
            title_blocks = [b for b in model.body_blocks if b.semantic_role == "title" and b.text]
            if title_blocks:
                title_text = title_blocks[0].text
                title_format = title_blocks[0].text_format
                title_block_processed = title_blocks[0]
                # 标记这个 block 对应的 raw_shape_id 已处理
                if title_blocks[0].raw_shape_id is not None:
                    processed_shape_ids.add(title_blocks[0].raw_shape_id)
        
        # 1. 填入标题（使用模板标题占位符的样式，原格式兜底）
        if title_text:
            title_placeholder = self._find_placeholder_by_role(slide, "title")
            if title_placeholder:
                self._fill_title_into_placeholder(slide, title_text, title_placeholder, title_format)
            else:
                # 备用：使用 slide.shapes.title
                self._fill_title_into_placeholder(slide, title_text, None, title_format)
        
        # 2. 根据语义角色分组正文内容
        # 剩余的 semantic_role="title" block（未被用作标题的）作为副标题处理
        body_main_blocks = []
        body_sidebar_blocks = []
        subtitle_blocks = []
        other_blocks = []
        
        for block in model.body_blocks:
            # 跳过被用作标题的那一个 block
            if title_block_processed is not None and block is title_block_processed:
                continue
            # 剩余的 role="title" block 作为副标题
            if block.semantic_role == "title" and block.text:
                subtitle_blocks.append(block)
                continue
            if block.type == "paragraph" and block.text:
                if block.semantic_role == "subtitle":
                    subtitle_blocks.append(block)
                elif block.semantic_role == "body_sidebar":
                    body_sidebar_blocks.append(block)
                elif block.semantic_role in ("body_main", "title", "unknown"):
                    body_main_blocks.append(block)
                else:
                    body_main_blocks.append(block)
            else:
                other_blocks.append(block)
        
        # 3. 填入副标题（如果有，原格式兜底）
        if subtitle_blocks:
            subtitle_placeholder = self._find_placeholder_by_role(slide, "subtitle")
            if subtitle_placeholder:
                self._fill_body_into_placeholder(slide, subtitle_blocks, subtitle_placeholder)
                # 标记这些 block 对应的 raw_shape_id 已处理
                for block in subtitle_blocks:
                    if block.raw_shape_id is not None:
                        processed_shape_ids.add(block.raw_shape_id)
            else:
                # 模板没有副标题占位符：将副标题合并到主正文的最前面
                body_main_blocks = subtitle_blocks + body_main_blocks
        
        # 4. 智能匹配正文占位符（优化多文本框的占位符分配）
        # 策略：
        # - 如果只有1个主正文shape → 填入body_main占位符
        # - 如果有多个主正文shape → 分析布局，智能匹配模板的多栏占位符
        # - 匹配成功的：填入占位符（只填内容，格式遵循模板）
        # - 匹配失败的：按原位置回填（保留原PPT布局）
        if body_main_blocks:
            # 按shape_id分组
            shape_blocks: dict[int, list] = {}
            for block in body_main_blocks:
                sid = block.raw_shape_id
                if sid is None:
                    sid = -1  # 无shape_id的放一起
                if sid not in shape_blocks:
                    shape_blocks[sid] = []
                shape_blocks[sid].append(block)
            
            main_shape_count = len([s for s in shape_blocks.keys() if s >= 0])
            
            if main_shape_count <= 1:
                # 只有一个主正文shape：填入body_main占位符
                body_placeholder = self._find_placeholder_by_role(slide, "body_main")
                if body_placeholder:
                    self._fill_body_into_placeholder(slide, body_main_blocks, body_placeholder)
                    for block in body_main_blocks:
                        if block.raw_shape_id is not None:
                            processed_shape_ids.add(block.raw_shape_id)
            else:
                # 多个主正文shape：尝试智能匹配模板占位符
                matched = self._smart_match_placeholders(slide, shape_blocks, processed_shape_ids)
                if not matched:
                    # 匹配失败：全部按原位置回填（在第7步处理）
                    pass
        
        # 5. 填入侧边栏内容（如果有对应的占位符）
        if body_sidebar_blocks:
            sidebar_placeholder = self._find_placeholder_by_role(slide, "body_sidebar")
            if sidebar_placeholder:
                self._fill_body_into_placeholder(slide, body_sidebar_blocks, sidebar_placeholder)
                for block in body_sidebar_blocks:
                    if block.raw_shape_id is not None:
                        processed_shape_ids.add(block.raw_shape_id)
            # 如果没有侧边栏占位符，不合并到主正文，而是按原位置回填（保留布局）
        
        # 6. 处理额外元素（图片、表格等），位置按比例适配，样式用模板配色
        for block in other_blocks:
            if block.type == "image":
                self._add_image_from_block(slide, block.content)
                if block.raw_shape_id is not None:
                    processed_shape_ids.add(block.raw_shape_id)
            elif block.type == "table":
                self._add_table_from_block(slide, block.content)
                if block.raw_shape_id is not None:
                    processed_shape_ids.add(block.raw_shape_id)
        
        # 7. 处理额外的自选图形和文本框（装饰性元素、未被占位符处理的内容）
        # 用 raw_shape_id 精确去重，避免内容重复处理
        # 注意：z-order顺序从底到顶：group/smartart → autoshape → text
        if model.raw_shapes:
            # 分离不同类型，确保正确的z-order
            group_list = []
            smartart_list = []
            autoshape_list = []
            text_shape_list = []
            other_shape_list = []
            
            for shape_data in model.raw_shapes:
                shape_id = shape_data.get("shape_id")
                
                # 精确去重：跳过已处理的 shape
                if shape_id is not None and shape_id in processed_shape_ids:
                    continue
                
                shape_type = shape_data.get("type")
                
                # 跳过图片/表格（已经在body_blocks中处理了）
                if shape_type in ("image", "table"):
                    continue
                
                if shape_type == "group":
                    group_list.append(shape_data)
                elif shape_type == "smartart":
                    smartart_list.append(shape_data)
                elif shape_type == "autoshape":
                    autoshape_list.append(shape_data)
                elif shape_type == "text":
                    text_shape_list.append(shape_data)
                else:
                    other_shape_list.append(shape_data)
            
            # 第1层：group和smartart（最底层）
            for shape_data in group_list:
                self._add_shape_from_data(slide, shape_data)
            for shape_data in smartart_list:
                self._add_shape_from_data(slide, shape_data)
            
            # 第2层：autoshape（背景形状）
            for shape_data in autoshape_list:
                self._add_extra_autoshape(slide, shape_data)
            
            # 第3层：其他形状（图片、图表等）
            for shape_data in other_shape_list:
                st = shape_data.get("type")
                if st == "image":
                    self._add_image_shape(slide, shape_data)
                elif st == "chart":
                    self._add_chart_shape(slide, shape_data)
                elif st == "ole":
                    self._add_ole_shape(slide, shape_data)
            
            # 第4层：text（文本层，确保在最上面）
            for shape_data in text_shape_list:
                self._add_extra_text_shape(slide, shape_data)

    def _find_placeholder_by_role(self, slide, role: str) -> object | None:
        """根据语义角色查找模板占位符
        
        匹配优先级（从高到低）：
        1. placeholder_mapping 中类型和索引都匹配
        2. placeholder_mapping 中类型匹配
        3. 占位符名称模糊匹配
        4. 角色特定的回退逻辑（title用slide.shapes.title，body找最大的文本占位符）
        """
        # 标题特殊处理：优先用 slide.shapes.title
        if role == "title":
            try:
                if slide.shapes.title is not None:
                    return slide.shapes.title
            except Exception:
                pass
        
        if self.placeholder_mapping:
            ph_info = self.placeholder_mapping.get(role)
            if ph_info:
                target_type = ph_info.get("placeholder_type")
                target_idx = ph_info.get("placeholder_idx")
                target_name = ph_info.get("name", "")
                
                # 第一优先级：类型和索引都匹配
                for shape in slide.placeholders:
                    try:
                        phf = shape.placeholder_format
                        if phf.type == target_type and phf.idx == target_idx:
                            return shape
                    except Exception:
                        continue
                
                # 第二优先级：类型匹配
                for shape in slide.placeholders:
                    try:
                        phf = shape.placeholder_format
                        if phf.type == target_type:
                            return shape
                    except Exception:
                        continue
                
                # 第三优先级：名称匹配
                if target_name:
                    for shape in slide.placeholders:
                        try:
                            if shape.name == target_name:
                                return shape
                        except Exception:
                            continue
        
        # 第四优先级：角色特定的回退逻辑
        if role == "title":
            # 找类型为1(TITLE)或3(CENTER_TITLE)的占位符
            for shape in slide.placeholders:
                try:
                    phf = shape.placeholder_format
                    if phf.type in (1, 3):
                        return shape
                except Exception:
                    continue
            # 最后尝试 slide.shapes.title（前面已经试过了，这里作为双保险）
            try:
                if slide.shapes.title is not None:
                    return slide.shapes.title
            except Exception:
                pass
        elif role == "body_main":
            # 找最大的文本/内容占位符（优先左半部分）
            best_placeholder = None
            best_area = 0
            slide_width = self.target_width or Emu(12192000)
            for shape in slide.placeholders:
                try:
                    phf = shape.placeholder_format
                    # 跳过页眉页脚
                    if phf.type in (13, 14, 15, 16):
                        continue
                    # 跳过标题
                    if phf.type in (1, 3):
                        continue
                    # 只考虑有文本框的占位符
                    if shape.has_text_frame:
                        area = (shape.width or 0) * (shape.height or 0)
                        # 主正文优先左半部分
                        left = shape.left or 0
                        if left < slide_width * 0.5:
                            area *= 1.5  # 左半部分加权
                        if area > best_area:
                            best_area = area
                            best_placeholder = shape
                except Exception:
                    continue
            if best_placeholder:
                return best_placeholder
        elif role == "body_sidebar":
            # 侧边栏：优先找右半部分的占位符
            slide_width = self.target_width or Emu(12192000)
            best_placeholder = None
            best_area = 0
            for shape in slide.placeholders:
                try:
                    phf = shape.placeholder_format
                    # 跳过页眉页脚
                    if phf.type in (13, 14, 15, 16):
                        continue
                    # 跳过标题
                    if phf.type in (1, 3):
                        continue
                    # 只考虑有文本框的占位符
                    if shape.has_text_frame:
                        left = shape.left or 0
                        # 侧边栏优先右半部分
                        if left >= slide_width * 0.5:
                            area = (shape.width or 0) * (shape.height or 0)
                            if area > best_area:
                                best_area = area
                                best_placeholder = shape
                except Exception:
                    continue
            if best_placeholder:
                return best_placeholder
        elif role == "subtitle":
            # 找类型为4(SUBTITLE)的占位符
            for shape in slide.placeholders:
                try:
                    phf = shape.placeholder_format
                    if phf.type == 4:
                        return shape
                except Exception:
                    continue
        
        return None
    
    def _smart_match_placeholders(self, slide, shape_blocks: dict[int, list], processed_shape_ids: set) -> bool:
        """智能匹配多文本框到模板占位符
        
        支持的布局模式：
        1. 左右双栏布局：模板2个占位符左右分布，原PPT文本框也分左右
        2. 上下布局：模板2个占位符上下分布，原PPT文本框也分上下
        3. 三栏布局：模板3个占位符水平分布，原PPT文本框也分三栏
        
        只有当占位符数量和文本shape数量匹配，且位置模式相似时才匹配
        否则保持原位置回填（保留原PPT布局）
        """
        try:
            # 获取模板的内容占位符（排除标题、页脚等）
            content_placeholders = []
            for ph in slide.placeholders:
                try:
                    phf = ph.placeholder_format
                    # 跳过标题、页脚、日期、页码
                    if phf.type in (1, 3, 4, 13, 14, 15, 16):
                        continue
                    # 只考虑有文本框的占位符（OBJECT和BODY类型都算）
                    if ph.has_text_frame:
                        content_placeholders.append(ph)
                except Exception:
                    continue
            
            if not content_placeholders:
                return False
            
            # 获取有位置信息的文本shape
            text_shapes_with_pos = []
            for shape_id, blocks in shape_blocks.items():
                if shape_id < 0:
                    continue
                # 从第一个block的shape_format获取位置信息
                block = blocks[0]
                left = None
                top = None
                width = None
                height = None
                if block.shape_format:
                    left = block.shape_format.left
                    top = block.shape_format.top
                    width = block.shape_format.width
                    height = block.shape_format.height
                if left is not None and top is not None:
                    text_shapes_with_pos.append((shape_id, blocks, left, top, width, height))
            
            if not text_shapes_with_pos:
                return False
            
            # 至少需要2个占位符和2个文本shape才尝试匹配
            if len(content_placeholders) < 2 or len(text_shapes_with_pos) < 2:
                return False
            
            slide_width = self.target_width or Emu(12192000)
            slide_height = self.target_height or Emu(6858000)
            
            # ========== 模式1：左右双栏布局 ==========
            if len(content_placeholders) >= 2:
                # 检查模板占位符是否是左右分布
                ph_lefts = [ph.left for ph in content_placeholders]
                ph_lefts_sorted = sorted(ph_lefts)
                # 如果最左和最右的占位符中心距离超过宽度的30%，认为是左右布局
                if len(ph_lefts_sorted) >= 2:
                    left_ph_center = ph_lefts_sorted[0] + content_placeholders[0].width / 2
                    right_ph_center = ph_lefts_sorted[-1] + content_placeholders[-1].width / 2
                    if right_ph_center - left_ph_center > slide_width * 0.3:
                        # 检查原PPT是否也是左右分栏
                        left_shapes = [s for s in text_shapes_with_pos if s[2] < slide_width * 0.5]
                        right_shapes = [s for s in text_shapes_with_pos if s[2] >= slide_width * 0.5]
                        
                        if left_shapes and right_shapes:
                            # 按left排序占位符
                            sorted_phs = sorted(content_placeholders, key=lambda x: x.left)
                            left_ph = sorted_phs[0]
                            right_ph = sorted_phs[-1]
                            
                            # 左栏内容（按top排序）
                            left_shapes.sort(key=lambda x: x[3])
                            left_blocks = []
                            for _, blocks, _, _, _, _ in left_shapes:
                                left_blocks.extend(blocks)
                            
                            # 右栏内容（按top排序）
                            right_shapes.sort(key=lambda x: x[3])
                            right_blocks = []
                            for _, blocks, _, _, _, _ in right_shapes:
                                right_blocks.extend(blocks)
                            
                            # 填入左栏
                            if left_blocks:
                                self._fill_body_into_placeholder(slide, left_blocks, left_ph)
                                for block in left_blocks:
                                    if block.raw_shape_id is not None:
                                        processed_shape_ids.add(block.raw_shape_id)
                            
                            # 填入右栏
                            if right_blocks:
                                self._fill_body_into_placeholder(slide, right_blocks, right_ph)
                                for block in right_blocks:
                                    if block.raw_shape_id is not None:
                                        processed_shape_ids.add(block.raw_shape_id)
                            
                            return True
            
            # ========== 模式2：上下布局 ==========
            if len(content_placeholders) >= 2:
                # 检查模板占位符是否是上下分布
                ph_tops = [ph.top for ph in content_placeholders]
                ph_tops_sorted = sorted(ph_tops)
                # 如果最上和最下的占位符中心距离超过高度的30%，认为是上下布局
                if len(ph_tops_sorted) >= 2:
                    top_ph_center = ph_tops_sorted[0] + content_placeholders[0].height / 2
                    bottom_ph_center = ph_tops_sorted[-1] + content_placeholders[-1].height / 2
                    if bottom_ph_center - top_ph_center > slide_height * 0.3:
                        # 检查原PPT是否也是上下分布
                        top_shapes = [s for s in text_shapes_with_pos if s[3] < slide_height * 0.5]
                        bottom_shapes = [s for s in text_shapes_with_pos if s[3] >= slide_height * 0.5]
                        
                        if top_shapes and bottom_shapes:
                            # 按top排序占位符
                            sorted_phs = sorted(content_placeholders, key=lambda x: x.top)
                            top_ph = sorted_phs[0]
                            bottom_ph = sorted_phs[-1]
                            
                            # 上栏内容（按left排序）
                            top_shapes.sort(key=lambda x: x[2])
                            top_blocks = []
                            for _, blocks, _, _, _, _ in top_shapes:
                                top_blocks.extend(blocks)
                            
                            # 下栏内容（按left排序）
                            bottom_shapes.sort(key=lambda x: x[2])
                            bottom_blocks = []
                            for _, blocks, _, _, _, _ in bottom_shapes:
                                bottom_blocks.extend(blocks)
                            
                            # 填入上栏
                            if top_blocks:
                                self._fill_body_into_placeholder(slide, top_blocks, top_ph)
                                for block in top_blocks:
                                    if block.raw_shape_id is not None:
                                        processed_shape_ids.add(block.raw_shape_id)
                            
                            # 填入下栏
                            if bottom_blocks:
                                self._fill_body_into_placeholder(slide, bottom_blocks, bottom_ph)
                                for block in bottom_blocks:
                                    if block.raw_shape_id is not None:
                                        processed_shape_ids.add(block.raw_shape_id)
                            
                            return True
            
            # ========== 模式3：三栏布局 ==========
            if len(content_placeholders) >= 3 and len(text_shapes_with_pos) >= 3:
                # 简单的三栏检测：3个占位符水平分布
                sorted_phs = sorted(content_placeholders, key=lambda x: x.left)
                # 检查是否大致三等分
                ph1_center = sorted_phs[0].left + sorted_phs[0].width / 2
                ph2_center = sorted_phs[1].left + sorted_phs[1].width / 2
                ph3_center = sorted_phs[2].left + sorted_phs[2].width / 2
                
                d1 = ph2_center - ph1_center
                d2 = ph3_center - ph2_center
                # 如果间距大致相等（差异<20%），认为是三栏布局
                if d1 > 0 and d2 > 0 and abs(d1 - d2) / max(d1, d2) < 0.3:
                    # 检查原PPT的文本框是否也分三栏
                    col1 = slide_width * 0.33
                    col2 = slide_width * 0.66
                    
                    col1_shapes = [s for s in text_shapes_with_pos if s[2] < col1]
                    col2_shapes = [s for s in text_shapes_with_pos if col1 <= s[2] < col2]
                    col3_shapes = [s for s in text_shapes_with_pos if s[2] >= col2]
                    
                    if col1_shapes and col2_shapes and col3_shapes:
                        # 收集每栏的内容
                        col1_blocks = []
                        col1_shapes.sort(key=lambda x: x[3])
                        for _, blocks, _, _, _, _ in col1_shapes:
                            col1_blocks.extend(blocks)
                        
                        col2_blocks = []
                        col2_shapes.sort(key=lambda x: x[3])
                        for _, blocks, _, _, _, _ in col2_shapes:
                            col2_blocks.extend(blocks)
                        
                        col3_blocks = []
                        col3_shapes.sort(key=lambda x: x[3])
                        for _, blocks, _, _, _, _ in col3_shapes:
                            col3_blocks.extend(blocks)
                        
                        # 填入三栏
                        for i, blocks in enumerate([col1_blocks, col2_blocks, col3_blocks]):
                            if blocks and i < len(sorted_phs):
                                self._fill_body_into_placeholder(slide, blocks, sorted_phs[i])
                                for block in blocks:
                                    if block.raw_shape_id is not None:
                                        processed_shape_ids.add(block.raw_shape_id)
                        
                        return True
            
            # 其他情况：不匹配，返回False
            return False
        except Exception:
            return False

    def _fill_simple_content(self, slide, model: SlideContentModel) -> None:
        if model.title:
            self._fill_title(slide, model.title)

        text_blocks = [b for b in model.body_blocks if b.type == "paragraph" and b.text]
        other_blocks = [b for b in model.body_blocks if b.type in ("table", "image")]

        if text_blocks:
            body_placeholder = self._find_body_placeholder(slide)
            if body_placeholder:
                self._fill_text_into_placeholder(body_placeholder, text_blocks)
            else:
                self._add_text_as_new_shape(slide, text_blocks)

        for block in other_blocks:
            if block.type == "image":
                self._add_image(slide, block.content)
            elif block.type == "table":
                self._add_table_from_content(slide, block.content)

    def _add_shape_from_data(self, slide, shape_data: dict) -> None:
        shape_type = shape_data.get("type")
        try:
            if shape_type == "text":
                # 判断是否是标题：位于页面顶部15%区域的文本
                is_title = False
                top = shape_data.get("top", 0)
                height = shape_data.get("height", 0)
                if top and self.target_height > 0:
                    if top + height < self.target_height * 0.2:
                        is_title = True
                self._add_text_shape(slide, shape_data, is_title=is_title)
            elif shape_type == "image":
                self._add_image_shape(slide, shape_data)
            elif shape_type == "table":
                self._add_table_shape(slide, shape_data)
            elif shape_type == "chart":
                self._add_chart_shape(slide, shape_data)
            elif shape_type == "ole":
                self._add_ole_shape(slide, shape_data)
            elif shape_type == "autoshape":
                self._add_auto_shape(slide, shape_data)
            elif shape_type == "group":
                for sub_shape in shape_data.get("shapes", []):
                    self._add_shape_from_data(slide, sub_shape)
            elif shape_type == "smartart":
                # SmartArt：作为group重建，内部形状应用颜色映射
                for sub_shape in shape_data.get("shapes", []):
                    self._add_shape_from_data(slide, sub_shape)
        except Exception:
            pass

    def _add_text_shape(self, slide, shape_data: dict, is_title: bool = False) -> None:
        left = self._adapt_position(shape_data.get("left", Emu(914400)), "x")
        top = self._adapt_position(shape_data.get("top", Emu(914400)), "y")
        width = self._adapt_position(shape_data.get("width", Emu(914400 * 8)), "x")
        height = self._adapt_position(shape_data.get("height", Emu(914400 * 2)), "y")

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.clear()

        # 应用文本框填充色和边框色（保留原格式）
        fill_color = shape_data.get("fill_color")
        if fill_color:
            try:
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = RGBColor.from_string(fill_color)
            except Exception:
                pass

        line_color = shape_data.get("line_color")
        line_width = shape_data.get("line_width")
        if line_color:
            try:
                textbox.line.color.rgb = RGBColor.from_string(line_color)
                if line_width:
                    textbox.line.width = line_width
            except Exception:
                pass

        # 获取模板格式规范
        template_fmt = self._get_template_format(is_title)

        paragraphs = shape_data.get("paragraphs", [])
        for i, para_data in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            runs = para_data.get("runs", [])
            if runs:
                for run_data in runs:
                    run = p.add_run()
                    run.text = run_data.get("text", "")
                    # 从run_data中提取原格式信息
                    original_format = None
                    if run_data:
                        from core.models import TextFormat
                        original_format = TextFormat(
                            font_name=run_data.get("font_name"),
                            font_size=run_data.get("font_size").pt if run_data.get("font_size") else None,
                            font_color=run_data.get("color"),
                            bold=run_data.get("bold"),
                            italic=run_data.get("italic"),
                            underline=run_data.get("underline"),
                        )
                    self._apply_run_format(run, run_data, template_fmt, is_title, original_format)
            else:
                p.text = para_data.get("text", "")
                # 应用模板格式到整段（原格式兜底）
                if template_fmt:
                    self._apply_template_format_to_paragraph(p, template_fmt)

            p.level = para_data.get("level", 0)
            if para_data.get("alignment") is not None:
                p.alignment = self._parse_alignment(para_data["alignment"])
            elif template_fmt and template_fmt.get("alignment"):
                p.alignment = self._parse_alignment(template_fmt["alignment"])

            try:
                if para_data.get("line_spacing") is not None:
                    p.line_spacing = para_data["line_spacing"]
                if para_data.get("space_before") is not None:
                    p.space_before = para_data["space_before"]
                if para_data.get("space_after") is not None:
                    p.space_after = para_data["space_after"]
            except Exception:
                pass

        # 自动调整文本框大小以适应内容，避免文本溢出
        self._auto_fit_textbox(textbox)

    def _get_template_format(self, is_title: bool) -> dict:
        """获取模板中标题或正文的格式规范"""
        if is_title:
            return self.template_formats.get("title", {})
        return self.template_formats.get("body", {})

    def _auto_fit_textbox(self, textbox) -> None:
        """自动调整文本框大小以适应内容，避免文本溢出
        
        策略：估算文本所需高度，如果超出当前容器高度，则扩大容器高度。
        """
        try:
            tf = textbox.text_frame
            if not tf.word_wrap:
                return

            box_width = textbox.width
            box_height = textbox.height
            if box_width <= 0 or box_height <= 0:
                return

            # 估算文本总高度
            total_text_height = 0
            for para in tf.paragraphs:
                text = para.text
                if not text:
                    total_text_height += Emu(914400 * 0.2)
                    continue

                # 获取字号
                font_size = Emu(914400 * 0.18)  # 默认14pt
                for run in para.runs:
                    if run.font.size:
                        font_size = run.font.size
                        break

                # 估算行数和行高
                char_width = font_size * 0.6
                chars_per_line = max(1, int(box_width / char_width))
                num_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
                line_height = font_size * 1.2
                total_text_height += num_lines * line_height

            # 如果文本高度超出容器高度，扩大容器
            if total_text_height > box_height:
                # 增加10%的余量
                new_height = int(total_text_height * 1.1)
                # 不超过页面高度的限制
                max_height = self.template_height * 0.8 if self.template_height else Emu(914400 * 7)
                if new_height > max_height:
                    new_height = max_height
                textbox.height = new_height
        except Exception:
            pass

    def _apply_run_format(self, run, run_data: dict, template_fmt: dict | None = None, is_title: bool = False, original_format=None) -> None:
        """应用run格式：模板格式优先，原格式兜底
        
        核心原则：
        - 字体：模板的主题字体优先，否则使用原格式的字体
        - 字号：模板定义的字号优先，否则使用原格式的字号
        - 颜色：模板定义的颜色优先，否则使用原格式的颜色
        - 粗体/斜体：模板要求的优先，否则保留原格式
        """
        font = run.font
        
        # 字体名称：模板主题字体优先，原格式兜底
        font_name = None
        if self.theme_fonts:
            font_name = self.theme_fonts.get("major" if is_title else "minor")
        if not font_name and template_fmt and template_fmt.get("font_name"):
            tmpl_font = template_fmt["font_name"]
            if tmpl_font.startswith("+mj") or tmpl_font.startswith("+mn"):
                if self.theme_fonts:
                    font_name = self.theme_fonts.get("major" if "mj" in tmpl_font else "minor")
            else:
                font_name = tmpl_font
        if not font_name and original_format and original_format.font_name:
            font_name = original_format.font_name
        if font_name:
            font.name = font_name
        
        # 强制使用模板主题字体（字体统一）
        self._apply_font(font, is_title)
        
        # 字号：模板要求优先，原格式兜底
        if template_fmt and template_fmt.get("font_size"):
            font.size = template_fmt["font_size"]
        elif original_format and original_format.font_size:
            font.size = Pt(original_format.font_size)
        
        # 粗体：模板要求优先，原格式兜底
        if template_fmt and template_fmt.get("bold") is not None:
            font.bold = template_fmt["bold"]
        elif original_format and original_format.bold is not None:
            font.bold = original_format.bold
        
        # 斜体：模板要求优先，原格式兜底
        if template_fmt and template_fmt.get("italic") is not None:
            font.italic = template_fmt["italic"]
        elif original_format and original_format.italic is not None:
            font.italic = original_format.italic
        
        # 下划线：模板要求优先，原格式兜底
        if template_fmt and template_fmt.get("underline") is not None:
            font.underline = template_fmt["underline"]
        elif original_format and original_format.underline is not None:
            font.underline = original_format.underline
        
        # 颜色：原格式优先（保留原设计意图），模板格式兜底，最后使用默认颜色
        color_applied = False
        if original_format and original_format.font_color:
            try:
                font.color.rgb = RGBColor.from_string(original_format.font_color)
                color_applied = True
            except Exception:
                pass
        elif template_fmt and template_fmt.get("color"):
            try:
                font.color.rgb = RGBColor.from_string(template_fmt["color"])
                color_applied = True
            except Exception:
                pass
        elif self.default_text_color:
            try:
                font.color.rgb = RGBColor.from_string(self.default_text_color)
                color_applied = True
            except Exception:
                pass

    def _apply_template_format_to_paragraph(self, paragraph, template_fmt: dict, original_format=None) -> None:
        """将模板格式应用到整个段落：模板格式优先，原格式兜底"""
        # 对齐：模板优先，原格式兜底
        if template_fmt.get("alignment"):
            paragraph.alignment = self._parse_alignment(template_fmt["alignment"])
        elif original_format and original_format.alignment is not None:
            from pptx.enum.text import PP_ALIGN
            alignment_map = {0: PP_ALIGN.LEFT, 1: PP_ALIGN.CENTER, 2: PP_ALIGN.RIGHT, 3: PP_ALIGN.JUSTIFY}
            paragraph.alignment = alignment_map.get(original_format.alignment)
        
        # 行距：模板优先，原格式兜底
        if original_format and original_format.line_spacing:
            try:
                paragraph.line_spacing = Pt(original_format.line_spacing)
            except Exception:
                pass
        
        for run in paragraph.runs:
            font = run.font
            # 字体名称：模板优先，原格式兜底
            if template_fmt.get("font_name"):
                font.name = template_fmt["font_name"]
            elif original_format and original_format.font_name:
                font.name = original_format.font_name
            
            # 强制使用模板主题字体（字体统一）
            self._apply_font(font, is_title=is_title)
            
            # 字号：模板优先，原格式兜底
            if template_fmt.get("font_size"):
                font.size = template_fmt["font_size"]
            elif original_format and original_format.font_size:
                font.size = Pt(original_format.font_size)
            
            # 粗体：模板优先，原格式兜底
            if template_fmt.get("bold") is not None:
                font.bold = template_fmt["bold"]
            elif original_format and original_format.bold is not None:
                font.bold = original_format.bold
            
            # 斜体：模板优先，原格式兜底
            if template_fmt.get("italic") is not None:
                font.italic = template_fmt["italic"]
            elif original_format and original_format.italic is not None:
                font.italic = original_format.italic
            
            # 颜色：原格式优先（保留原设计意图），模板格式兜底
            if original_format and original_format.font_color:
                try:
                    font.color.rgb = RGBColor.from_string(original_format.font_color)
                except Exception:
                    pass
            elif template_fmt.get("color"):
                try:
                    font.color.rgb = RGBColor.from_string(template_fmt["color"])
                except Exception:
                    pass

    def _build_color_mapping(self, content_models: list, accent_color: str | None) -> dict:
        """分析原PPT主色调，建立颜色映射表
        
        映射规则：
        1. 找出原PPT中出现频率最高的非中性色（主色）
        2. 找出原PPT的浅色版本（主色的浅色填充）
        3. 将原主色映射为模板强调色，原浅色映射为模板强调色的浅色版本
        4. 其他强调色也映射到模板强调色系
        """
        if not accent_color:
            return {}
        
        try:
            from collections import Counter
            
            # 收集所有填充色和描边色
            all_colors = []
            for model in content_models:
                for shape in model.raw_shapes:
                    if shape.get("fill_color"):
                        all_colors.append(shape["fill_color"].upper())
                    if shape.get("line_color"):
                        all_colors.append(shape["line_color"].upper())
            
            if not all_colors:
                return {}
            
            # 统计颜色频率
            color_counts = Counter(all_colors)
            
            # 过滤掉中性色（黑白灰），找出主色
            neutral_colors = {"FFFFFF", "000000", "CCCCCC", "999999", "666666", "333333", "F5F5F5", "E0E0E0"}
            non_neutral = [(c, cnt) for c, cnt in color_counts.items() if c not in neutral_colors and len(c) == 6]
            
            if not non_neutral:
                return {}
            
            # 按频率排序，取前几个主色
            non_neutral.sort(key=lambda x: x[1], reverse=True)
            
            color_mapping = {}
            processed_colors = set()
            
            # 第一个主色 → 模板强调色
            primary_color = non_neutral[0][0]
            color_mapping[primary_color] = accent_color.upper()
            processed_colors.add(primary_color)
            
            # 找主色的浅色版本（同一色系，亮度更高）
            primary_rgb = self._hex_to_rgb(primary_color)
            if primary_rgb:
                for color, cnt in non_neutral[1:]:
                    if color in processed_colors:
                        continue
                    color_rgb = self._hex_to_rgb(color)
                    if not color_rgb:
                        continue
                    # 判断是否是同一色系（色相接近）
                    if self._is_same_color_family(primary_rgb, color_rgb):
                        # 浅色 → 强调色的浅色版本
                        light_accent = self._lighten_color(accent_color, 0.7)
                        color_mapping[color] = light_accent
                        processed_colors.add(color)
            
            # 其他强调色 → 也映射到模板强调色系（不同深浅）
            for color, cnt in non_neutral:
                if color in processed_colors:
                    continue
                color_rgb = self._hex_to_rgb(color)
                if not color_rgb:
                    continue
                # 根据颜色亮度决定映射到哪个深浅
                luminance = self._get_luminance(color_rgb)
                if luminance > 0.6:
                    # 浅色 → 强调色的更浅版本
                    color_mapping[color] = self._lighten_color(accent_color, 0.8)
                elif luminance < 0.3:
                    # 深色 → 强调色的深色版本
                    color_mapping[color] = self._darken_color(accent_color, 0.7)
                else:
                    color_mapping[color] = accent_color.upper()
                processed_colors.add(color)
            
            return color_mapping
        except Exception:
            return {}
    
    def _map_color(self, color: str | None) -> str | None:
        """应用颜色映射，将原颜色转换为模板配色"""
        if not color:
            return None
        color_upper = color.upper()
        if color_upper in self.color_mapping:
            return self.color_mapping[color_upper]
        return color
    
    def _hex_to_rgb(self, hex_color: str) -> tuple | None:
        """将十六进制颜色转换为RGB元组"""
        try:
            hex_color = hex_color.strip("#").upper()
            if len(hex_color) != 6:
                return None
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return (r, g, b)
        except Exception:
            return None
    
    def _rgb_to_hex(self, r: int, g: int, b: int) -> str:
        """将RGB转换为十六进制颜色"""
        return f"{min(255, max(0, r)):02X}{min(255, max(0, g)):02X}{min(255, max(0, b)):02X}"
    
    def _get_luminance(self, rgb: tuple) -> float:
        """计算颜色的相对亮度"""
        r, g, b = rgb
        return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
    
    def _is_same_color_family(self, rgb1: tuple, rgb2: tuple, threshold: float = 0.2) -> bool:
        """判断两个颜色是否属于同一色系（基于归一化后的RGB比例）"""
        r1, g1, b1 = rgb1
        r2, g2, b2 = rgb2
        sum1 = max(r1 + g1 + b1, 1)
        sum2 = max(r2 + g2 + b2, 1)
        # 归一化
        nr1, ng1, nb1 = r1/sum1, g1/sum1, b1/sum1
        nr2, ng2, nb2 = r2/sum2, g2/sum2, b2/sum2
        # 计算差异
        diff = abs(nr1 - nr2) + abs(ng1 - ng2) + abs(nb1 - nb2)
        return diff < threshold
    
    def _lighten_color(self, hex_color: str, factor: float) -> str:
        """提亮颜色（factor: 0-1，越大越亮）"""
        rgb = self._hex_to_rgb(hex_color)
        if not rgb:
            return hex_color
        r, g, b = rgb
        r = int(r + (255 - r) * factor)
        g = int(g + (255 - g) * factor)
        b = int(b + (255 - b) * factor)
        return self._rgb_to_hex(r, g, b)
    
    def _darken_color(self, hex_color: str, factor: float) -> str:
        """加深颜色（factor: 0-1，越大越深）"""
        rgb = self._hex_to_rgb(hex_color)
        if not rgb:
            return hex_color
        r, g, b = rgb
        r = int(r * factor)
        g = int(g * factor)
        b = int(b * factor)
        return self._rgb_to_hex(r, g, b)
    
    def _apply_font(self, font, is_title: bool = False) -> None:
        """应用模板字体统一
        
        如果启用了强制模板字体，则将字体替换为模板的主题字体
        标题用major font，正文用minor font
        """
        if not self.force_template_font:
            return
        
        target_font = None
        if is_title and self.template_title_font:
            target_font = self.template_title_font
        elif not is_title and self.template_body_font:
            target_font = self.template_body_font
        
        if target_font:
            try:
                font.name = target_font
            except Exception:
                pass
    
    def _is_light_color(self, hex_color: str) -> bool:
        """判断颜色是否为浅色（亮度 > 0.5 视为浅色）"""
        if not hex_color or len(hex_color) != 6:
            return False
        try:
            r = int(hex_color[0:2], 16) / 255.0
            g = int(hex_color[2:4], 16) / 255.0
            b = int(hex_color[4:6], 16) / 255.0
            # 相对亮度计算
            luminance = 0.299 * r + 0.587 * g + 0.114 * b
            return luminance > 0.5
        except Exception:
            return False

    def _get_text_color_for_bg(self, bg_color: str | None) -> str:
        """根据背景颜色获取匹配的文字颜色（color pairing）
        
        - 浅色背景 → 深色文字（dk1: #0A2F24）
        - 深色背景 → 浅色文字（白色: #FFFFFF）
        - 无背景 → 使用默认文字颜色
        """
        if bg_color and self._is_light_color(bg_color):
            # 浅色背景：用深色文字
            return "0A2F24"
        elif bg_color and not self._is_light_color(bg_color):
            # 深色背景：用白色文字
            return "FFFFFF"
        else:
            # 无背景：使用默认文字颜色
            return self.default_text_color or "000000"

    def _is_decoration_shape(self, shape_data: dict) -> bool:
        """判断形状是否为装饰性形状（竖杆、小色块等）
        
        判断规则：
        - 宽度很小（< 0.5in）且高宽比 > 3 → 竖杆类装饰
        - 没有文字内容 → 装饰
        - 面积很小 → 装饰
        """
        try:
            width = shape_data.get("width", 0) or 0
            height = shape_data.get("height", 0) or 0
            width_in = width / 914400
            height_in = height / 914400
            
            # 检查是否有文字
            has_text = False
            for para in shape_data.get("paragraphs", []):
                for run in para.get("runs", []):
                    if run.get("text", "").strip():
                        has_text = True
                        break
                if has_text:
                    break
            
            # 窄竖条（宽度<0.5in，高宽比>3）→ 装饰
            if width_in < 0.5 and height_in > 0 and height_in / max(width_in, 0.01) > 3:
                return True
            
            # 没有文字的小形状 → 装饰
            if not has_text and width_in < 2 and height_in < 2:
                return True
            
            return False
        except Exception:
            return False

    def _parse_alignment(self, align_str: str):
        """将字符串对齐方式转为PP_ALIGN枚举"""
        if not align_str:
            return None
        align_map = {
            "LEFT": PP_ALIGN.LEFT,
            "CENTER": PP_ALIGN.CENTER,
            "RIGHT": PP_ALIGN.RIGHT,
            "JUSTIFY": PP_ALIGN.JUSTIFY,
            "DISTRIBUTE": PP_ALIGN.DISTRIBUTE,
        }
        for key, val in align_map.items():
            if key in align_str.upper():
                return val
        return None

    def _add_image_shape(self, slide, shape_data: dict) -> None:
        blob = shape_data.get("blob")
        if not blob:
            return
        stream = BytesIO(blob)
        pic = slide.shapes.add_picture(
            stream,
            left=self._adapt_position(shape_data.get("left", Emu(914400)), "x"),
            top=self._adapt_position(shape_data.get("top", Emu(914400)), "y"),
            width=self._adapt_position(shape_data.get("width"), "x"),
            height=self._adapt_position(shape_data.get("height"), "y"),
        )
        # 应用裁剪
        try:
            if shape_data.get("crop_left"):
                pic.crop_left = shape_data["crop_left"]
            if shape_data.get("crop_top"):
                pic.crop_top = shape_data["crop_top"]
            if shape_data.get("crop_right"):
                pic.crop_right = shape_data["crop_right"]
            if shape_data.get("crop_bottom"):
                pic.crop_bottom = shape_data["crop_bottom"]
        except Exception:
            pass

    def _add_table_shape(self, slide, shape_data: dict) -> None:
        table_data = shape_data.get("data", {})
        if not table_data:
            return

        # 兼容新旧格式
        if isinstance(table_data, dict):
            self._add_table_full(slide, shape_data)
        else:
            # 旧格式：二维数组
            self._add_table_simple(slide, shape_data, table_data)

    def _add_table_full(self, slide, shape_data: dict) -> None:
        """完整表格重放：保留单元格背景色、字体格式、合并单元格"""
        table_data = shape_data.get("data", {})
        rows = table_data.get("rows", 0)
        cols = table_data.get("cols", 0)
        cells_data = table_data.get("cells", [])

        if rows == 0 or cols == 0:
            return

        table_shape = slide.shapes.add_table(
            rows, cols,
            shape_data.get("left", Emu(914400)),
            shape_data.get("top", Emu(914400 * 3)),
            shape_data.get("width", Emu(914400 * 8)),
            shape_data.get("height", Emu(914400 * 2)),
        )
        table = table_shape.table

        # 设置列宽
        col_widths = table_data.get("col_widths", [])
        for i, width in enumerate(col_widths):
            if i < len(table.columns) and width:
                table.columns[i].width = width

        # 设置行高
        row_heights = table_data.get("row_heights", [])
        for i, height in enumerate(row_heights):
            if i < len(table.rows) and height:
                table.rows[i].height = height

        # 填充单元格内容和格式
        for cell_data in cells_data:
            r = cell_data["row"]
            c = cell_data["col"]
            if r >= rows or c >= cols:
                continue

            cell = table.cell(r, c)
            cell.text = cell_data.get("text", "")

            # 背景色（应用颜色映射，遵循模板配色）
            if cell_data.get("fill_color"):
                try:
                    fill_color = self._map_color(cell_data["fill_color"])
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(fill_color)
                except Exception:
                    pass

            # 对齐和锚点
            try:
                if cell_data.get("alignment"):
                    align = self._parse_alignment(cell_data["alignment"])
                    if align and cell.text_frame.paragraphs:
                        cell.text_frame.paragraphs[0].alignment = align
            except Exception:
                pass

            try:
                if cell_data.get("vertical_anchor"):
                    anchor_map = {
                        "TOP": MSO_ANCHOR.TOP,
                        "MIDDLE": MSO_ANCHOR.MIDDLE,
                        "BOTTOM": MSO_ANCHOR.BOTTOM,
                    }
                    for key, val in anchor_map.items():
                        if key in cell_data["vertical_anchor"].upper():
                            cell.vertical_anchor = val
                            break
            except Exception:
                pass

            # 边距
            try:
                if cell_data.get("margin_left") is not None:
                    cell.margin_left = cell_data["margin_left"]
                if cell_data.get("margin_right") is not None:
                    cell.margin_right = cell_data["margin_right"]
                if cell_data.get("margin_top") is not None:
                    cell.margin_top = cell_data["margin_top"]
                if cell_data.get("margin_bottom") is not None:
                    cell.margin_bottom = cell_data["margin_bottom"]
            except Exception:
                pass

            # 字体格式
            try:
                if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                    run = cell.text_frame.paragraphs[0].runs[0]
                    font = run.font
                    if cell_data.get("font_name"):
                        font.name = cell_data["font_name"]
                    # 强制使用模板主题字体（字体统一，表格用正文字体）
                    self._apply_font(font, is_title=False)
                    if cell_data.get("font_size"):
                        font.size = cell_data["font_size"]
                    if cell_data.get("bold") is not None:
                        font.bold = cell_data["bold"]
                    if cell_data.get("italic") is not None:
                        font.italic = cell_data["italic"]
                    if cell_data.get("color"):
                        try:
                            # 应用颜色映射（遵循模板配色）
                            text_color = self._map_color(cell_data["color"])
                            font.color.rgb = RGBColor.from_string(text_color)
                        except Exception:
                            pass
            except Exception:
                pass

        # 合并单元格
        merged_cells = table_data.get("merged_cells", [])
        for merge_info in merged_cells:
            try:
                start_r = merge_info["start_row"]
                start_c = merge_info["start_col"]
                row_span = merge_info.get("row_span", 1)
                col_span = merge_info.get("col_span", 1)

                if row_span > 1 or col_span > 1:
                    end_r = start_r + row_span - 1
                    end_c = start_c + col_span - 1
                    cell_a = table.cell(start_r, start_c)
                    cell_b = table.cell(end_r, end_c)
                    cell_a.merge(cell_b)
            except Exception:
                pass
        
        # Color Pairing：确保表格文字在背景上可见
        try:
            if self.default_text_color:
                for r in range(rows):
                    for c in range(cols):
                        cell = table.cell(r, c)
                        # 检查单元格是否有背景色
                        cell_fill_color = None
                        try:
                            if cell.fill.type == 1:  # solid
                                cell_fill_color = str(cell.fill.fore_color.rgb)
                        except:
                            pass
                        
                        # 决定文字颜色
                        if cell_fill_color:
                            # 有单元格背景色：根据单元格背景色调整
                            text_color = self._get_text_color_for_bg(cell_fill_color)
                        else:
                            # 没有单元格背景色：用默认文字颜色
                            text_color = self.default_text_color
                        
                        # 应用文字颜色
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    run.font.color.rgb = RGBColor.from_string(text_color)
                                except Exception:
                                    pass
        except Exception:
            pass

    def _add_table_simple(self, slide, shape_data: dict, data: list) -> None:
        """简化版表格重放（向后兼容）"""
        if not data:
            return
        rows = len(data)
        cols = len(data[0]) if data[0] else 1
        table_shape = slide.shapes.add_table(
            rows, cols,
            shape_data.get("left", Emu(914400)),
            shape_data.get("top", Emu(914400 * 3)),
            shape_data.get("width", Emu(914400 * 8)),
            shape_data.get("height", Emu(914400 * 2)),
        )
        table = table_shape.table
        for r, row_data in enumerate(data):
            for c, cell_text in enumerate(row_data):
                if c < cols:
                    table.cell(r, c).text = str(cell_text)
        
        # Color Pairing：确保表格文字在背景上可见
        try:
            if self.default_text_color:
                for r in range(rows):
                    for c in range(cols):
                        cell = table.cell(r, c)
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    run.font.color.rgb = RGBColor.from_string(self.default_text_color)
                                except Exception:
                                    pass
        except Exception:
            pass

    def _add_chart_shape(self, slide, shape_data: dict) -> None:
        """完整图表重放：保留图例、坐标轴、系列颜色"""
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

            chart_data_list = shape_data.get("data", [])
            categories = shape_data.get("categories", [])

            if not chart_data_list or not categories:
                return

            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = categories

            for series in chart_data_list:
                series_name = series.get("name", "Series")
                values = series.get("values", [])
                if values:
                    chart_data_obj.add_series(series_name, values)

            left = shape_data.get("left", Emu(914400))
            top = shape_data.get("top", Emu(914400))
            width = shape_data.get("width", Emu(914400 * 8))
            height = shape_data.get("height", Emu(914400 * 4))

            # 解析图表类型
            chart_type_str = shape_data.get("chart_type", "")
            chart_type = self._parse_chart_type(chart_type_str)

            chart_frame = slide.shapes.add_chart(
                chart_type,
                left, top, width, height,
                chart_data_obj
            )
            chart = chart_frame.chart

            # 设置图表标题
            if shape_data.get("chart_title"):
                chart.has_title = True
                chart.chart_title.text_frame.text = shape_data["chart_title"]
            else:
                chart.has_title = False

            # 设置图例
            if shape_data.get("has_legend"):
                chart.has_legend = True
                pos_str = shape_data.get("legend_position", "")
                if "BOTTOM" in pos_str:
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                elif "TOP" in pos_str:
                    chart.legend.position = XL_LEGEND_POSITION.TOP
                elif "LEFT" in pos_str:
                    chart.legend.position = XL_LEGEND_POSITION.LEFT
                elif "RIGHT" in pos_str:
                    chart.legend.position = XL_LEGEND_POSITION.RIGHT
            else:
                chart.has_legend = False

            # 设置系列颜色（应用颜色映射，遵循模板配色）
            for i, series in enumerate(chart.series):
                if i < len(chart_data_list):
                    color = chart_data_list[i].get("color")
                    if color:
                        try:
                            mapped_color = self._map_color(color)
                            series.format.fill.solid()
                            series.format.fill.fore_color.rgb = RGBColor.from_string(mapped_color)
                        except Exception:
                            pass

            # 设置坐标轴
            x_axis_data = shape_data.get("x_axis", {})
            y_axis_data = shape_data.get("y_axis", {})
            try:
                if chart.category_axis and x_axis_data:
                    ca = chart.category_axis
                    if x_axis_data.get("has_title") and x_axis_data.get("title"):
                        ca.has_title = True
                        ca.axis_title.text_frame.text = x_axis_data["title"]
                    if x_axis_data.get("visible") is not None:
                        ca.visible = x_axis_data["visible"]
            except Exception:
                pass
            try:
                if chart.value_axis and y_axis_data:
                    va = chart.value_axis
                    if y_axis_data.get("has_title") and y_axis_data.get("title"):
                        va.has_title = True
                        va.axis_title.text_frame.text = y_axis_data["title"]
                    if y_axis_data.get("visible") is not None:
                        va.visible = y_axis_data["visible"]
                    if y_axis_data.get("minimum_scale") is not None:
                        va.minimum_scale = y_axis_data["minimum_scale"]
                    if y_axis_data.get("maximum_scale") is not None:
                        va.maximum_scale = y_axis_data["maximum_scale"]
            except Exception:
                pass

        except Exception:
            pass

    def _parse_chart_type(self, chart_type_str: str):
        """将图表类型字符串转为XL_CHART_TYPE枚举"""
        from pptx.enum.chart import XL_CHART_TYPE
        type_map = {
            "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "COLUMN_STACKED": XL_CHART_TYPE.COLUMN_STACKED,
            "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
            "LINE": XL_CHART_TYPE.LINE,
            "LINE_MARKERS": XL_CHART_TYPE.LINE_MARKERS,
            "PIE": XL_CHART_TYPE.PIE,
            "DOUGHNUT": XL_CHART_TYPE.DOUGHNUT,
            "AREA": XL_CHART_TYPE.AREA,
            "XY_SCATTER": XL_CHART_TYPE.XY_SCATTER,
            "BUBBLE": XL_CHART_TYPE.BUBBLE,
            "RADAR": XL_CHART_TYPE.RADAR,
        }
        for key, val in type_map.items():
            if key in chart_type_str.upper():
                return val
        return XL_CHART_TYPE.COLUMN_CLUSTERED

    def _add_ole_shape(self, slide, shape_data: dict) -> None:
        try:
            blob = shape_data.get("blob")
            if not blob:
                return
            stream = BytesIO(blob)
            left = shape_data.get("left", Emu(914400))
            top = shape_data.get("top", Emu(914400))
            width = shape_data.get("width", Emu(914400 * 6))
            height = shape_data.get("height", Emu(914400 * 4))
            slide.shapes.add_ole_object(stream, left, top, width, height)
        except Exception:
            pass

    def _add_auto_shape(self, slide, shape_data: dict) -> None:
        """重放自选图形：保留填充色、边框色等样式"""
        try:
            from pptx.enum.shapes import MSO_SHAPE
            shape_type_str = shape_data.get("shape_type", "")
            mso_shape = MSO_SHAPE.RECTANGLE
            try:
                if "OVAL" in shape_type_str or "ELLIPSE" in shape_type_str:
                    mso_shape = MSO_SHAPE.OVAL
                elif "ROUNDED_RECTANGLE" in shape_type_str:
                    mso_shape = MSO_SHAPE.ROUNDED_RECTANGLE
                elif "LINE" in shape_type_str:
                    mso_shape = MSO_SHAPE.LINE
                elif "ARROW" in shape_type_str:
                    mso_shape = MSO_SHAPE.RIGHT_ARROW
                elif "DIAMOND" in shape_type_str:
                    mso_shape = MSO_SHAPE.DIAMOND
                elif "TRIANGLE" in shape_type_str:
                    mso_shape = MSO_SHAPE.RIGHT_TRIANGLE
                elif "STAR" in shape_type_str:
                    mso_shape = MSO_SHAPE.FIVE_POINTED_STAR
            except Exception:
                pass

            left = self._adapt_position(shape_data.get("left", Emu(914400)), "x")
            top = self._adapt_position(shape_data.get("top", Emu(914400)), "y")
            width = self._adapt_position(shape_data.get("width", Emu(914400 * 4)), "x")
            height = self._adapt_position(shape_data.get("height", Emu(914400 * 2)), "y")

            shape = slide.shapes.add_shape(
                mso_shape, left, top, width, height
            )

            try:
                shape.rotation = shape_data.get("rotation", 0)
            except Exception:
                pass

            fill_color = shape_data.get("fill_color")
            if fill_color:
                try:
                    shape.fill.solid()
                    # 深色背景模板：装饰性形状改为白色
                    if self.default_text_color == "FFFFFF":
                        shape.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
                    else:
                        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
                except Exception:
                    pass

            line_color = shape_data.get("line_color")
            line_width = shape_data.get("line_width")
            if line_color:
                try:
                    # 深色背景模板：装饰性线条改为白色
                    if self.default_text_color == "FFFFFF":
                        shape.line.color.rgb = RGBColor.from_string("FFFFFF")
                    else:
                        shape.line.color.rgb = RGBColor.from_string(line_color)
                    if line_width:
                        shape.line.width = line_width
                except Exception:
                    pass

            if shape_data.get("text") and shape.has_text_frame:
                shape.text_frame.text = shape_data["text"]
                # 自动调整自选图形大小以适应文本
                self._auto_fit_textbox(shape)
        except Exception:
            pass

    def _fill_title_into_placeholder(self, slide, title_text: str, title_placeholder=None, original_format=None) -> None:
        """将标题填入模板的标题占位符，并应用模板的标题样式
        
        如果指定了 title_placeholder，则填入该占位符；
        否则查找 slide.shapes.title 或第一个标题占位符。
        
        original_format: 原PPT中标题的格式信息（TextFormat对象），用于模板未定义时兜底
        """
        try:
            # 优先使用传入的占位符
            if title_placeholder is not None:
                tf = title_placeholder.text_frame
                template_fmt = self._get_template_format(True)
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    for run in p.runs:
                        run.text = ""
                    if p.runs:
                        run = p.runs[0]
                        run.text = title_text
                        self._apply_run_format(run, {}, template_fmt, True, original_format)
                    else:
                        run = p.add_run()
                        run.text = title_text
                        self._apply_run_format(run, {}, template_fmt, True, original_format)
                else:
                    tf.text = title_text
                return
            
            # 备用：使用 slide.shapes.title
            if slide.shapes.title is not None:
                tf = slide.shapes.title.text_frame
                template_fmt = self._get_template_format(True)
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    for run in p.runs:
                        run.text = ""
                    if p.runs:
                        run = p.runs[0]
                        run.text = title_text
                        self._apply_run_format(run, {}, template_fmt, True, original_format)
                    else:
                        run = p.add_run()
                        run.text = title_text
                        self._apply_run_format(run, {}, template_fmt, True, original_format)
                else:
                    tf.text = title_text
                return
        except Exception:
            pass
        
        # 备用：通过placeholder查找
        for shape in slide.placeholders:
            try:
                phf = shape.placeholder_format
                if phf.type in (0, 3):  # title, ctrTitle
                    tf = shape.text_frame
                    template_fmt = self._get_template_format(True)
                    if tf.paragraphs:
                        p = tf.paragraphs[0]
                        for run in p.runs:
                            run.text = ""
                        if p.runs:
                            run = p.runs[0]
                            run.text = title_text
                            self._apply_run_format(run, {}, template_fmt, True, original_format)
                        else:
                            run = p.add_run()
                            run.text = title_text
                            self._apply_run_format(run, {}, template_fmt, True, original_format)
                    return
            except Exception:
                continue

    def _fill_body_into_placeholder(self, slide, text_blocks: list, body_placeholder=None) -> None:
        """将正文填入模板的正文占位符：模板格式优先，原格式兜底
        
        如果指定了 body_placeholder，则填入该占位符；
        否则使用 _find_body_placeholder 查找。
        """
        if not body_placeholder:
            body_placeholder = self._find_body_placeholder(slide)
        if not body_placeholder:
            return
        
        try:
            tf = body_placeholder.text_frame
            template_fmt = self.template_formats.get("body", {})
            
            # 清空现有内容（但保留第一段的样式）
            tf.clear()
            
            # 填入新内容，每段应用模板格式 + 原格式兜底
            first = True
            for block in text_blocks:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.text = block.text
                p.level = block.level
                
                # 获取该block的原始格式信息
                original_format = block.text_format if hasattr(block, 'text_format') else None
                
                # 应用格式：模板格式优先，原格式兜底
                if template_fmt or original_format:
                    self._apply_template_format_to_paragraph(p, template_fmt, original_format)
        except Exception:
            pass

    def _add_extra_text_shape(self, slide, shape_data: dict) -> None:
        """添加额外的文本框：位置适配，模板格式优先，原格式兜底
        
        Color Pairing 规则：
        - 装饰形状（竖杆、小色块）：用模板强调色（浅色背景→深绿，深色背景→白色）
        - 内容框（有文字的大形状）：保留原背景色，文字根据背景色深浅自动匹配
        """
        try:
            left = self._adapt_position(shape_data.get("left", Emu(914400)), "x")
            top = self._adapt_position(shape_data.get("top", Emu(914400)), "y")
            width = self._adapt_position(shape_data.get("width", Emu(914400 * 4)), "x")
            height = self._adapt_position(shape_data.get("height", Emu(914400 * 2)), "y")
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()
            
            fill_color = shape_data.get("fill_color")
            is_decoration = self._is_decoration_shape(shape_data)
            
            if is_decoration:
                # 装饰形状：用模板强调色
                textbox.fill.solid()
                if self.default_text_color == "FFFFFF":
                    textbox.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
                else:
                    textbox.fill.fore_color.rgb = RGBColor.from_string("3DCD58")
            elif fill_color:
                # 内容框：保留原背景色（因为背景色决定了内容框的视觉效果）
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = RGBColor.from_string(fill_color)
            
            # 处理线条颜色
            line_color = shape_data.get("line_color")
            if line_color:
                if is_decoration:
                    if self.default_text_color == "FFFFFF":
                        textbox.line.color.rgb = RGBColor.from_string("FFFFFF")
                    else:
                        textbox.line.color.rgb = RGBColor.from_string("3DCD58")
                else:
                    textbox.line.color.rgb = RGBColor.from_string(line_color)
            
            # 判断是否是标题区域的文本
            is_title = top < self.target_height * 0.2
            template_fmt = self._get_template_format(is_title)
            
            # 确定文字颜色：根据背景色做 color pairing
            actual_bg_color = None
            if is_decoration:
                actual_bg_color = "FFFFFF" if self.default_text_color == "FFFFFF" else "3DCD58"
            elif fill_color:
                actual_bg_color = fill_color
            else:
                # 没有填充色的文本框：根据master_style确定背景色
                # F1/F2: 浅色背景，F3/F4: 深色背景
                master_style = getattr(self, 'master_style', 'F1').upper()
                if master_style in ("F3", "F4"):
                    # 深色背景模板：给文本框加浅色背景，文字用深色
                    textbox.fill.solid()
                    textbox.fill.fore_color.rgb = RGBColor.from_string("E7FFD9")
                    actual_bg_color = "E7FFD9"
                else:
                    # 浅色背景模板：使用白色背景，文字用深色
                    actual_bg_color = "FFFFFF"
            
            text_color_override = self._get_text_color_for_bg(actual_bg_color) if actual_bg_color else None
            
            paragraphs = shape_data.get("paragraphs", [])
            for i, para_data in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                runs = para_data.get("runs", [])
                if runs:
                    for run_data in runs:
                        run = p.add_run()
                        run.text = run_data.get("text", "")
                        # 从run_data中提取原格式信息
                        original_format = None
                        if run_data:
                            from core.models import TextFormat
                            original_format = TextFormat(
                                font_name=run_data.get("font_name"),
                                font_size=run_data.get("font_size").pt if run_data.get("font_size") else None,
                                font_color=run_data.get("color"),
                                bold=run_data.get("bold"),
                                italic=run_data.get("italic"),
                                underline=run_data.get("underline"),
                            )
                        self._apply_run_format(run, run_data, template_fmt, is_title, original_format)
                        # 应用 color pairing 的文字颜色（覆盖模板样式）
                        if text_color_override and run.text.strip():
                            try:
                                run.font.color.rgb = RGBColor.from_string(text_color_override)
                            except Exception:
                                pass
                else:
                    p.text = para_data.get("text", "")
                    if template_fmt:
                        self._apply_template_format_to_paragraph(p, template_fmt)
                    if text_color_override and p.text.strip():
                        try:
                            for run in p.runs:
                                run.font.color.rgb = RGBColor.from_string(text_color_override)
                        except Exception:
                            pass
                
                p.level = para_data.get("level", 0)
            
            # 自动调整文本框大小以适应内容，避免文本溢出
            self._auto_fit_textbox(textbox)
        except Exception:
            pass

    def _add_extra_autoshape(self, slide, shape_data: dict) -> None:
        """添加额外的装饰形状：位置适配，模板格式优先，原格式兜底
        
        Color Pairing 规则：
        - 装饰形状（小、无文字）：用模板强调色
        - 内容框（大、有文字）：保留原背景色，文字根据背景色匹配
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE
            
            shape_type_str = shape_data.get("shape_type", "")
            mso_shape = MSO_SHAPE.RECTANGLE
            try:
                if "OVAL" in shape_type_str or "ELLIPSE" in shape_type_str:
                    mso_shape = MSO_SHAPE.OVAL
                elif "ROUNDED_RECTANGLE" in shape_type_str:
                    mso_shape = MSO_SHAPE.ROUNDED_RECTANGLE
                elif "LINE" in shape_type_str:
                    mso_shape = MSO_SHAPE.LINE
                elif "ARROW" in shape_type_str:
                    mso_shape = MSO_SHAPE.RIGHT_ARROW
            except Exception:
                pass
            
            left = self._adapt_position(shape_data.get("left", Emu(914400)), "x")
            top = self._adapt_position(shape_data.get("top", Emu(914400)), "y")
            width = self._adapt_position(shape_data.get("width", Emu(914400 * 4)), "x")
            height = self._adapt_position(shape_data.get("height", Emu(914400 * 2)), "y")
            
            shape = slide.shapes.add_shape(mso_shape, left, top, width, height)
            
            try:
                shape.rotation = shape_data.get("rotation", 0)
            except Exception:
                pass
            
            fill_color = shape_data.get("fill_color")
            is_decoration = self._is_decoration_shape(shape_data)
            
            # 应用颜色映射（遵循模板配色要求）
            if fill_color:
                fill_color = self._map_color(fill_color)
            line_color = shape_data.get("line_color")
            if line_color:
                line_color = self._map_color(line_color)
            
            # 确定背景色：保留原格式优先，已应用颜色映射
            actual_bg_color = None
            if is_decoration:
                # 装饰形状：用模板强调色
                if self.template_accent_color:
                    accent = self.template_accent_color
                    if self.default_text_color == "FFFFFF":
                        # 深色背景下用浅色装饰
                        accent = self._lighten_color(accent, 0.5)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(accent)
                    actual_bg_color = accent
                    try:
                        shape.line.color.rgb = RGBColor.from_string(accent)
                    except Exception:
                        pass
                elif self.default_text_color == "FFFFFF":
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
                    actual_bg_color = "FFFFFF"
                    try:
                        shape.line.color.rgb = RGBColor.from_string("FFFFFF")
                    except Exception:
                        pass
                else:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string("3DCD58")
                    actual_bg_color = "3DCD58"
                    try:
                        shape.line.color.rgb = RGBColor.from_string("3DCD58")
                    except Exception:
                        pass
            elif fill_color:
                # 内容框：使用已应用颜色映射后的背景色
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
                actual_bg_color = fill_color
                if line_color:
                    try:
                        shape.line.color.rgb = RGBColor.from_string(line_color)
                    except Exception:
                        pass
            else:
                # 内容框没有填充色：根据整体背景色决定是否加背景
                if self.default_text_color == "FFFFFF":
                    # 深色背景下，给内容框加浅色背景，确保文字可见
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string("E7FFD9")
                    actual_bg_color = "E7FFD9"
                else:
                    # 浅色背景下，内容框透明也没关系，文字是深色的
                    actual_bg_color = None
            
            # 如果有文字，应用模板样式 + 原格式兜底 + color pairing
            if shape_data.get("text") and shape.has_text_frame:
                shape.text_frame.text = shape_data["text"]
                template_fmt = self._get_template_format(False)
                if template_fmt and shape.text_frame.paragraphs:
                    for p in shape.text_frame.paragraphs:
                        self._apply_template_format_to_paragraph(p, template_fmt)
                # color pairing：根据背景色调整文字颜色
                if actual_bg_color:
                    text_color = self._get_text_color_for_bg(actual_bg_color)
                    try:
                        for p in shape.text_frame.paragraphs:
                            for run in p.runs:
                                run.font.color.rgb = RGBColor.from_string(text_color)
                    except Exception:
                        pass
        except Exception:
            pass

    def _add_image_from_block(self, slide, content: dict) -> None:
        """从body block添加图片：位置适配"""
        try:
            from io import BytesIO
            stream = BytesIO(content.get("blob", b""))
            left = self._adapt_position(content.get("left", Emu(914400)), "x")
            top = self._adapt_position(content.get("top", Emu(914400)), "y")
            width = self._adapt_position(content.get("width", Emu(914400 * 4)), "x")
            height = self._adapt_position(content.get("height", Emu(914400 * 3)), "y")
            slide.shapes.add_picture(stream, left, top, width, height)
        except Exception:
            pass

    def _add_table_from_block(self, slide, content: dict) -> None:
        """从body block添加表格"""
        try:
            self._add_table_full(slide, {
                "left": content.get("left", Emu(914400)),
                "top": content.get("top", Emu(914400)),
                "width": content.get("width", Emu(914400 * 8)),
                "height": content.get("height", Emu(914400 * 4)),
                "data": content,
            })
        except Exception:
            pass

    def _fill_title(self, slide, title_text: str) -> None:
        try:
            if slide.shapes.title is not None:
                slide.shapes.title.text = title_text
                # 应用模板标题格式
                template_fmt = self.template_formats.get("title", {})
                if template_fmt and slide.shapes.title.text_frame.paragraphs:
                    p = slide.shapes.title.text_frame.paragraphs[0]
                    self._apply_template_format_to_paragraph(p, template_fmt)
                return
        except Exception:
            pass

        for shape in slide.placeholders:
            try:
                phf = shape.placeholder_format
                if phf.type == 0:
                    shape.text_frame.text = title_text
                    template_fmt = self.template_formats.get("title", {})
                    if template_fmt and shape.text_frame.paragraphs:
                        p = shape.text_frame.paragraphs[0]
                        self._apply_template_format_to_paragraph(p, template_fmt)
                    return
            except Exception:
                continue

    def _find_body_placeholder(self, slide) -> object | None:
        for shape in slide.placeholders:
            try:
                phf = shape.placeholder_format
                ph_type = phf.type

                if ph_type == 2:
                    return shape
                if ph_type in (7, 8, 9, 10):
                    return shape
                if ph_type not in (0, 1, 3, 4):
                    return shape
            except Exception:
                continue
        return None

    def _fill_text_into_placeholder(self, placeholder, text_blocks) -> None:
        try:
            tf = placeholder.text_frame
            tf.clear()
            tf.word_wrap = True

            # 获取模板正文格式
            template_fmt = self.template_formats.get("body", {})

            first = True
            for block in text_blocks:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.text = block.text
                p.level = block.level

                # 应用模板正文格式
                if template_fmt:
                    self._apply_template_format_to_paragraph(p, template_fmt)
        except Exception:
            pass

    def _add_text_as_new_shape(self, slide, text_blocks) -> None:
        try:
            left = Emu(914400)
            top = Emu(914400 * 3)
            width = Emu(914400 * 10)
            height = Emu(914400 * 6)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.word_wrap = True
            tf = textbox.text_frame

            template_fmt = self.template_formats.get("body", {})

            for i, block in enumerate(text_blocks):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = block.text
                p.level = block.level

                if template_fmt:
                    self._apply_template_format_to_paragraph(p, template_fmt)
            
            # 自动调整文本框大小以适应内容，避免文本溢出
            self._auto_fit_textbox(textbox)
        except Exception:
            pass

    def _add_image(self, slide, content: dict) -> None:
        try:
            blob = content.get("blob")
            if not blob:
                return
            stream = BytesIO(blob)
            slide.shapes.add_picture(
                stream,
                left=content.get("left") or Emu(914400),
                top=content.get("top") or Emu(914400),
                width=content.get("width"),
                height=content.get("height"),
            )
        except Exception:
            pass

    def _add_table_from_content(self, slide, content) -> None:
        """从ContentBlock.content添加表格"""
        if isinstance(content, dict) and "cells" in content:
            shape_data = {"data": content, "left": Emu(914400), "top": Emu(914400 * 3)}
            self._add_table_full(slide, shape_data)
        elif isinstance(content, list):
            shape_data = {"data": content, "left": Emu(914400), "top": Emu(914400 * 3)}
            self._add_table_simple(slide, shape_data, content)

    def _add_table(self, slide, data: list[list[str]]) -> None:
        shape_data = {"data": data}
        self._add_table_simple(slide, shape_data, data)

    def _check_and_fix_overflow(self, slide) -> None:
        try:
            slide_width = self.target_width
            slide_height = self.target_height
        except Exception:
            slide_width = Emu(12192000)
            slide_height = Emu(6858000)

        for shape in slide.shapes:
            try:
                # 跳过满屏背景图
                if (shape.shape_type == 13 and  # MSO_SHAPE_TYPE.PICTURE
                    shape.left == 0 and shape.top == 0 and
                    abs(shape.width - slide_width) < Emu(100000) and
                    abs(shape.height - slide_height) < Emu(100000)):
                    continue

                if shape.left < 0:
                    shape.left = Emu(91440)
                if shape.top < 0:
                    shape.top = Emu(91440)
                if shape.left + shape.width > slide_width + Emu(10000):
                    shape.width = slide_width - shape.left - Emu(91440)
                if shape.top + shape.height > slide_height + Emu(10000):
                    shape.height = slide_height - shape.top - Emu(91440)
            except Exception:
                pass

    def _determine_layout(self, model: SlideContentModel) -> str:
        """根据内容模型和布局特征确定布局类型
        
        优先级：
        1. original_layout_type（如果是封面/章节/结尾等特殊布局）
        2. layout_features 中的特征（表格/图片/多列等）
        3. 默认内容页布局
        
        注意：普通内容页的 original_layout_type 应该是 "content"，
        不要因为布局名称包含 "title" 就误判为 "cover"
        """
        # 特殊布局类型：封面、章节、结尾、议程 - 保留原始布局
        special_layouts = {"cover", "section", "closing", "agenda"}
        if model.original_layout_type and model.original_layout_type in special_layouts:
            return model.original_layout_type
        
        # 根据布局特征判断
        features = model.layout_features or {}
        
        # 多列布局
        if features.get("columns", 1) >= 2:
            return "two_column"
        
        # 表格页
        if features.get("has_table"):
            return "table"
        
        # 图片页（有图片且文字少）
        if features.get("has_image") and features.get("text_density", 1) < 0.3:
            return "image"
        
        # 默认内容页
        return "content"

    def _remove_watermarks_from_slide(self, slide):
        """删除水印文字和形状"""
        from core.watermark.detector import WatermarkDetector
        detector = WatermarkDetector()
        
        indices_to_remove = []
        
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                result = detector.detect_text(text, 0)
                if result and result.detected:
                    indices_to_remove.append(i)
        
        for i in reversed(indices_to_remove):
            shape = slide.shapes[i]
            sp = shape._element
            sp.getparent().remove(sp)

    def _remove_original_footer_from_slide(self, slide):
        """删除原PPT的页脚占位符及页脚区域的所有内容"""
        footer_threshold = self.target_height * 0.88
        
        indices_to_remove = []
        
        for i, shape in enumerate(slide.shapes):
            top = shape.top or 0
            if top > footer_threshold:
                indices_to_remove.append(i)
        
        for i in reversed(indices_to_remove):
            shape = slide.shapes[i]
            sp = shape._element
            sp.getparent().remove(sp)

    def _apply_background_to_slide(self, slide):
        """应用背景色或渐变 - 技术适配模式

        技术适配原则：直接在幻灯片上应用背景颜色，确保正确显示
        
        问题：幻灯片使用的是原PPT的母版，而不是新复制的模板母版
        解决方案：直接设置幻灯片背景颜色，而不是依赖母版继承
        """
        try:
            bg = slide.background
            fill = bg.fill
            
            if self.background_color:
                if self.background_color.lower() == 'gradient':
                    # 应用渐变背景
                    try:
                        grad_config = getattr(self, 'background_gradient', None)
                        if grad_config:
                            fill.gradient()
                            stops = fill.gradient_stops
                            start_color = grad_config.get('start_color', '1A237E')
                            end_color = grad_config.get('end_color', '0D47A1')
                            if len(stops) >= 2:
                                stops[0].color.rgb = RGBColor.from_string(start_color)
                                stops[1].color.rgb = RGBColor.from_string(end_color)
                        else:
                            # 默认渐变
                            fill.gradient()
                            stops = fill.gradient_stops
                            if len(stops) >= 2:
                                stops[0].color.rgb = RGBColor(0x0A, 0x2F, 0x24)
                                stops[1].color.rgb = RGBColor(0x3D, 0xCD, 0x58)
                    except Exception:
                        # 渐变失败时回退到纯色
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(0x0A, 0x2F, 0x24)
                else:
                    fill.solid()
                    try:
                        r, g, b = int(self.background_color[0:2], 16), int(self.background_color[2:4], 16), int(self.background_color[4:6], 16)
                        fill.fore_color.rgb = RGBColor(r, g, b)
                    except Exception:
                        try:
                            fill.fore_color.theme_color = getattr(MSO_THEME_COLOR_INDEX, 'BACKGROUND_1', 0)
                        except Exception:
                            pass
        except Exception:
            pass

    def _get_color_from_theme(self, theme_color_name: str) -> str:
        """从主题文件中提取实际的RGB颜色值
        
        Args:
            theme_color_name: 主题颜色名称，如 'bg1', 'bg2', 'dk1', 'accent1' 等
        
        Returns:
            RGB颜色值，如 '0A2F24'，如果未找到则返回空字符串
        """
        if not self.template_path:
            return ""
        
        try:
            import zipfile
            from lxml import etree
            
            ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            
            with zipfile.ZipFile(self.template_path, 'r') as zf:
                # 确定主题文件路径
                theme_file = None
                
                # 如果有母版关联的主题文件路径，使用它
                if self.template_theme_path:
                    # 处理相对路径，如 '../theme/theme3.xml' -> 'ppt/theme/theme3.xml'
                    if self.template_theme_path.startswith('../'):
                        theme_file = f'ppt/{self.template_theme_path[3:]}'
                    else:
                        theme_file = f'ppt/{self.template_theme_path}'
                
                # 如果没有，查找第一个主题文件
                if not theme_file or theme_file not in zf.namelist():
                    theme_files = [f for f in zf.namelist() if f.startswith('ppt/theme/theme') and f.endswith('.xml')]
                    if theme_files:
                        theme_file = theme_files[0]
                
                if not theme_file or theme_file not in zf.namelist():
                    return ""
                
                # 读取主题文件
                theme_xml = zf.read(theme_file)
                theme_elem = etree.fromstring(theme_xml)
                
                # 查找颜色方案
                clrScheme = theme_elem.find(f'.//{{{ns_a}}}clrScheme', namespaces={'a': ns_a})
                if clrScheme is None:
                    return ""
                
                # 查找对应的颜色定义
                color_map = {
                    'bg1': 'dk1',
                    'bg2': 'lt1',
                    'dk1': 'dk1',
                    'dk2': 'dk2',
                    'lt1': 'lt1',
                    'lt2': 'lt2',
                }
                
                target_name = color_map.get(theme_color_name.lower(), theme_color_name.lower())
                
                color_elem = clrScheme.find(f'.//{{{ns_a}}}{target_name}', namespaces={'a': ns_a})
                if color_elem is not None:
                    # 先查找直接子元素
                    srgbClr = color_elem.find(f'{{{ns_a}}}srgbClr', namespaces={'a': ns_a})
                    if srgbClr is None:
                        # 如果没有找到，查找后代元素
                        srgbClr = color_elem.find(f'.//{{{ns_a}}}srgbClr', namespaces={'a': ns_a})
                    if srgbClr is not None:
                        return srgbClr.get('val', '')
            
            return ""
        except Exception:
            return ""

    def _unify_fonts_and_colors_on_slide(self, slide):
        """统一字体和颜色"""
        if not self.template_title_font and not self.template_body_font:
            return
        
        title_font = self.template_title_font or "Arial"
        body_font = self.template_body_font or "Arial"
        
        for shape in slide.shapes:
            if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                continue
            
            tf = shape.text_frame
            
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        font_size = run.font.size
                        is_bold = run.font.bold
                        
                        is_title = font_size and font_size >= Pt(24) or is_bold
                        
                        if is_title:
                            run.font.name = title_font
                        else:
                            run.font.name = body_font
                        
                        run.font._element.set('eastAsian', run.font.name)
                        
                        if is_title and self.title_text_color:
                            self._set_font_color_on_run(run, self.title_text_color)
                        elif self.default_text_color:
                            self._set_font_color_on_run(run, self.default_text_color)

    def _set_font_color_on_run(self, run, hex_color: str):
        """设置字体颜色"""
        try:
            run.font.color.rgb = RGBColor.from_string(hex_color.lstrip("#"))
        except:
            pass

    def convert_with_classification(
        self,
        source_path: str,
        output_path: str,
        background_style: str = "dark_green",
    ):
        """
        Hybrid mode conversion main entry point.
        Returns (output_path, qa_items).
        """
        import shutil
        import yaml
        from pathlib import Path
        from pptx import Presentation
        from core.classifier.slide_classifier import SlideClassifier
        from core.migrator.slide_migrator import SlideMigrator
        from core.models import QAReportItem

        # Load style config
        config_path = Path(__file__).parent.parent.parent / "config" / "master_styles.yaml"
        with open(config_path, "r", encoding="utf-8") as f:
            styles_config = yaml.safe_load(f)

        style_cfg = styles_config.get("master_styles", {}).get(background_style, {})
        master_idx = style_cfg.get("master_index", 2)

        # 1. Classify all slides
        source_prs = Presentation(source_path)
        classifier = SlideClassifier()
        classifications = classifier.classify_all(source_prs)

        # 2. Create output from template
        shutil.copy(self.template_path, output_path)
        target_prs = Presentation(output_path)

        # Remove existing slides from template
        while len(target_prs.slides._sldIdLst) > 0:
            rId = target_prs.slides._sldIdLst[0].rId
            target_prs.part.drop_rel(rId)
            del target_prs.slides._sldIdLst[0]

        migrator = SlideMigrator(self.template_path, master_idx)
        qa_items = []

        for idx, cls in enumerate(classifications):
            source_slide = source_prs.slides[idx]

            if cls.migration_mode == "migration":
                # Migration path: create new slide from template layout
                new_slide = migrator.migrate_slide(
                    source_slide=source_slide,
                    target_prs=target_prs,
                    slide_type=cls.slide_type,
                    layout_index=cls.target_layout_index,
                )
                # Apply font normalization
                replaced_fonts = self._normalize_fonts(new_slide)

                qa_item = QAReportItem(
                    slide_no=idx + 1,
                    detected_type=cls.slide_type,
                    applied_layout=cls.target_layout_name,
                    migration_mode="migration",
                    font_replaced=" → ".join(replaced_fonts) if replaced_fonts else "",
                    objects_moved=0,
                    objects_deleted=0,
                    overflow_risk="None",
                    need_manual_review=False,
                    comment=f"Migration: {cls.slide_type}",
                )
            else:
                # Adaptation path: create new slide and copy content
                new_slide = self._adapt_single_slide(
                    source_slide, target_prs, master_idx, cls
                )
                replaced_fonts = self._normalize_fonts(new_slide)

                qa_item = QAReportItem(
                    slide_no=idx + 1,
                    detected_type=cls.slide_type,
                    applied_layout=cls.target_layout_name,
                    migration_mode="adaptation",
                    font_replaced=" → ".join(replaced_fonts) if replaced_fonts else "",
                    objects_moved=0,
                    objects_deleted=0,
                    overflow_risk="Low",
                    need_manual_review=True,
                    comment=f"Adaptation: {cls.slide_type}，建议人工检查",
                )

            qa_items.append(qa_item)

        target_prs.save(output_path)
        return output_path, qa_items

    def _adapt_single_slide(self, source_slide, target_prs, master_idx, classification):
        """Adaptation path: use specified layout + position matching + overflow adjustment."""
        from core.migrator.object_migrator import ObjectMigrator
        from core.migrator.position_matcher import PositionMatcher
        from core.migrator.overflow_adjuster import OverflowAdjuster
        from core.migrator.coordinate_mapper import CoordinateMapper

        master = target_prs.slide_masters[min(master_idx, len(target_prs.slide_masters) - 1)]
        layout_idx = min(classification.target_layout_index, len(master.slide_layouts) - 1)
        target_layout = master.slide_layouts[layout_idx]
        new_slide = target_prs.slides.add_slide(target_layout)

        self._clear_placeholder_defaults(new_slide)

        bg_dark = master_idx == 2
        text_color = "#FFFFFF" if bg_dark else "#0A2F24"

        tgt_width = target_prs.slide_width
        tgt_height = target_prs.slide_height

        try:
            src_prs = source_slide.part.package.presentation_part.presentation
            src_width = src_prs.slide_width
            src_height = src_prs.slide_height
        except Exception:
            src_width = tgt_width
            src_height = tgt_height

        position_matcher = PositionMatcher(target_prs, master_idx)
        overflow_adjuster = OverflowAdjuster(tgt_width, tgt_height)
        coordinate_mapper = CoordinateMapper(src_width, src_height, tgt_width, tgt_height)

        object_migrator = ObjectMigrator(
            text_color=text_color,
            bg_dark=bg_dark,
            position_matcher=position_matcher,
            overflow_adjuster=overflow_adjuster,
            coordinate_mapper=coordinate_mapper
        )
        object_migrator.migrate_objects(source_slide, new_slide, tgt_width, tgt_height)

        return new_slide

    def _clear_placeholder_defaults(self, slide):
        """Clear default placeholder text like 'Click to edit master title style'."""
        for ph in slide.placeholders:
            if ph.has_text_frame:
                tf = ph.text_frame
                for para in tf.paragraphs:
                    para.text = ""

    def _normalize_fonts(self, slide) -> set:
        """
        Normalize fonts on a slide. Replace old fonts with target fonts.
        Returns set of replaced font names.
        """
        import yaml
        from pathlib import Path

        replaced = set()

        config_path = Path(__file__).parent.parent.parent / "config" / "font_mapping.yaml"
        if not config_path.exists():
            return replaced

        with open(config_path, "r", encoding="utf-8") as f:
            font_config = yaml.safe_load(f)

        # Build replacement map
        replace_map = {}
        for group in ["chinese", "english"]:
            cfg = font_config.get("font_replacements", {}).get(group, {})
            target = cfg.get("target_font", "Poppins")
            for src in cfg.get("source_fonts", []):
                replace_map[src.lower()] = target

        # Normalize all shapes
        for shape in slide.shapes:
            self._normalize_shape_fonts(shape, replace_map, replaced)

        return replaced

    def _normalize_shape_fonts(self, shape, replace_map: dict, replaced: set):
        """Recursively normalize fonts in a shape."""
        if not hasattr(shape, 'has_text_frame'):
            return
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        old_font = run.font.name
                        new_font = replace_map.get(old_font.lower())
                        if new_font and old_font != new_font:
                            run.font.name = new_font
                            replaced.add(old_font)

        if hasattr(shape, 'has_table') and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                old_font = run.font.name
                                new_font = replace_map.get(old_font.lower())
                                if new_font and old_font != new_font:
                                    run.font.name = new_font
                                    replaced.add(old_font)
