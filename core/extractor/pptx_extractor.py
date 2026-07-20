from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from core.models import SlideContentModel, ContentBlock, WatermarkReport, TextFormat, ShapeFormat
from core.watermark.detector import WatermarkDetector


class PptxExtractor:
    """PPT 内容抽取器：从源 PPT 抽取内容模型，保留原始格式，自动过滤水印"""

    def __init__(self):
        self.watermark_detector = WatermarkDetector()

    def extract(self, pptx_path: str) -> list[SlideContentModel]:
        if not Path(pptx_path).exists():
            raise FileNotFoundError(f"File not found: {pptx_path}")

        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        models = []

        for idx, slide in enumerate(prs.slides):
            model = self._extract_slide(slide, idx, total_slides)
            models.append(model)

        return models

    def _extract_slide(self, slide, slide_index: int, total_slides: int = 0) -> SlideContentModel:
        title = None
        title_source = {}
        title_format = None
        body_blocks = []
        notes_text = None
        raw_shapes = []
        # shape_id 映射表：id(shape) → raw_shapes 中的 shape_id
        # 用于建立 ContentBlock 与 raw_shape 的精确映射
        shape_id_map = {}
        next_shape_id = 0
        body_placeholder_count = 0
        text_shapes_positions = []
        layout_features = {
            "has_title": False,
            "has_subtitle": False,
            "has_table": False,
            "has_image": False,
            "has_chart": False,
            "columns": 1,
            "text_density": 0.0,
            "shape_count": 0,
            "placeholder_count": 0,
            "has_multi_body": False,
        }

        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip() or None
        except Exception:
            pass

        slide_height = slide.part.package.presentation_part.presentation.slide_height
        slide_width = slide.part.package.presentation_part.presentation.slide_width
        # 页眉区域：顶部8%
        header_threshold = slide_height * 0.08
        # 页脚区域：底部12%
        footer_threshold = slide_height * 0.88

        def _is_header_or_footer(shape) -> bool:
            """判断形状是否在页眉或页脚区域
            
            判断逻辑：
            1. 占位符类型为 title/body/subtitle 的不算页眉页脚
            2. 完全在顶部8%区域 → 页眉
            3. 完全在底部12%区域 → 页脚（无论文本内容如何，只要不是正文占位符）
            4. 注意：用户要求舍弃原PPT的所有页眉页脚，只用模板的
            """
            try:
                top = shape.top or 0
                height = shape.height or 0
                bottom = top + height
                # 标题/正文/副标题占位符不算页眉页脚
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # 1=title, 2=body, 3=ctrTitle, 4=subTitle, 7=text
                    if ph_type in (1, 2, 3, 4, 7):
                        return False
                    # 页脚类型的占位符：13=SLIDE_NUMBER, 14=HEADER, 15=FOOTER, 16=DATE
                    if ph_type in (13, 14, 15, 16):
                        return True
                # 完全在顶部8%区域 → 页眉
                if bottom < header_threshold:
                    return True
                # 完全在底部12%区域 → 页脚
                # 注意：无论文本内容如何，只要在底部12%且不是正文占位符，都视为页脚
                # 因为用户要求完全舍弃原PPT的页眉页脚，只用模板的
                if top > footer_threshold:
                    return True
                return False
            except Exception:
                return False

        def _infer_semantic_role(shape, slide_width: int) -> str:
            """根据形状的占位符类型和位置推断语义角色"""
            try:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    ph_idx = shape.placeholder_format.idx
                    # 1=title, 3=ctrTitle
                    if ph_type in (1, 3):
                        return "title"
                    # 4=subtitle
                    if ph_type == 4:
                        return "subtitle"
                    # 2=body, 7=text
                    if ph_type in (2, 7):
                        # 根据位置判断主正文还是侧边栏
                        left = shape.left or 0
                        if left > slide_width * 0.5:
                            return "body_sidebar"
                        return "body_main"
                    # 其他占位符
                    return "unknown"
                
                # 非占位符形状：根据位置和类型推断
                top = shape.top or 0
                if top < slide_height * 0.15:
                    return "title"
                return "body_main"
            except Exception:
                return "unknown"

        for shape in slide.shapes:
            layout_features["shape_count"] += 1
            
            # 跳过页眉/页脚区域的形状（不抽取原PPT的页眉页脚，完全继承模板的）
            if _is_header_or_footer(shape):
                continue

            # 为每个形状分配唯一的 shape_id
            current_shape_id = next_shape_id
            next_shape_id += 1
            shape_id_map[id(shape)] = current_shape_id

            shape_data = self._extract_shape(shape, slide_index)
            if shape_data:
                # 在 shape_data 中记录 shape_id
                shape_data["shape_id"] = current_shape_id
                raw_shapes.append(shape_data)

            # 统计占位符
            if shape.is_placeholder:
                layout_features["placeholder_count"] += 1
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (2, 7):
                        body_placeholder_count += 1
                        text_shapes_positions.append(shape.left or 0)
                except Exception:
                    pass

            if shape == slide.shapes.title:
                # 标题已提取为 model.title，从 raw_shapes 中移除，避免回填时重复
                if raw_shapes and raw_shapes[-1].get("shape_id") == current_shape_id:
                    raw_shapes.pop()
                title_text = self._get_shape_text(shape)
                if title_text:
                    watermark_report = self.watermark_detector.detect_text(title_text, slide_index)
                    if not watermark_report.detected:
                        # 清理标题中可能包含的水印关键词（即使不判定为水印段落）
                        title = self.watermark_detector.clean_text(title_text)
                        layout_features["has_title"] = True
                        # 记录标题来源占位符信息
                        try:
                            title_source = {
                                "placeholder_type": shape.placeholder_format.type,
                                "placeholder_idx": shape.placeholder_format.idx,
                                "left": shape.left,
                                "top": shape.top,
                                "width": shape.width,
                                "height": shape.height,
                            }
                        except Exception:
                            title_source = {}
                        # 提取标题格式信息（含继承解析）
                        title_format = self._extract_title_format(shape)
                continue

            if shape.has_text_frame:
                # 有填充色的自选图形（内容框）：不提取为body_block，保留形状整体
                # 这样可以保留背景色，用于color pairing
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        if shape.fill.type == 1:  # solid fill
                            continue
                    except Exception:
                        pass
                
                # 推断语义角色
                semantic_role = _infer_semantic_role(shape, slide_width)
                if semantic_role == "subtitle":
                    layout_features["has_subtitle"] = True
                
                blocks = self._extract_text_blocks(shape, slide_index, semantic_role, shape)
                # 为每个 block 关联 raw_shape_id
                for block in blocks:
                    block.raw_shape_id = current_shape_id
                body_blocks.extend(blocks)
                continue

            if shape.has_table:
                table_data = self._extract_table_full(shape.table)
                if table_data:
                    layout_features["has_table"] = True
                    body_blocks.append(ContentBlock(
                        type="table",
                        text=None,
                        content=table_data,
                        level=0,
                        semantic_role="body_main",
                        source_shape_id=id(shape),
                        raw_shape_id=current_shape_id,
                    ))
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    layout_features["has_image"] = True
                    body_blocks.append(ContentBlock(
                        type="image",
                        text=None,
                        content={
                            "blob": image.blob,
                            "ext": image.ext,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                        },
                        level=0,
                        semantic_role="body_main",
                        source_shape_id=id(shape),
                        raw_shape_id=current_shape_id,
                    ))
                except Exception:
                    pass
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                blocks = self._extract_group(shape, slide_index)
                # 为 group 中的 block 也关联 raw_shape_id
                for block in blocks:
                    block.raw_shape_id = current_shape_id
                body_blocks.extend(blocks)
                continue

        # 分析布局特征
        if body_placeholder_count > 1:
            layout_features["has_multi_body"] = True
            # 根据位置分布判断列数
            if len(text_shapes_positions) >= 2:
                text_shapes_positions.sort()
                # 如果有两个文本框，一个在左半部分，一个在右半部分，则为两列布局
                mid_point = slide_width / 2
                left_count = sum(1 for p in text_shapes_positions if p < mid_point)
                right_count = sum(1 for p in text_shapes_positions if p >= mid_point)
                if left_count > 0 and right_count > 0:
                    layout_features["columns"] = 2
        
        # 计算文本密度（文本形状数量 / 总形状数量）
        if layout_features["shape_count"] > 0:
            text_shape_count = sum(1 for b in body_blocks if b.type == "paragraph")
            layout_features["text_density"] = text_shape_count / layout_features["shape_count"]

        # 过滤掉与text形状完全重叠且无文本的autoshape（这些是文本框的背景矩形）
        # 问题：原PPT中可能有一个有填充色的矩形作为文本框的背景，两者完全重叠
        # 转换后这个矩形会遮挡文本内容
        raw_shapes = self._filter_overlapping_shapes(raw_shapes)

        # 最终防御：用 locals() 确保所有变量都已绑定，避免 UnboundLocalError
        _locals = locals()
        _title = _locals.get("title")
        _title_source = _locals.get("title_source", {})
        _title_format = _locals.get("title_format")
        _body_blocks = _locals.get("body_blocks", [])
        _notes_text = _locals.get("notes_text")
        _raw_shapes = _locals.get("raw_shapes", [])
        _layout_features = _locals.get("layout_features", {})

        # 如果没有标准标题占位符，从body_blocks中提取semantic_role="title"的第一个块作为标题
        if not _title and _body_blocks:
            title_blocks = [b for b in _body_blocks if b.semantic_role == "title" and b.text and b.type == "paragraph"]
            if title_blocks:
                first_title_block = title_blocks[0]
                _title = first_title_block.text
                _title_format = first_title_block.text_format
                _title_source = {
                    "from_body_block": True,
                    "semantic_role": "title",
                }
                _layout_features["has_title"] = True
                _body_blocks = [b for b in _body_blocks if b is not first_title_block]

        return SlideContentModel(
            slide_index=slide_index,
            title=_title,
            body_blocks=_body_blocks,
            notes=_notes_text,
            original_layout_type=self._detect_layout_type(slide, slide_index, total_slides),
            raw_shapes=_raw_shapes,
            layout_features=_layout_features,
            title_source=_title_source,
            title_format=_title_format,
        )

    def _extract_shape(self, shape, slide_index: int) -> dict | None:
        try:
            shape_type = shape.shape_type

            # 优先判断：自选图形（有填充色/边框）即使有文字，也当作autoshape处理
            # 这样可以保留内容框的背景形状，用于 color pairing
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                has_fill = False
                try:
                    if shape.fill.type == 1:  # solid fill
                        has_fill = True
                except Exception:
                    pass
                if has_fill:
                    return self._extract_auto_shape(shape)

            if shape.has_text_frame:
                text_result = self._extract_text_shape(shape, slide_index)
                if text_result:
                    return text_result
                # 如果文本形状没有内容，回退到提取为autoshape（保留几何形状）
                if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    return self._extract_auto_shape(shape)
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                return self._extract_image_shape(shape)
            elif shape.has_table:
                return self._extract_table_shape(shape)
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                # 检查是否是SmartArt（diagram）
                is_smartart = False
                try:
                    xml = shape._element.xml
                    if 'dgm:' in xml or 'diagram' in xml.lower() or 'smartArt' in xml.lower():
                        is_smartart = True
                except Exception:
                    pass
                
                if is_smartart:
                    return self._extract_smartart_shape(shape, slide_index)
                else:
                    return self._extract_group_shape(shape, slide_index)
            elif shape_type == MSO_SHAPE_TYPE.CHART:
                return self._extract_chart_shape(shape)
            elif shape_type == MSO_SHAPE_TYPE.OLE_OBJECT:
                return self._extract_ole_shape(shape)
            else:
                # 尝试提取自选图形的几何信息
                return self._extract_auto_shape(shape)
        except Exception:
            return None

    def _filter_overlapping_shapes(self, raw_shapes: list[dict]) -> list[dict]:
        """过滤掉与text形状完全重叠且无填充、无文本的autoshape
        
        问题：原PPT中可能有一个无填充的矩形作为文本框的背景，两者完全重叠
        转换后这个矩形可能是多余的
        
        判断规则：
        1. 找到所有text类型的形状
        2. 找到所有autoshape类型的形状，且无文本、无填充色、无描边
        3. 如果一个autoshape与一个text形状位置和大小完全相同，则移除该autoshape
        4. 注意：有填充色的autoshape是设计元素（如内容卡片背景），必须保留
        """
        if not raw_shapes:
            return raw_shapes
        
        # 找出所有text形状的位置
        text_shapes = []
        for shape in raw_shapes:
            if shape.get("type") == "text":
                left = shape.get("left", 0)
                top = shape.get("top", 0)
                width = shape.get("width", 0)
                height = shape.get("height", 0)
                text_shapes.append((left, top, width, height))
        
        if not text_shapes:
            return raw_shapes
        
        # 过滤掉与text形状重叠的无填充、无文本、无描边的autoshape
        filtered = []
        tolerance = 1000  # 容差，单位EMU
        
        for shape in raw_shapes:
            shape_type = shape.get("type")
            if shape_type == "autoshape":
                # 检查是否有文本
                has_text = False
                paragraphs = shape.get("paragraphs", [])
                for para in paragraphs:
                    if para.get("text", "").strip():
                        has_text = True
                        break
                if shape.get("text", "").strip():
                    has_text = True
                
                # 检查是否有填充色
                has_fill = bool(shape.get("fill_color"))
                
                # 检查是否有描边
                has_line = bool(shape.get("line_color") and shape.get("line_width") and shape.get("line_width") > 0)
                
                # 有填充、有文本或有描边的autoshape：保留（是设计元素）
                if has_fill or has_text or has_line:
                    filtered.append(shape)
                    continue
                
                # 无填充、无文本、无描边的autoshape：检查是否与text形状重叠
                left = shape.get("left", 0)
                top = shape.get("top", 0)
                width = shape.get("width", 0)
                height = shape.get("height", 0)
                
                is_overlapping = False
                for t_left, t_top, t_width, t_height in text_shapes:
                    if (abs(left - t_left) < tolerance and
                        abs(top - t_top) < tolerance and
                        abs(width - t_width) < tolerance and
                        abs(height - t_height) < tolerance):
                        is_overlapping = True
                        break
                
                if is_overlapping:
                    # 与text形状完全重叠的空autoshape，跳过
                    continue
            
            filtered.append(shape)
        
        return filtered

    def _extract_text_shape(self, shape, slide_index: int) -> dict:
        paragraphs = []
        tf = shape.text_frame

        for paragraph in tf.paragraphs:
            text = paragraph.text.strip()
            if text:
                watermark_report = self.watermark_detector.detect_text(text, slide_index)
                if watermark_report.detected:
                    continue

            para_data = {
                "text": self.watermark_detector.clean_text(paragraph.text),
                "level": paragraph.level or 0,
                "alignment": str(paragraph.alignment) if paragraph.alignment else None,
                "runs": [],
                "line_spacing": None,
                "space_before": None,
                "space_after": None,
            }

            try:
                if paragraph.line_spacing is not None:
                    para_data["line_spacing"] = paragraph.line_spacing
                if paragraph.space_before is not None:
                    para_data["space_before"] = paragraph.space_before
                if paragraph.space_after is not None:
                    para_data["space_after"] = paragraph.space_after
            except Exception:
                pass

            for run in paragraph.runs:
                # 清理水印关键词
                run_text = self.watermark_detector.clean_text(run.text)
                run_data = {
                    "text": run_text,
                    "font_name": run.font.name,
                    "font_size": run.font.size,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "underline": run.font.underline,
                    "color": None,
                }
                try:
                    if run.font.color and run.font.color.rgb:
                        run_data["color"] = str(run.font.color.rgb)
                except Exception:
                    pass
                para_data["runs"].append(run_data)

            paragraphs.append(para_data)

        if not paragraphs:
            return None

        fill_color = None
        try:
            fill = shape.fill
            if fill.type == 1:
                if fill.fore_color and fill.fore_color.rgb:
                    fill_color = str(fill.fore_color.rgb)
        except Exception:
            pass

        line_color = None
        line_width = None
        try:
            line = shape.line
            if line.color and line.color.rgb:
                line_color = str(line.color.rgb)
            if line.width:
                line_width = line.width
        except Exception:
            pass

        return {
            "type": "text",
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "paragraphs": paragraphs,
            "shape_name": shape.name,
            "fill_color": fill_color,
            "line_color": line_color,
            "line_width": line_width,
        }

    def _extract_image_shape(self, shape) -> dict:
        try:
            image = shape.image
            return {
                "type": "image",
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "blob": image.blob,
                "ext": image.ext,
                "crop_left": getattr(shape, "crop_left", None),
                "crop_top": getattr(shape, "crop_top", None),
                "crop_right": getattr(shape, "crop_right", None),
                "crop_bottom": getattr(shape, "crop_bottom", None),
            }
        except Exception:
            return None

    def _extract_table_shape(self, shape) -> dict:
        try:
            table_data = self._extract_table_full(shape.table)
            return {
                "type": "table",
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "data": table_data,
            }
        except Exception:
            return None

    def _extract_table_full(self, table) -> dict:
        """完整提取表格：单元格文本、背景色、字体格式、合并单元格"""
        cells = []
        col_widths = []
        row_heights = []

        try:
            for col_idx in range(len(table.columns)):
                col_widths.append(table.columns[col_idx].width)
        except Exception:
            pass

        try:
            for row_idx in range(len(table.rows)):
                row_heights.append(table.rows[row_idx].height)
        except Exception:
            pass

        merged_cells = []
        for r, row in enumerate(table.rows):
            for c, cell in enumerate(row.cells):
                cell_data = {
                    "row": r,
                    "col": c,
                    "text": self.watermark_detector.clean_text(cell.text.strip()),
                    "fill_color": None,
                    "font_name": None,
                    "font_size": None,
                    "bold": None,
                    "italic": None,
                    "color": None,
                    "alignment": None,
                    "vertical_anchor": None,
                    "margin_left": None,
                    "margin_right": None,
                    "margin_top": None,
                    "margin_bottom": None,
                }

                try:
                    fill = cell.fill
                    if fill.type == 1:  # solid
                        if fill.fore_color and fill.fore_color.rgb:
                            cell_data["fill_color"] = str(fill.fore_color.rgb)
                except Exception:
                    pass

                try:
                    if cell.vertical_anchor is not None:
                        cell_data["vertical_anchor"] = str(cell.vertical_anchor)
                    cell_data["margin_left"] = cell.margin_left
                    cell_data["margin_right"] = cell.margin_right
                    cell_data["margin_top"] = cell.margin_top
                    cell_data["margin_bottom"] = cell.margin_bottom
                except Exception:
                    pass

                try:
                    tf = cell.text_frame
                    if tf.paragraphs:
                        p = tf.paragraphs[0]
                        if p.alignment is not None:
                            cell_data["alignment"] = str(p.alignment)
                        if p.runs:
                            run = p.runs[0]
                            font = run.font
                            if font.name:
                                cell_data["font_name"] = font.name
                            if font.size:
                                cell_data["font_size"] = font.size
                            if font.bold is not None:
                                cell_data["bold"] = font.bold
                            if font.italic is not None:
                                cell_data["italic"] = font.italic
                            try:
                                if font.color and font.color.rgb:
                                    cell_data["color"] = str(font.color.rgb)
                            except Exception:
                                pass
                except Exception:
                    pass

                # 检测合并单元格
                try:
                    if cell.span_h > 1 or cell.span_v > 1:
                        merged_cells.append({
                            "start_row": r,
                            "start_col": c,
                            "row_span": cell.span_v,
                            "col_span": cell.span_h,
                        })
                except Exception:
                    pass

                cells.append(cell_data)

        return {
            "cells": cells,
            "rows": len(table.rows),
            "cols": len(table.columns),
            "col_widths": col_widths,
            "row_heights": row_heights,
            "merged_cells": merged_cells,
        }

    def _extract_group_shape(self, group, slide_index: int) -> dict:
        shapes_data = []
        for shape in group.shapes:
            shape_data = self._extract_shape(shape, slide_index)
            if shape_data:
                shapes_data.append(shape_data)
        return {
            "type": "group",
            "left": group.left,
            "top": group.top,
            "width": group.width,
            "height": group.height,
            "shapes": shapes_data,
        }

    def _extract_smartart_shape(self, shape, slide_index: int) -> dict:
        """提取SmartArt图形
        
        SmartArt在python-pptx中以group形式存在，内部包含多个自选图形和文本框
        我们提取其内部的所有形状和文本，作为smartart类型存储
        
        回填时，以group形式重建，并应用颜色映射
        """
        shapes_data = []
        # 提取内部文本内容（用于内容识别）
        all_text_parts = []
        
        # 递归提取group中的所有形状
        def extract_shapes_recursive(shapes):
            for inner_shape in shapes:
                shape_data = self._extract_shape(inner_shape, slide_index)
                if shape_data:
                    shapes_data.append(shape_data)
                    # 收集文本
                    if shape_data.get("type") == "text":
                        for para in shape_data.get("paragraphs", []):
                            if para.get("text"):
                                all_text_parts.append(para["text"])
                    elif shape_data.get("type") == "autoshape" and shape_data.get("text"):
                        all_text_parts.append(shape_data["text"])
                # 递归处理子group
                try:
                    if hasattr(inner_shape, 'shapes') and inner_shape.shapes:
                        if inner_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                            extract_shapes_recursive(inner_shape.shapes)
                except Exception:
                    pass
        
        try:
            extract_shapes_recursive(shape.shapes)
        except Exception:
            pass
        
        return {
            "type": "smartart",
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "shapes": shapes_data,
            "all_text": " ".join(all_text_parts),
        }

    def _extract_chart_shape(self, shape) -> dict:
        """完整提取图表：数据、类别、系列颜色、图例、坐标轴"""
        try:
            chart = shape.chart
            chart_data = []
            categories = []
            series_colors = []
            chart_title = None
            has_legend = False
            legend_position = None

            try:
                if chart.has_title and chart.chart_title:
                    chart_title = self.watermark_detector.clean_text(chart.chart_title.text_frame.text)
            except Exception:
                pass

            try:
                if chart.has_legend:
                    has_legend = True
                    legend_position = str(chart.legend.position) if chart.legend else None
            except Exception:
                pass

            try:
                for series in chart.series:
                    series_data = {
                        "name": series.name,
                        "values": [],
                        "color": None,
                    }
                    try:
                        values = series.values
                        if values:
                            series_data["values"] = list(values)
                    except Exception:
                        for point in series.points:
                            try:
                                series_data["values"].append(point.value)
                            except Exception:
                                series_data["values"].append(None)

                    try:
                        if series.format.fill.type == 1:
                            if series.format.fill.fore_color and series.format.fill.fore_color.rgb:
                                series_data["color"] = str(series.format.fill.fore_color.rgb)
                    except Exception:
                        pass

                    chart_data.append(series_data)
            except Exception:
                pass

            try:
                if chart.plots and chart.plots[0].categories:
                    for cat in chart.plots[0].categories:
                        try:
                            categories.append(str(cat))
                        except Exception:
                            pass
            except Exception:
                pass

            if not categories and chart_data:
                max_len = max(len(s.get("values", [])) for s in chart_data) if chart_data else 0
                categories = [f"类别{i+1}" for i in range(max_len)]

            # 提取坐标轴信息
            x_axis = {}
            y_axis = {}
            try:
                if chart.category_axis:
                    ca = chart.category_axis
                    x_axis = {
                        "has_title": ca.has_title,
                        "title": self.watermark_detector.clean_text(ca.axis_title.text_frame.text) if ca.has_title and ca.axis_title else None,
                        "visible": ca.visible,
                    }
            except Exception:
                pass
            try:
                if chart.value_axis:
                    va = chart.value_axis
                    y_axis = {
                        "has_title": va.has_title,
                        "title": self.watermark_detector.clean_text(va.axis_title.text_frame.text) if va.has_title and va.axis_title else None,
                        "visible": va.visible,
                        "minimum_scale": va.minimum_scale,
                        "maximum_scale": va.maximum_scale,
                    }
            except Exception:
                pass

            return {
                "type": "chart",
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "chart_type": str(chart.chart_type),
                "data": chart_data,
                "categories": categories,
                "chart_title": chart_title,
                "has_legend": has_legend,
                "legend_position": legend_position,
                "x_axis": x_axis,
                "y_axis": y_axis,
            }
        except Exception:
            return None

    def _extract_ole_shape(self, shape) -> dict:
        try:
            ole_format = shape.ole_format
            if ole_format and ole_format.binary:
                return {
                    "type": "ole",
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "prog_id": ole_format.prog_id,
                    "blob": ole_format.binary,
                }
        except Exception:
            pass
        return None

    def _extract_auto_shape(self, shape) -> dict | None:
        """提取自选图形的几何信息和样式"""
        try:
            shape_type = shape.shape_type

            fill_color = None
            try:
                fill = shape.fill
                if fill.type == 1:
                    if fill.fore_color and fill.fore_color.rgb:
                        fill_color = str(fill.fore_color.rgb)
            except Exception:
                pass

            line_color = None
            line_width = None
            try:
                line = shape.line
                if line.color and line.color.rgb:
                    line_color = str(line.color.rgb)
                if line.width:
                    line_width = line.width
            except Exception:
                pass

            text = None
            if shape.has_text_frame:
                try:
                    text = self.watermark_detector.clean_text(shape.text_frame.text)
                except Exception:
                    pass

            return {
                "type": "autoshape",
                "shape_type": str(shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "rotation": getattr(shape, "rotation", 0),
                "fill_color": fill_color,
                "line_color": line_color,
                "line_width": line_width,
                "text": text,
            }
        except Exception:
            return None

    def _extract_text_blocks(self, shape, slide_index: int, semantic_role: str = "unknown", source_shape=None) -> list[ContentBlock]:
        blocks = []
        tf = shape.text_frame
        
        # 记录占位符信息
        ph_type = None
        ph_idx = None
        shape_id = id(source_shape) if source_shape else id(shape)
        if source_shape and source_shape.is_placeholder:
            try:
                ph_type = source_shape.placeholder_format.type
                ph_idx = source_shape.placeholder_format.idx
            except Exception:
                pass
        
        # 提取形状格式
        shape_format = self._extract_shape_format(source_shape or shape)
        
        for paragraph in tf.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            watermark_report = self.watermark_detector.detect_text(text, slide_index)
            if watermark_report.detected:
                continue
            
            # 清理文本中的水印关键词（即使整个段落不是水印，也要移除其中的水印词）
            cleaned_text = self.watermark_detector.clean_text(text)
            if not cleaned_text.strip():
                continue
            
            # 提取段落格式（取第一个run的格式，含继承解析）
            text_format = self._extract_text_format(paragraph, source_shape or shape)

            blocks.append(ContentBlock(
                type="paragraph",
                text=cleaned_text,
                level=paragraph.level or 0,
                semantic_role=semantic_role,
                original_placeholder_type=ph_type,
                original_placeholder_idx=ph_idx,
                source_shape_id=shape_id,
                text_format=text_format,
                shape_format=shape_format if blocks == [] else None,
            ))
        return blocks

    def _extract_title_format(self, shape) -> TextFormat | None:
        """从标题形状中提取文本格式信息（含继承解析）"""
        try:
            if not shape.has_text_frame:
                return None
            tf = shape.text_frame
            if not tf.paragraphs:
                return None
            return self._extract_text_format(tf.paragraphs[0], shape)
        except Exception:
            return None

    def _extract_text_format(self, paragraph, shape=None) -> TextFormat | None:
        """从段落中提取文本格式信息（含继承解析）
        
        python-pptx 的格式有继承机制（主题 → 布局 → 占位符 → run），
        run.font.size 为 None 时实际继承自上层。
        本方法递归解析继承值，获取实际渲染值。
        """
        try:
            text_format = TextFormat()
            
            # 获取第一个run的格式
            if paragraph.runs:
                run = paragraph.runs[0]
                font = run.font
                
                # 字体名称：run级别 → 占位符级别
                if font.name:
                    text_format.font_name = font.name
                elif shape:
                    # 尝试从占位符继承
                    try:
                        if shape.is_placeholder and shape.text_frame.paragraphs:
                            for p in shape.text_frame.paragraphs:
                                if p.runs:
                                    inherited_name = p.runs[0].font.name
                                    if inherited_name:
                                        text_format.font_name = inherited_name
                                        break
                    except Exception:
                        pass
                
                # 字号：run级别 → 段落默认 → 占位符级别
                if font.size:
                    text_format.font_size = font.size.pt
                elif shape:
                    # 尝试从占位符的其他段落继承
                    try:
                        if shape.is_placeholder and shape.text_frame.paragraphs:
                            for p in shape.text_frame.paragraphs:
                                if p.runs and p.runs[0].font.size:
                                    text_format.font_size = p.runs[0].font.size.pt
                                    break
                    except Exception:
                        pass
                
                # 粗体
                if font.bold is not None:
                    text_format.bold = font.bold
                
                # 斜体
                if font.italic is not None:
                    text_format.italic = font.italic
                
                # 下划线
                if font.underline is not None:
                    text_format.underline = font.underline
                
                # 颜色：run级别 → 主题颜色
                if font.color and font.color.rgb:
                    text_format.font_color = str(font.color.rgb)
                elif font.color and font.color.type is not None:
                    # 尝试获取主题颜色
                    try:
                        if font.color.theme_color:
                            # 主题颜色映射到近似RGB值
                            theme_color_map = {
                                'tx1': '000000', 'dk1': '000000',
                                'tx2': '44546A', 'dk2': '44546A',
                                'bg1': 'FFFFFF', 'lt1': 'FFFFFF',
                                'bg2': 'E7E6E6', 'lt2': 'E7E6E6',
                                'accent1': '4472C4', 'accent2': 'ED7D31',
                                'accent3': 'A5A5A5', 'accent4': 'FFC000',
                                'accent5': '5B9BD5', 'accent6': '70AD47',
                            }
                            tc_str = str(font.color.theme_color)
                            for key, val in theme_color_map.items():
                                if key in tc_str:
                                    text_format.font_color = val
                                    break
                    except Exception:
                        pass
            
            # 段落对齐
            if paragraph.alignment is not None:
                text_format.alignment = paragraph.alignment.value
            
            # 行距
            try:
                if paragraph.line_spacing is not None:
                    text_format.line_spacing = float(paragraph.line_spacing)
            except Exception:
                pass
            
            return text_format
        except Exception:
            return None

    def _extract_shape_format(self, shape) -> ShapeFormat | None:
        """从形状中提取形状格式信息"""
        try:
            shape_format = ShapeFormat()
            
            # 位置和大小
            if shape.left is not None:
                shape_format.left = shape.left
            if shape.top is not None:
                shape_format.top = shape.top
            if shape.width is not None:
                shape_format.width = shape.width
            if shape.height is not None:
                shape_format.height = shape.height
            
            # 旋转
            if shape.rotation is not None:
                shape_format.rotation = shape.rotation
            
            # 形状类型
            try:
                shape_format.shape_type = str(shape.shape_type)
            except Exception:
                pass
            
            # 填充
            try:
                fill = shape.fill
                if fill.type == 1:  # solid
                    shape_format.fill_type = "solid"
                    if fill.fore_color and fill.fore_color.rgb:
                        shape_format.fill_color = str(fill.fore_color.rgb)
                elif fill.type == 2:  # gradient
                    shape_format.fill_type = "gradient"
                elif fill.type == 3:  # picture
                    shape_format.fill_type = "picture"
                else:
                    shape_format.fill_type = "none"
            except Exception:
                pass
            
            # 线条
            try:
                line = shape.line
                if line.color and line.color.rgb:
                    shape_format.line_color = str(line.color.rgb)
                if line.width is not None:
                    shape_format.line_width = line.width
            except Exception:
                pass
            
            return shape_format
        except Exception:
            return None

    def _extract_group(self, group, slide_index: int) -> list[ContentBlock]:
        blocks = []
        for shape in group.shapes:
            if shape.has_text_frame:
                blocks.extend(self._extract_text_blocks(shape, slide_index))
        return blocks

    def _extract_table(self, table) -> list[list[str]]:
        """简化版表格提取（向后兼容）"""
        data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            data.append(row_data)
        return data

    def _get_shape_text(self, shape) -> str:
        try:
            return shape.text_frame.text.strip()
        except Exception:
            return ""

    def _detect_layout_type(self, slide, slide_index: int, total_slides: int = 0) -> str:
        """检测原PPT的布局类型
        
        注意：不能仅凭layout名称判断，因为很多内容页布局名称也包含"Title"
        需要结合：
        1. layout名称中的关键词（精确匹配封面/章节/结尾）
        2. 占位符类型和数量
        3. 内容特征（文字多少、是否有感谢语/议程关键词等）
        """
        try:
            layout_name = (slide.slide_layout.name or "").lower()
            
            # 精确匹配封面布局（只有封面才会叫这些名字）
            cover_keywords = [
                "封面", "cover", "title slide", "opening",
                "title slide simple", "title slide with image",
                "title slide offer",
            ]
            for kw in cover_keywords:
                if kw in layout_name:
                    return "cover"
            
            # 章节页
            if "章节" in layout_name or "section break" in layout_name or "section header" in layout_name:
                return "section"
            
            # 结尾页
            if "结尾" in layout_name or "closing" in layout_name or "thank you" in layout_name or "end slide" in layout_name:
                return "closing"
            
            # 议程页
            if "议程" in layout_name or "agenda" in layout_name or "table of contents" in layout_name or "目录" in layout_name:
                return "agenda"
        except Exception:
            pass
        
        # 根据内容特征判断
        try:
            # 收集所有文本内容（排除页眉页脚区域的）
            all_text = ""
            text_shape_count = 0
            total_shapes = len(slide.shapes)
            
            slide_height = slide.part.package.presentation_part.presentation.slide_height
            header_th = slide_height * 0.08
            footer_th = slide_height * 0.88
            
            for shape in slide.shapes:
                top = shape.top or 0
                height = shape.height or 0
                bottom = top + height
                # 跳过页眉页脚区域的形状
                if bottom < header_th or top > footer_th:
                    continue
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        all_text += text + "\n"
                        text_shape_count += 1
            
            all_text_lower = all_text.lower()
            text_lines = [l.strip() for l in all_text.split('\n') if l.strip()]
            text_line_count = len(text_lines)
            
            # 计算总字符数（用于判断内容密度）
            total_chars = len(all_text.replace('\n', '').replace(' ', ''))
            
            # ========== 封面识别 ==========
            if slide_index == 0:
                # 封面特征：有大标题、有副标题、内容少
                if text_line_count <= 10 and text_shape_count <= 8:
                    cover_indicators = [
                        "presenter", "演讲者", "by:", "date:", "日期", 
                        "v0.", "version", "团队", "team", "speaker",
                        "presented by", "作者", "制作"
                    ]
                    has_cover_indicator = any(kw in all_text_lower for kw in cover_indicators)
                    if has_cover_indicator or text_line_count <= 6:
                        return "cover"
            
            # ========== 结尾页识别 ==========
            is_last_page = (total_slides > 0 and slide_index == total_slides - 1)
            closing_keywords = [
                "thank you", "thanks", "谢谢", "感谢",
                "q&a", "q & a", "qa", "question", "问题",
                "the end", "结束", "结语", "总结",
                "contact", "联系我们", "联系方式",
                "next step", "下一步", "decision", "决策"
            ]
            has_closing_kw = any(kw in all_text_lower for kw in closing_keywords)
            
            if has_closing_kw and text_line_count <= 10:
                return "closing"
            
            # 最后一页且内容少 → 可能是结尾页
            if is_last_page and text_line_count <= 5:
                return "closing"
            
            # ========== 议程页识别 ==========
            agenda_title_keywords = ["agenda", "议程", "目录", "outline", "overview", "table of content"]
            title_text = " ".join(text_lines[:2]).lower()
            has_agenda_title = any(kw in title_text for kw in agenda_title_keywords)
            
            # 检查是否有编号列表（支持多种格式）
            has_numbered_list = False
            numbered_count = 0
            for line in text_lines:
                line_stripped = line.strip()
                if not line_stripped:
                    continue
                # 格式1: 数字+标点 (1. 2. 3. / 1、2、)
                if len(line_stripped) >= 2 and line_stripped[0].isdigit() and line_stripped[1] in '.、)':
                    numbered_count += 1
                # 格式2: 两位数字开头 (01 02 03)
                elif len(line_stripped) >= 2 and line_stripped[:2].isdigit():
                    numbered_count += 1
                # 格式3: 中文数字
                elif line_stripped.startswith(('一、', '二、', '三、', '四、', '五、', '六、', '七、', '八、', '九、', '十、')):
                    numbered_count += 1
            if numbered_count >= 3:
                has_numbered_list = True
            
            if has_agenda_title and has_numbered_list:
                return "agenda"
            
            # ========== 章节页识别 ==========
            # 章节页特征：
            # 1. 内容很少（只有标题+可能的副标题）
            # 2. 文字少（字符数少）
            # 3. 形状少
            # 4. 标题通常比较大（但这里不好判断字体大小，用内容量替代）
            is_content_light = (
                text_line_count <= 4 
                and text_shape_count <= 4 
                and total_shapes <= 12
                and total_chars < 100
            )
            
            if is_content_light and text_line_count <= 3:
                # 可能是章节/分隔页
                # 排除封面和结尾（已经判断过了）
                if slide_index > 0 and not (is_last_page and has_closing_kw):
                    return "section"
        
        except Exception:
            pass
        
        # 默认是内容页
        return "content"
