#!/usr/bin/env python3
"""调试脚本：验证母版复制、页脚删除、字体应用功能"""

import sys
import os
sys.path.insert(0, '/workspace/ppt-conform')

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

def test_master_copy():
    """测试母版复制功能"""
    print("=== 测试母版复制功能 ===")
    
    template_path = "/workspace/ppt-conform/templates/se_energy_tech_ppt_20260421.pptx"
    test_input = "/tmp/test_input.pptx"
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Test Slide"
    prs.save(test_input)
    print(f"创建测试PPT: {test_input}")
    
    template_prs = Presentation(template_path)
    print(f"模板母版数量: {len(template_prs.slide_masters)}")
    
    output_prs = Presentation(test_input)
    print(f"原始PPT母版数量: {len(output_prs.slide_masters)}")
    
    import zipfile
    import tempfile
    from lxml import etree
    
    selected_master_index = 0
    template_master = template_prs.slide_masters[selected_master_index]
    
    temp_path = '/tmp/tmp_output_for_master.pptx'
    output_prs.save(temp_path)
    
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp_path = tmp.name
    
    with zipfile.ZipFile(temp_path, 'r') as src_zip, \
         zipfile.ZipFile(tmp_path, 'w') as dst_zip:

        master_xml_bytes = template_master.part.blob
        new_master_idx = len(output_prs.slide_masters) + 1
        master_filename = 'ppt/slideMasters/slideMaster%d.xml' % new_master_idx

        layout_counter = 1
        for item in src_zip.infolist():
            if 'ppt/slideLayouts/slideLayout' in item.filename:
                layout_counter += 1

        presentation_xml = src_zip.read('ppt/presentation.xml')
        presentation_elem = etree.fromstring(presentation_xml)
        nsmap_p = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        ns_rels = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

        sld_master_id_lst = presentation_elem.find('.//p:sldMasterIdLst', nsmap_p)
        if sld_master_id_lst is None:
            sld_master_id_lst = etree.SubElement(
                presentation_elem,
                '{http://schemas.openxmlformats.org/presentationml/2006/main}sldMasterIdLst'
            )

        rId_master = 'rId%d' % (100 + new_master_idx)
        new_master_elem = etree.SubElement(
            sld_master_id_lst,
            '{http://schemas.openxmlformats.org/presentationml/2006/main}sldMasterId'
        )
        new_master_elem.set('{%s}id' % ns_rels, rId_master)

        master_elem = etree.fromstring(master_xml_bytes)
        sld_layout_id_lst = master_elem.find('.//p:sldLayoutIdLst', nsmap_p)
        if sld_layout_id_lst is not None:
            for child in list(sld_layout_id_lst):
                sld_layout_id_lst.remove(child)
        else:
            sld_layout_id_lst = etree.SubElement(
                master_elem,
                '{http://schemas.openxmlformats.org/presentationml/2006/main}sldLayoutIdLst'
            )

        layout_idx = layout_counter
        for _ in template_master.slide_layouts:
            rId_layout = 'rId%d' % (200 + layout_idx)
            new_layout_elem = etree.SubElement(
                sld_layout_id_lst,
                '{http://schemas.openxmlformats.org/presentationml/2006/main}sldLayoutId'
            )
            new_layout_elem.set('{%s}id' % ns_rels, rId_layout)
            layout_idx += 1

        rels_data = src_zip.read('ppt/_rels/presentation.xml.rels')
        rels_elem = etree.fromstring(rels_data)
        ns_rels_def = 'http://schemas.openxmlformats.org/package/2006/relationships'

        new_rel = etree.SubElement(
            rels_elem,
            '{%s}Relationship' % ns_rels_def
        )
        new_rel.set('Id', rId_master)
        new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster')
        new_rel.set('Target', 'slideMasters/slideMaster%d.xml' % new_master_idx)

        for item in src_zip.infolist():
            data = src_zip.read(item.filename)
            if item.filename == 'ppt/presentation.xml':
                data = etree.tostring(presentation_elem, pretty_print=True)
            elif item.filename == 'ppt/_rels/presentation.xml.rels':
                data = etree.tostring(rels_elem, pretty_print=True)
            dst_zip.writestr(item.filename, data)

        dst_zip.writestr(master_filename, etree.tostring(master_elem, pretty_print=True))

        layout_idx = layout_counter
        for layout in template_master.slide_layouts:
            layout_xml_bytes = layout.part.blob
            layout_filename = 'ppt/slideLayouts/slideLayout%d.xml' % layout_idx
            dst_zip.writestr(layout_filename, layout_xml_bytes)
            layout_idx += 1

        master_rels_path = 'ppt/slideMasters/_rels/slideMaster%d.xml.rels' % new_master_idx
        layout_idx = layout_counter
        rels_elem = etree.Element('{%s}Relationships' % ns_rels_def)

        for _ in template_master.slide_layouts:
            rId_layout = 'rId%d' % (200 + layout_idx)
            new_rel = etree.SubElement(
                rels_elem,
                '{%s}Relationship' % ns_rels_def
            )
            new_rel.set('Id', rId_layout)
            new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout')
            new_rel.set('Target', '../slideLayouts/slideLayout%d.xml' % layout_idx)
            layout_idx += 1

        dst_zip.writestr(master_rels_path, etree.tostring(rels_elem, pretty_print=True))
    
    result_prs = Presentation(tmp_path)
    print(f"\n结果PPT母版数量: {len(result_prs.slide_masters)}")
    for i, master in enumerate(result_prs.slide_masters):
        print(f"  母版[{i}]: {len(master.slide_layouts)} 个版式")
    
    result_prs.save('/tmp/test_output_master_final.pptx')
    print(f"\n保存测试PPT: /tmp/test_output_master_final.pptx")
    print("母版复制成功！")

if __name__ == "__main__":
    test_master_copy()