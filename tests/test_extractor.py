import pytest
from pathlib import Path
from pptx import Presentation
from core.extractor.pptx_extractor import PptxExtractor
from core.models import SlideContentModel


def test_extract_basic():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Test Title"

    temp_path = Path("/tmp/test_extract.pptx")
    prs.save(temp_path)

    extractor = PptxExtractor()
    models = extractor.extract(str(temp_path))

    assert len(models) == 1
    assert models[0].title == "Test Title"
    temp_path.unlink()


def test_extract_with_watermark_filter():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "正常标题"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "TRAE AI 生成"

    temp_path = Path("/tmp/test_watermark.pptx")
    prs.save(temp_path)

    extractor = PptxExtractor()
    models = extractor.extract(str(temp_path))

    assert len(models) == 1
    assert models[0].title == "正常标题"
    temp_path.unlink()