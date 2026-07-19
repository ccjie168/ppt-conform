import pytest
from pathlib import Path
from pptx import Presentation
from core.replayer.content_replayer import ContentReplayer
from core.models import SlideContentModel, ContentBlock, UserConfig
from core.registry.template_registry import TemplateRegistry


def test_replay_basic():
    registry = TemplateRegistry()

    content_models = [
        SlideContentModel(
            slide_index=0,
            title="测试封面",
            body_blocks=[],
            original_layout_type="cover"
        ),
        SlideContentModel(
            slide_index=1,
            title="测试内容",
            body_blocks=[ContentBlock(type="paragraph", text="这是正文内容")],
            original_layout_type="content"
        )
    ]

    # 创建临时输入PPT文件
    temp_input = Path("/tmp/input_test.pptx")
    if not temp_input.exists():
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        prs.save(str(temp_input))

    config = UserConfig(
        input_path=str(temp_input),
        output_path="/tmp/output.pptx",
        master_style="F1"
    )

    replayer = ContentReplayer(registry)
    output_path = replayer.replay(content_models, config)

    assert Path(output_path).exists()

    prs = Presentation(output_path)
    assert len(prs.slides) == 2
    assert prs.slides[0].shapes.title.text == "测试封面"

    Path(output_path).unlink()
    temp_input.unlink()