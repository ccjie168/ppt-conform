import pytest
import tempfile
import os
import shutil
from pptx import Presentation
from core.migrator.slide_migrator import SlideMigrator

TEMPLATE_PATH = "templates/2026 se template eng.pptx"


def _make_source_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "My Title"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "My Subtitle"
    prs.save(path)
    return path


def test_migrate_cover_slide():
    with tempfile.TemporaryDirectory() as tmpdir:
        source_path = os.path.join(tmpdir, "source.pptx")
        _make_source_pptx(source_path)

        output_path = os.path.join(tmpdir, "output.pptx")
        shutil.copy(TEMPLATE_PATH, output_path)

        migrator = SlideMigrator(TEMPLATE_PATH, master_index=2)  # Dark Green

        source_prs = Presentation(source_path)
        target_prs = Presentation(output_path)

        new_slide = migrator.migrate_slide(
            source_slide=source_prs.slides[0],
            target_prs=target_prs,
            slide_type="Cover",
            layout_index=0,
        )

        assert new_slide is not None
        assert new_slide.shapes.title is not None
        assert "My Title" in new_slide.shapes.title.text

        target_prs.save(output_path)

        assert os.path.exists(output_path)
        prs2 = Presentation(output_path)
        assert len(prs2.slides) >= 1
