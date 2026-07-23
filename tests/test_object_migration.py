import os
import tempfile
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from core.replayer.content_replayer import ContentReplayer
from core.registry.template_registry import TemplateRegistry


TEMPLATE_PATH = os.path.join(
    os.path.dirname(__file__), "..", "templates", "2026 se template eng.pptx"
)


def _make_source_with_objects(path):
    """Create a source PPT with picture, table, text-box, auto-shape, line."""
    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)

    # 1. Title text box (non-placeholder)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(4), Inches(0.8)
    )
    title_box.text_frame.text = "Test Title"

    # 2. Body text box
    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(4), Inches(1)
    )
    body_box.text_frame.text = "Line 1\nLine 2\nLine 3"

    # 3. Auto-shape (rectangle)
    rect = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(5), Inches(0.5), Inches(2), Inches(1)
    )

    # 4. Line
    line = slide.shapes.add_shape(
        13,  # MSO_SHAPE.LINE_INVERSE
        Inches(5), Inches(2), Inches(2), Inches(0.1)
    )

    # 5. Table
    table = slide.shapes.add_table(
        2, 3,
        Inches(0.5), Inches(3), Inches(4), Inches(1.5)
    ).table
    table.cell(0, 0).text = "A1"
    table.cell(0, 1).text = "B1"
    table.cell(0, 2).text = "C1"
    table.cell(1, 0).text = "A2"
    table.cell(1, 1).text = "B2"
    table.cell(1, 2).text = "C2"

    # 6. Decorative footer text (should be filtered out)
    footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(2), Inches(0.3)
    )
    footer.text_frame.text = "se.com"

    prs.save(path)
    return path


def test_adaptation_preserves_objects():
    """Adaptation path should keep pictures, tables, shapes, text boxes."""
    with tempfile.TemporaryDirectory() as tmpdir:
        src = os.path.join(tmpdir, "source.pptx")
        out = os.path.join(tmpdir, "output.pptx")
        _make_source_with_objects(src)

        registry = TemplateRegistry()
        replayer = ContentReplayer(registry, template_path=TEMPLATE_PATH)
        out_path, qa_items = replayer.convert_with_classification(
            src, out, background_style="dark_green"
        )

        assert len(qa_items) == 1
        result_prs = Presentation(out_path)
        assert len(result_prs.slides) == 1
        slide = result_prs.slides[0]

        # Count shape types
        type_counts = {}
        for shape in slide.shapes:
            st = shape.shape_type
            type_counts[st] = type_counts.get(st, 0) + 1

        # Must have table
        assert MSO_SHAPE_TYPE.TABLE in type_counts, \
            f"Table missing. Types: {type_counts}"

        # Must have at least one auto-shape (the rectangle)
        assert MSO_SHAPE_TYPE.AUTO_SHAPE in type_counts, \
            f"AutoShape missing. Types: {type_counts}"

        # Verify table content preserved
        tables = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        assert len(tables) == 1
        tbl = tables[0].table
        assert tbl.cell(0, 0).text == "A1"
        assert tbl.cell(1, 2).text == "C2"

        # Footer "se.com" should be filtered out
        footer_texts = [s.text_frame.text for s in slide.shapes
                        if s.has_text_frame and "se.com" in s.text_frame.text]
        assert len(footer_texts) == 0, "Decorative footer was not filtered"
