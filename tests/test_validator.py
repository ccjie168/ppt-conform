import pytest
from pathlib import Path
from pptx import Presentation
from core.validator.validator import Validator
from core.models import ValidationReport


def test_validator_basic():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test"

    temp_path = Path("/tmp/test_validate.pptx")
    prs.save(temp_path)

    validator = Validator()
    report = validator.validate(str(temp_path))

    assert isinstance(report, ValidationReport)
    temp_path.unlink()


def test_watermark_validation_fail():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.placeholders[1]
    body.text_frame.text = "TRAE AI 生成"

    temp_path = Path("/tmp/test_watermark_fail.pptx")
    prs.save(temp_path)

    validator = Validator()
    report = validator.validate(str(temp_path))

    assert report.passed is False
    assert any(issue.rule_id == "R040" for issue in report.issues)
    temp_path.unlink()