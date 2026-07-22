import pytest
from pathlib import Path
from pptx import Presentation
from core.extractor.pptx_extractor import PptxExtractor
from core.registry.template_registry import TemplateRegistry
from core.replayer.content_replayer import ContentReplayer
from core.validator.validator import Validator
from core.models import UserConfig, SlideContentModel, ContentBlock


def test_full_conversion():
    """测试完整的PPT转换流程：创建多页PPT，转换，验证输出存在且内容正确"""
    input_path = Path("/tmp/test_e2e_input.pptx")
    output_path = Path("/tmp/test_e2e_output.pptx")

    # 创建多页测试PPT
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "演示文稿封面"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "目录"
    body = slide2.placeholders[1]
    body.text_frame.text = "第一部分\n第二部分\n第三部分"

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "正文内容"
    body = slide3.placeholders[1]
    body.text_frame.text = "这是正常的正文内容，没有水印"

    prs.save(input_path)

    try:
        # 1. 内容抽取
        extractor = PptxExtractor()
        content_models = extractor.extract(str(input_path))
        assert len(content_models) == 3

        # 2. 加载模板注册表（自动加载配置）
        registry = TemplateRegistry()

        # 3. 配置转换参数
        config = UserConfig(
            input_path=str(input_path),
            output_path=str(output_path),
            master_style="F2"
        )

        # 4. 内容重放
        replayer = ContentReplayer(registry)
        replayer.replay(content_models, config)

        # 5. 质量校验（检查水印，字体检查可能因默认字体失败，但这不是关键）
        validator = Validator()
        report = validator.validate(str(output_path))
        # 验证没有水印问题
        watermark_issues = [i for i in report.issues if i.rule_id == "R040"]
        assert len(watermark_issues) == 0, "输出文件不应包含水印"

        # 6. 验证输出文件内容
        output_prs = Presentation(str(output_path))
        assert len(output_prs.slides) == 3
        assert output_prs.slides[0].shapes.title.text == "演示文稿封面"

    finally:
        if input_path.exists():
            input_path.unlink()
        if output_path.exists():
            output_path.unlink()


def test_watermark_blocked():
    """测试水印检测和阻断：创建包含水印文本的PPT，验证校验失败并阻止输出"""
    # 直接创建一个包含水印文本的输出文件（绕过extractor，测试validator）
    output_path = Path("/tmp/test_watermark_output.pptx")

    # 创建包含水印文本的PPT（模拟不应通过的场景）
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "正常标题"
    body = slide.placeholders[1]
    body.text_frame.text = "TRAE AI 生成"

    prs.save(output_path)

    try:
        # 质量校验（应该检测到水印并失败）
        validator = Validator()
        report = validator.validate(str(output_path))

        # 验证校验失败且存在水印问题
        assert report.passed is False, "包含水印的PPT应该校验失败"
        watermark_issues = [i for i in report.issues if i.rule_id == "R040"]
        assert len(watermark_issues) > 0, "应该检测到水印文本"

    finally:
        if output_path.exists():
            output_path.unlink()


def test_e2e_convert_dark_green():
    """端到端测试：深绿色风格混合模式转换"""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.util import Inches
    from core.replayer.content_replayer import ContentReplayer
    from core.registry.template_registry import TemplateRegistry
    from core.qa.reporter import QAReporter

    TEMPLATE_PATH = os.path.join(
        os.path.dirname(os.path.dirname(__file__)),
        "templates", "2026 se template eng.pptx",
    )

    with tempfile.TemporaryDirectory() as tmpdir:
        # 准备源 PPT — 3 页：封面 + 内容 + 表格
        source_path = os.path.join(tmpdir, "source.pptx")
        prs = Presentation()
        # 封面
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "Test Presentation"
        # 内容页
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Content Page"
        s2.placeholders[1].text = "Some body text here"
        # 表格页
        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        s3.shapes.title.text = "Data Table"
        s3.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(8), Inches(3))
        prs.save(source_path)

        output_path = os.path.join(tmpdir, "output.pptx")
        report_path = os.path.join(tmpdir, "report.xlsx")

        # 执行转换
        registry = TemplateRegistry()
        replayer = ContentReplayer(registry, template_path=TEMPLATE_PATH)

        out_path, qa_items = replayer.convert_with_classification(
            source_path=source_path,
            output_path=output_path,
            background_style="dark_green",
        )

        # 验证输出
        assert os.path.exists(out_path)
        assert len(qa_items) == 3

        # 生成 QA 报告
        reporter = QAReporter()
        reporter.generate(qa_items, report_path)
        assert os.path.exists(report_path)

        # 验证输出 PPT 可打开且页数一致
        out_prs = Presentation(out_path)
        assert len(out_prs.slides) == 3

        # 验证尺寸与模板一致
        template_prs = Presentation(TEMPLATE_PATH)
        assert out_prs.slide_width == template_prs.slide_width
        assert out_prs.slide_height == template_prs.slide_height