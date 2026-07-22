import pytest
import tempfile
import os
from pptx import Presentation
from core.precheck.analyzer import PreCheckAnalyzer


def _make_simple_pptx(path):
    prs = Presentation()
    # 强制 16:9 比例
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides[0].shapes.title.text = "Test Title"
    prs.save(path)
    return path


def test_precheck_basic_info():
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "test.pptx")
        _make_simple_pptx(pptx_path)
        
        analyzer = PreCheckAnalyzer()
        result = analyzer.analyze(pptx_path)
        
        assert result.slide_count == 1
        assert result.master_count >= 1
        assert result.is_4_3_ratio is False  # default is 16:9


def test_precheck_fonts_detected():
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "test.pptx")
        _make_simple_pptx(pptx_path)
        
        analyzer = PreCheckAnalyzer()
        result = analyzer.analyze(pptx_path)
        
        assert isinstance(result.fonts_used, list)


def test_precheck_no_media_no_animation_simple_pptx():
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "test.pptx")
        _make_simple_pptx(pptx_path)
        
        analyzer = PreCheckAnalyzer()
        result = analyzer.analyze(pptx_path)
        
        assert result.has_media is False
        assert result.has_animation is False
        assert result.has_smartart is False
        assert result.has_embedded_chart is False
