import pytest
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches
from core.replayer.content_replayer import ContentReplayer
from core.registry.template_registry import TemplateRegistry

TEMPLATE_PATH = "templates/2026 se template eng.pptx"


def test_convert_with_classification():
    with tempfile.TemporaryDirectory() as tmpdir:
        # Create source PPT
        source_path = os.path.join(tmpdir, "source.pptx")
        prs = Presentation()
        # Cover slide
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "Test Presentation"
        # Content slide
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Content Page"
        s2.placeholders[1].text = "Some body text here with enough content"
        # Table slide
        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        s3.shapes.title.text = "Data Table"
        s3.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(8), Inches(3))
        prs.save(source_path)

        output_path = os.path.join(tmpdir, "output.pptx")

        registry = TemplateRegistry()
        replayer = ContentReplayer(registry, template_path=TEMPLATE_PATH)

        out_path, qa_items = replayer.convert_with_classification(
            source_path=source_path,
            output_path=output_path,
            background_style="dark_green",
        )

        assert os.path.exists(out_path)
        assert len(qa_items) == 3
        assert qa_items[0].slide_no == 1
        assert qa_items[0].detected_type == "Cover"
        assert qa_items[0].migration_mode == "migration"

        # Verify output PPT
        out_prs = Presentation(out_path)
        assert len(out_prs.slides) == 3

        # Verify size matches template
        tpl_prs = Presentation(TEMPLATE_PATH)
        assert out_prs.slide_width == tpl_prs.slide_width
        assert out_prs.slide_height == tpl_prs.slide_height


def test_convert_white_style():
    with tempfile.TemporaryDirectory() as tmpdir:
        source_path = os.path.join(tmpdir, "source.pptx")
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.slides[0].shapes.title.text = "White Style Test"
        prs.save(source_path)

        output_path = os.path.join(tmpdir, "output.pptx")
        registry = TemplateRegistry()
        replayer = ContentReplayer(registry, template_path=TEMPLATE_PATH)

        out_path, qa_items = replayer.convert_with_classification(
            source_path=source_path,
            output_path=output_path,
            background_style="white",
        )

        assert os.path.exists(out_path)
        assert len(qa_items) == 1
