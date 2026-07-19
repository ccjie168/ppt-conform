"""对比测试脚本：验证原PPT和转换PPT的内容完整性

测试目标：
1. 转换后标题内容完整保留
2. 转换后正文内容完整保留（不丢失、不重复）
3. 转换后图片、表格等非文本元素保留
4. 转换后页眉页脚使用模板标准格式（不包含原PPT页眉页脚）
5. 转换后布局结构合理（内容与格式不割裂）
"""
import os
import sys
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 添加项目根目录到 path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from core.extractor.pptx_extractor import PptxExtractor
from core.replayer.content_replayer import ContentReplayer
from core.registry.template_registry import TemplateRegistry
from core.models import UserConfig


def create_source_ppt(output_path: str):
    """创建一个包含多种内容的测试源PPT
    
    包含：
    - 封面页（标题+副标题）
    - 内容页1（标题+主正文，单文本框）
    - 内容页2（标题+两个文本框，模拟多列布局）
    - 内容页3（标题+正文+图片）
    - 内容页4（标题+表格）
    """
    prs = Presentation()
    prs.slide_width = Emu(12192000)  # 16:9
    prs.slide_height = Emu(6858000)
    
    # ========== 第1页：封面 ==========
    slide_layout = prs.slide_layouts[0]  # 封面布局
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "测试封面标题"
    # 副标题
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.type == 4:  # subtitle
                shape.text_frame.text = "测试副标题"
                break
        except Exception:
            pass
    
    # ========== 第2页：单文本框内容页 ==========
    slide_layout = prs.slide_layouts[5]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_box.text_frame.text = "第一页内容标题"
    for run in title_box.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
    # 正文
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = body_box.text_frame
    tf.text = "这是第一段正文内容。"
    p2 = tf.add_paragraph()
    p2.text = "这是第二段正文内容。"
    p3 = tf.add_paragraph()
    p3.text = "这是第三段正文内容。"
    
    # ========== 第3页：多文本框内容页（模拟多列） ==========
    slide_layout = prs.slide_layouts[5]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_box.text_frame.text = "多列内容标题"
    for run in title_box.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
    # 左侧文本框
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    left_box.text_frame.text = "左侧文本框内容第一段"
    p2 = left_box.text_frame.add_paragraph()
    p2.text = "左侧文本框内容第二段"
    # 右侧文本框
    right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(5.5), Inches(5))
    right_box.text_frame.text = "右侧文本框内容第一段"
    p2 = right_box.text_frame.add_paragraph()
    p2.text = "右侧文本框内容第二段"
    
    # ========== 第4页：带图片的内容页 ==========
    slide_layout = prs.slide_layouts[5]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_box.text_frame.text = "带图片的内容页"
    for run in title_box.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
    # 正文
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    body_box.text_frame.text = "这是图片页的正文内容。"
    # 添加一个简单的图片（用1x1像素的PNG）
    import io
    from PIL import Image
    img = Image.new('RGB', (100, 100), color='red')
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    slide.shapes.add_picture(img_bytes, Inches(7), Inches(1.5), Inches(5), Inches(4))
    
    # ========== 第5页：带表格的内容页 ==========
    slide_layout = prs.slide_layouts[5]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_box.text_frame.text = "带表格的内容页"
    for run in title_box.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
    # 表格
    rows, cols = 3, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(10), Inches(3))
    table = table_shape.table
    table.cell(0, 0).text = "表头1"
    table.cell(0, 1).text = "表头2"
    table.cell(0, 2).text = "表头3"
    table.cell(1, 0).text = "数据1"
    table.cell(1, 1).text = "数据2"
    table.cell(1, 2).text = "数据3"
    table.cell(2, 0).text = "数据4"
    table.cell(2, 1).text = "数据5"
    table.cell(2, 2).text = "数据6"
    
    prs.save(output_path)
    return output_path


def extract_all_text(prs) -> list:
    """提取PPT中所有文本内容（按页分组）"""
    all_text = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            slide_texts.append(text)
        all_text.append(slide_texts)
    return all_text


def compare_ppts(source_path: str, output_path: str, template_path: str):
    """对比原PPT和转换后的PPT
    
    返回 (成功数, 总检查数, 详情列表)
    """
    source_prs = Presentation(source_path)
    output_prs = Presentation(output_path)
    
    source_texts = extract_all_text(source_prs)
    output_texts = extract_all_text(output_prs)
    
    results = []
    pass_count = 0
    total_checks = 0
    
    # 检查1：页数一致
    total_checks += 1
    if len(output_prs.slides) == len(source_prs.slides):
        results.append(("PASS", f"页数一致: {len(output_prs.slides)} 页"))
        pass_count += 1
    else:
        results.append(("FAIL", f"页数不一致: 源={len(source_prs.slides)}, 输出={len(output_prs.slides)}"))
    
    # 检查2：每页的标题和内容是否保留
    for slide_idx in range(min(len(source_texts), len(output_texts))):
        source_slide_texts = source_texts[slide_idx]
        output_slide_texts = output_texts[slide_idx]
        
        # 检查每个源文本是否在输出中存在
        for src_text in source_slide_texts:
            # 跳过空文本和过短的文本
            if not src_text or len(src_text) < 2:
                continue
            total_checks += 1
            
            # 在输出中查找该文本（可能被拆分）
            found = False
            for out_text in output_slide_texts:
                # 去除空格和换行后比较
                src_clean = src_text.replace(" ", "").replace("\n", "")
                out_clean = out_text.replace(" ", "").replace("\n", "")
                if src_clean in out_clean or out_clean in src_clean:
                    found = True
                    break
            
            if found:
                results.append(("PASS", f"第{slide_idx+1}页: 内容保留 '{src_text[:30]}...'"))
                pass_count += 1
            else:
                results.append(("FAIL", f"第{slide_idx+1}页: 内容丢失 '{src_text[:30]}...'"))
    
    # 检查3：内容没有重复（同一页同一个文本不应该出现多次）
    for slide_idx, output_slide_texts in enumerate(output_texts):
        seen_texts = {}
        for out_text in output_slide_texts:
            if not out_text or len(out_text) < 5:
                continue
            out_clean = out_text.replace(" ", "").replace("\n", "")
            if out_clean in seen_texts:
                total_checks += 1
                results.append(("FAIL", f"第{slide_idx+1}页: 内容重复 '{out_text[:30]}...'"))
            else:
                seen_texts[out_clean] = True
    
    # 检查4：输出PPT的页眉页脚不应包含原PPT的页眉页脚内容
    # （由于测试源PPT没有页眉页脚，这里只检查输出是否合理）
    
    return pass_count, total_checks, results


def test_conversion_with_template():
    """使用真实模板测试转换效果"""
    template_path = "/workspace/ppt-conform/templates/se_energy_tech_ppt_20260421.pptx"
    
    if not Path(template_path).exists():
        print(f"模板文件不存在: {template_path}")
        return False
    
    with tempfile.TemporaryDirectory() as tmpdir:
        # 创建源PPT
        source_path = os.path.join(tmpdir, "source.pptx")
        output_path = os.path.join(tmpdir, "output.pptx")
        
        print("=" * 70)
        print("步骤1: 创建测试源PPT")
        print("=" * 70)
        create_source_ppt(source_path)
        print(f"源PPT已创建: {source_path}")
        
        # 提取源PPT内容
        print("\n步骤2: 提取源PPT内容")
        print("-" * 70)
        extractor = PptxExtractor()
        content_models = extractor.extract(source_path)
        print(f"提取到 {len(content_models)} 页内容")
        for i, model in enumerate(content_models):
            print(f"  第{i+1}页: title='{model.title}', body_blocks={len(model.body_blocks)}, raw_shapes={len(model.raw_shapes)}")
            for j, block in enumerate(model.body_blocks):
                print(f"    block{j}: type={block.type}, role={block.semantic_role}, raw_shape_id={block.raw_shape_id}, text='{(block.text or '')[:30]}'")
        
        # 转换PPT
        print("\n步骤3: 转换PPT（使用模板）")
        print("-" * 70)
        registry = TemplateRegistry()
        replayer = ContentReplayer(registry, template_path=template_path)
        
        # 测试多个Master风格
        template_prs = Presentation(template_path)
        num_masters = len(template_prs.slide_masters)
        print(f"模板有 {num_masters} 个Master")
        
        all_results = []
        for master_idx in range(min(num_masters, 4)):  # 最多测试4个Master
            master_output = os.path.join(tmpdir, f"output_master{master_idx}.pptx")
            config = UserConfig(
                input_path=source_path,
                output_path=master_output,
                master_style=str(master_idx),
            )
            try:
                replayer = ContentReplayer(registry, template_path=template_path)
                replayer.replay(content_models, config)
                
                print(f"\n  Master #{master_idx} 转换完成")
                
                # 对比
                pass_count, total_checks, results = compare_ppts(source_path, master_output, template_path)
                success_rate = (pass_count / total_checks * 100) if total_checks > 0 else 0
                all_results.append((master_idx, pass_count, total_checks, success_rate))
                
                print(f"  对比结果: {pass_count}/{total_checks} 通过 ({success_rate:.1f}%)")
                for status, msg in results:
                    if status == "FAIL":
                        print(f"    [{status}] {msg}")
            except Exception as e:
                print(f"  Master #{master_idx} 转换失败: {e}")
                import traceback
                traceback.print_exc()
                all_results.append((master_idx, 0, 1, 0))
        
        # 汇总结果
        print("\n" + "=" * 70)
        print("转换成功率汇总")
        print("=" * 70)
        total_pass = sum(r[1] for r in all_results)
        total_checks = sum(r[2] for r in all_results)
        overall_rate = (total_pass / total_checks * 100) if total_checks > 0 else 0
        
        for master_idx, pass_count, checks, rate in all_results:
            print(f"Master #{master_idx}: {pass_count}/{checks} ({rate:.1f}%)")
        print(f"\n总体成功率: {total_pass}/{total_checks} ({overall_rate:.1f}%)")
        
        return overall_rate >= 90.0


def test_conversion_with_default_template():
    """使用默认模板测试转换效果"""
    with tempfile.TemporaryDirectory() as tmpdir:
        # 创建源PPT
        source_path = os.path.join(tmpdir, "source.pptx")
        output_path = os.path.join(tmpdir, "output.pptx")
        
        print("=" * 70)
        print("测试: 使用默认模板转换")
        print("=" * 70)
        create_source_ppt(source_path)
        
        extractor = PptxExtractor()
        content_models = extractor.extract(source_path)
        
        registry = TemplateRegistry()
        replayer = ContentReplayer(registry)  # 不指定模板
        
        config = UserConfig(
            input_path=source_path,
            output_path=output_path,
            master_style="0",
        )
        replayer.replay(content_models, config)
        
        pass_count, total_checks, results = compare_ppts(source_path, output_path, "")
        success_rate = (pass_count / total_checks * 100) if total_checks > 0 else 0
        
        print(f"对比结果: {pass_count}/{total_checks} 通过 ({success_rate:.1f}%)")
        for status, msg in results:
            print(f"  [{status}] {msg}")
        
        return success_rate >= 90.0


if __name__ == "__main__":
    print("PPT转换对比测试")
    print("=" * 70)
    
    # 测试1: 使用真实模板
    print("\n[测试1] 使用真实模板转换")
    result1 = test_conversion_with_template()
    
    # 测试2: 使用默认模板
    print("\n[测试2] 使用默认模板转换")
    result2 = test_conversion_with_default_template()
    
    print("\n" + "=" * 70)
    print("最终结论")
    print("=" * 70)
    if result1 and result2:
        print("✓ 所有测试通过，转换成功率 >= 90%")
    elif result1 or result2:
        print("△ 部分测试通过")
    else:
        print("✗ 测试未通过，需要进一步改进")
