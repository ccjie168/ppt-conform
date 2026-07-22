import pytest
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from core.classifier.slide_classifier import SlideClassifier


def _make_cover_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide.shapes.title.text = "Presentation Title"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Subtitle"
    prs.save(path)
    return path


def _make_table_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Data Table"
    rows, cols = 3, 3
    left = top = Inches(2)
    width = height = Inches(4)
    slide.shapes.add_table(rows, cols, left, top, width, height)
    prs.save(path)
    return path


def test_classify_cover():
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "test.pptx")
        _make_cover_pptx(pptx_path)
        prs = Presentation(pptx_path)
        
        classifier = SlideClassifier()
        results = classifier.classify_all(prs)
        
        assert len(results) == 1
        assert results[0].slide_type == "Cover"
        assert results[0].migration_mode == "migration"


def test_classify_table_page():
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "test.pptx")
        _make_table_pptx(pptx_path)
        prs = Presentation(pptx_path)
        
        classifier = SlideClassifier()
        results = classifier.classify_all(prs)
        
        assert len(results) == 1
        assert results[0].slide_type == "Table_Page"
        assert results[0].migration_mode == "adaptation"


def test_classify_content_default():
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "test.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = "Some Content"
        slide.placeholders[1].text = "Body text here with enough content to not be a section divider"
        prs.save(pptx_path)
        
        classifier = SlideClassifier()
        results = classifier.classify_all(prs)
        
        assert results[0].slide_type == "Content"
